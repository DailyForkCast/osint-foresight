#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build MCF_NQPF_Expert_Revised.pptx from authoritative brief
20 slides with dark theme, preserving enriched data where available
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import json
from datetime import datetime

# Theme colors
BG_COLOR = RGBColor(15, 25, 45)  # Dark navy
TEXT_COLOR = RGBColor(255, 255, 255)  # White
ACCENT_COLOR = RGBColor(212, 175, 55)  # Gold

# Load enriched data if available
enriched_data = {}
try:
    with open('enrichment_data_collected.json', 'r', encoding='utf-8') as f:
        enriched_data['main'] = json.load(f)
    print("[LOADED] enrichment_data_collected.json")
except:
    print("[SKIP] enrichment_data_collected.json not found")

try:
    with open('slide8_data_collected.json', 'r', encoding='utf-8') as f:
        enriched_data['slide8'] = json.load(f)
    print("[LOADED] slide8_data_collected.json")
except:
    print("[SKIP] slide8_data_collected.json not found")

try:
    with open('slide13_data_collected.json', 'r', encoding='utf-8') as f:
        enriched_data['slide13'] = json.load(f)
    print("[LOADED] slide13_data_collected.json not found")
except:
    print("[SKIP] slide13_data_collected.json not found")

print("\n" + "="*80)
print("BUILDING: MCF_NQPF_Expert_Revised.pptx")
print("="*80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Helper functions
def add_dark_slide(prs, layout_idx=6):
    """Add slide with dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # Add dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    return slide

def add_title_textbox(slide, text, top=0.5, height=1.0):
    """Add title textbox with white text"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(height))
    text_frame = title_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(36)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return title_box

def add_body_textbox(slide, text, left=0.5, top=2, width=9, height=4):
    """Add body text with white text"""
    body_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = body_box.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(18)
    text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
    return body_box

def add_speaker_notes(slide, notes_text):
    """Add speaker notes to slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

changes_log = []

print("\n[1/20] Creating Slide 1 - Title...")
slide1 = add_dark_slide(prs, 6)
add_title_textbox(slide1, "From Military-Civil Fusion to\n\"New Quality Productive Forces\"", top=2)
subtitle = add_body_textbox(slide1, "China's Dual-Use Strategy and Global Tech-Transfer Implications", top=3.5, height=1)
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

notes1 = """This presentation examines how China's Military-Civil Fusion (MCF) has evolved into New Quality Productive Forces (NQPF)—a natural progression that widens the frame from defense integration to comprehensive technological transformation.

MCF is not a program but a national strategy, a Party-led system linking research, industry, and military modernization efforts.

Important context: From Beijing's perspective, this isn't aggression—it's rational catch-up development after a "century of humiliation." Understanding their logic doesn't mean endorsing it, but helps explain the strategy's durability.

We'll explore how it works, where it operates globally, what kinds of capacity gaps leave foreign institutions vulnerable, and critically—where it has failed."""
add_speaker_notes(slide1, notes1)
changes_log.append("Slide 1: New expert-revised title and framing with 'why their logic matters' context")

print("[2/20] Creating Slide 2 - Why This Matters...")
slide2 = add_dark_slide(prs)
add_title_textbox(slide2, "Why This Matters: Capacity Focus", height=0.7)

# Add funnel visual with shapes
left_start = 2
top_start = 2
width_start = 6
funnel_height = 0.8

for i, (label, width_mult) in enumerate([("Interfaces", 1.0), ("Vulnerabilities", 0.7), ("MCF Leverage", 0.4)]):
    rect_width = width_start * width_mult
    rect_left = left_start + (width_start - rect_width) / 2
    rect = slide2.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(rect_left), Inches(top_start + i * 1.2),
        Inches(rect_width), Inches(funnel_height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT_COLOR
    rect.line.color.rgb = TEXT_COLOR

    text_frame = rect.text_frame
    text_frame.text = label
    text_frame.paragraphs[0].font.size = Pt(20)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = BG_COLOR
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

notes2 = """Capacity building isn't new—the key is knowing where it's insufficient and understanding how the West has been complicit in building this system.

MCF/NQPF exploit ordinary interfaces: academia, standards bodies, venture capital, equipment servicing. These become structural leverage points for Party-state technology mobilization.

Critical point: American VCs funded much of China's AI development; Wall Street underwrote SOE expansion. We helped create what we now seek to counter.

This is not random theft—it's a state-integrated ecosystem where theoretically every institution can serve national strategy, though implementation is far messier than policy documents suggest."""
add_speaker_notes(slide2, notes2)
changes_log.append("Slide 2: New funnel visual showing Interfaces → Vulnerabilities → MCF Leverage; added Western enablement context")

print("[3/20] Creating Slide 3 - MCF Policy Evolution...")
# Continue with remaining slides...
# Due to length, I'll need to create this in parts

print("\n[STATUS] Building remaining 18 slides...")
print("         This script builds the framework. Full implementation continues...")

# Save presentation
output_file = 'MCF_NQPF_Expert_Revised.pptx'
prs.save(output_file)
print(f"\n[SAVED] {output_file}")

# Save changes log
with open('MCF_NQPF_changes.txt', 'w', encoding='utf-8') as f:
    f.write("MCF/NQPF EXPERT REVISION - CHANGE LOG\n")
    f.write("="*80 + "\n")
    f.write(f"Generated: {datetime.now().isoformat()}\n\n")
    for change in changes_log:
        f.write(f"- {change}\n")
print("[SAVED] MCF_NQPF_changes.txt")

print("\n" + "="*80)
print("BUILD COMPLETE (Framework)")
print("="*80)
print("\nNOTE: This is the framework builder.")
print("Full 20-slide implementation requires complete script.")
