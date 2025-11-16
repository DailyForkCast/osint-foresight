#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Build final slides 10-16 for MCF->NQPF presentation
Loads existing 9-slide presentation and adds remaining slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import math

# Color scheme
NAVY = RGBColor(15, 25, 45)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
GRAY = RGBColor(100, 100, 100)

def add_dark_background(slide):
    """Apply navy background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def add_notes(slide, notes_text):
    """Add speaker notes to slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

# Load existing presentation
prs = Presentation('MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx')
blank_layout = prs.slide_layouts[6]

print("Loaded existing 9-slide presentation")
print("Building slides 10-16...")

# ========== SLIDE 10: Mechanisms Abroad ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Mechanisms Abroad: Civil & Semi-Civil Channels"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 5 Swimlanes with detailed content
swimlanes = [
    ("Academic / R&D Collaboration", [
        "Joint labs & co-authored publications",
        "HIT/NPU: Co-authorships continued after Entity-List designation; formal MoUs largely pre-dated listings",
        "Shadow joint programs via third-country institutions"
    ]),
    ("Talent Recruitment Programs", [
        "1000 Talents Plan and variants (provincial, municipal)",
        "Incentivized return of Western-trained researchers",
        "Shadow recruitment under research grants"
    ]),
    ("Innovation & Transfer Hubs", [
        "Zhongguancun (ZGC) overseas centers",
        "CITTC (China International Technology Transfer Center)",
        "Provincial S&T bureaus with foreign liaison offices"
    ]),
    ("Investment Channels", [
        "VC funds with PLA or SOE backing",
        "Mergers & acquisitions masking end-users",
        "Technology licensing with hidden sub-licensing"
    ]),
    ("Standards Platforms", [
        "Participation in SDOs (ITU, 3GPP, IEEE)",
        "Pre-positioning for 6G, AI governance standards",
        "De-facto control via market dominance"
    ])
]

lane_height = 1.0
y_start = 1.2
for idx, (lane_title, items) in enumerate(swimlanes):
    y_pos = y_start + idx * lane_height

    # Lane container
    lane_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(y_pos),
        Inches(9), Inches(0.9)
    )
    lane_box.fill.solid()
    lane_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
    lane_box.line.color.rgb = WHITE
    lane_box.line.width = Pt(1)

    # Lane title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.05), Inches(8.8), Inches(0.25))
    tf = title_box.text_frame
    tf.text = lane_title
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Lane items (compact bullets)
    items_text = " | ".join(items)
    items_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_pos + 0.32), Inches(8.8), Inches(0.55))
    tf = items_box.text_frame
    tf.text = items_text
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE

add_notes(slide, """Slide 10: Mechanisms Abroad - Civil & Semi-Civil Channels

HIT/NPU PHRASING RULE APPLIED: "Co-authorships continued after Entity-List designation; formal MoUs largely pre-dated listings."

Academic/R&D: Universities abroad engage with PLA-affiliated institutions, often unaware of military linkages. Entity List designations have limited impact on ongoing collaboration.

Talent Recruitment: China's talent programs (1000 Talents, provincial variants) offer substantial incentives for overseas researchers to return, often bringing proprietary knowledge or research.

Innovation Hubs: Zhongguancun (ZGC) operates overseas centers to facilitate technology transfer. CITTC acts as a bridge for international IP licensing. Provincial S&T bureaus maintain foreign liaison offices.

Investment: VC funds with PLA or SOE backing acquire Western startups or IP. M&A activity masks end-users through complex ownership structures.

Standards: Chinese entities participate heavily in standards development organizations (SDOs) like ITU, 3GPP, IEEE to shape 6G, AI governance, and other emerging technology standards.

SOURCES:
- CSET reports on Chinese talent recruitment
- ASPI China Defence Universities Tracker
- BIS Entity List designations
- Open-source ZGC/CITTC documentation

Placeholder data used for specific program names; replace with validated project data when available.
""")

print("  [10/16] Mechanisms Abroad - 5 swimlanes with HIT/NPU phrasing")

# ========== SLIDE 11: Global Examples ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Global Examples: Dual-Use Tech Transfer"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Minimalist world map representation (simplified continents as shapes)
# Africa
africa = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(4.5), Inches(3.5),
    Inches(1.2), Inches(1.5)
)
africa.fill.solid()
africa.fill.fore_color.rgb = RGBColor(60, 70, 90)
africa.line.color.rgb = RGBColor(100, 110, 130)
africa.line.width = Pt(1)

# Europe (smaller oval to the left)
europe = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(4.2), Inches(2.3),
    Inches(0.8), Inches(0.7)
)
europe.fill.solid()
europe.fill.fore_color.rgb = RGBColor(60, 70, 90)
europe.line.color.rgb = RGBColor(100, 110, 130)
europe.line.width = Pt(1)

# Asia (large oval to the right)
asia = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(6.5), Inches(2.5),
    Inches(1.8), Inches(1.3)
)
asia.fill.solid()
asia.fill.fore_color.rgb = RGBColor(60, 70, 90)
asia.line.color.rgb = RGBColor(100, 110, 130)
asia.line.width = Pt(1)

# South America
south_america = slide.shapes.add_shape(
    MSO_SHAPE.PENTAGON,
    Inches(2.5), Inches(4.0),
    Inches(0.8), Inches(1.2)
)
south_america.fill.solid()
south_america.fill.fore_color.rgb = RGBColor(60, 70, 90)
south_america.line.color.rgb = RGBColor(100, 110, 130)
south_america.line.width = Pt(1)

# Pin locations with callouts
pins = [
    ("Kenya (Nairobi)", 5.0, 3.8, "Safe City (Huawei/Safaricom)\nFR + CCTV integration"),
    ("Serbia (Belgrade)", 4.6, 2.5, "Smart Surveillance\nHuawei FR network"),
    ("Argentina (Neuquen*)", 2.7, 4.5, "Deep-Space Ground Station\n*Dual-use debates context"),
    ("Pakistan (CPEC)", 7.0, 3.0, "BeiDou Augmentation\nPakSat-MM1")
]

for pin_label, x, y, desc in pins:
    # Pin (downward triangle)
    pin = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x), Inches(y),
        Inches(0.15), Inches(0.15)
    )
    pin.rotation = 180
    pin.fill.solid()
    pin.fill.fore_color.rgb = GOLD
    pin.line.color.rgb = WHITE
    pin.line.width = Pt(1)

    # Callout box
    callout_x = x + 0.3 if x < 5 else x - 2.5
    callout_y = y - 0.3
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(callout_x), Inches(callout_y),
        Inches(2.2), Inches(0.6)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(30, 40, 60)
    callout.line.color.rgb = GOLD
    callout.line.width = Pt(1.5)

    # Callout text
    tf = callout.text_frame
    tf.text = f"{pin_label}\n{desc}"
    tf.word_wrap = True
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GOLD

# Footnote for Argentina
footnote = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
tf = footnote.text_frame
tf.text = "*Argentina (Neuquen): Frequently cited case in dual-use debates; insert lease/treaty details when available."
p = tf.paragraphs[0]
p.font.size = Pt(10)
p.font.italic = True
p.font.color.rgb = GRAY

add_notes(slide, """Slide 11: Global Examples - Dual-Use Tech Transfer

ARGENTINA FOOTNOTE APPLIED: "Argentina (Neuquen*): Frequently cited case in dual-use debates; insert lease/treaty details when available."

Kenya (Nairobi): Huawei Safe City project with Safaricom integrates facial recognition (FR) and CCTV surveillance. Data flows and access protocols remain opaque. Dual-use concern: law enforcement surveillance tech can inform social control systems.

Serbia (Belgrade): Huawei deployed a comprehensive FR-enabled surveillance network. Serbian government cooperation with Chinese tech firms raises concerns about data sovereignty and potential PLA access to biometric data.

Argentina (Neuquen*): Deep-space ground station operated by Chinese entities. Originally framed as civilian space research, the station's capabilities include military satellite tracking and space-domain awareness. Lease terms and oversight remain contentious.

Pakistan (CPEC/CPEC+): BeiDou augmentation system and PakSat-MM1 satellite provide dual-use positioning, navigation, and timing (PNT) as well as communications infrastructure. PLA integration potential is high.

SOURCES:
- Reuters investigations on Huawei Safe City projects
- CSIS reports on Chinese space infrastructure abroad
- ASPI analysis of BRI/DSR dual-use projects
- Open-source lease agreements (where available)

Placeholder used for specific lease/treaty details; replace with validated project data when available.
""")

print("  [11/16] Global Examples - map with Argentina (Neuquen*) footnote")

# ========== SLIDE 12: Global Implications - Capacity Needs ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Global Implications: Capacity Needs"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 5x2 matrix
matrix_data = [
    ("Data/Privacy", "Foreign 'safe city' cloud tie-ins", "Tech audit & data governance"),
    ("Supply-Chains", "Opaque sub-tier sourcing", "Sub-tier mapping & end-use checks"),
    ("Research", "Dual appointments/hidden affiliations", "Research-security workflows"),
    ("Standards", "Under-participation in SDOs", "Engineer/regulator participation"),
    ("Finance", "Beneficial ownership opacity", "Investment disclosures & screening")
]

row_height = 0.9
y_start = 1.3
col_widths = [2.0, 3.5, 3.5]

# Headers
headers = ["Domain", "Vulnerability", "Needed Capacity"]
for col_idx, header in enumerate(headers):
    x_pos = 0.5 + sum(col_widths[:col_idx])
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_pos), Inches(y_start),
        Inches(col_widths[col_idx]), Inches(0.5)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = GOLD
    header_box.line.color.rgb = WHITE
    header_box.line.width = Pt(1)

    tf = header_box.text_frame
    tf.text = header
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

# Data rows
for row_idx, (domain, vuln, capacity) in enumerate(matrix_data):
    y_pos = y_start + 0.5 + row_idx * row_height
    row_data = [domain, vuln, capacity]

    for col_idx, text in enumerate(row_data):
        x_pos = 0.5 + sum(col_widths[:col_idx])
        cell_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(col_widths[col_idx]), Inches(row_height)
        )
        cell_box.fill.solid()
        cell_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
        cell_box.line.color.rgb = WHITE
        cell_box.line.width = Pt(1)

        tf = cell_box.text_frame
        tf.text = text
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        if col_idx == 0:  # Domain column
            p.font.bold = True
            p.font.color.rgb = GOLD

add_notes(slide, """Slide 12: Global Implications - Capacity Needs Matrix

This slide maps five critical domains where MCF/NQPF activities exploit vulnerabilities, and identifies the institutional capacity needed to address each gap.

Data/Privacy: Foreign "safe city" systems create data dependencies with opaque cloud tie-ins. Host countries lack technical audit capabilities to assess data flows, retention, and access. NEEDED: Tech audit frameworks and data sovereignty governance.

Supply-Chains: Sub-tier sourcing in complex supply chains (semiconductors, aerospace) is often opaque. End-use verification is weak. MCF/NQPF entities exploit this to acquire dual-use components. NEEDED: Sub-tier mapping tools and robust end-use checks.

Research: Dual appointments, hidden affiliations, and shadow joint programs obscure the ultimate beneficiaries of research collaboration. Universities lack screening frameworks. NEEDED: Research security workflows including affiliation disclosure and risk-tiering.

Standards: Under-participation by democratic nations in SDOs (ITU, 3GPP, IEEE) allows Chinese entities to shape standards for 6G, AI governance, quantum comms. NEEDED: Increased engineer and regulator participation in SDOs.

Finance: Beneficial ownership opacity in VC funds and M&A transactions masks PLA or SOE backing. Due diligence is weak. NEEDED: Investment disclosure requirements and screening mechanisms (e.g., CFIUS-like bodies).

SOURCES:
- CSET analysis of research security gaps
- RAND reports on SDO participation
- BIS sub-tier sourcing guidance
- CSIS investment screening studies

All data here is placeholder; replace with validated project data when available.
""")

print("  [12/16] Global Implications - 5x2 capacity needs matrix")

# ========== SLIDE 13: Gray-Zone Tech Acquisition ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Gray-Zone Tech Acquisition (MCF Support)"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Two columns
left_col_x = 1.0
right_col_x = 5.5
col_width = 4.0
col_height = 5.0
y_start = 1.3

# Left column: Legitimate Activities
left_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(left_col_x), Inches(y_start),
    Inches(col_width), Inches(col_height)
)
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
left_box.line.color.rgb = WHITE
left_box.line.width = Pt(1.5)

# Left column title
left_title = slide.shapes.add_textbox(Inches(left_col_x + 0.2), Inches(y_start + 0.1), Inches(col_width - 0.4), Inches(0.4))
tf = left_title.text_frame
tf.text = "Legitimate Activities"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# Left column content
left_content = [
    "Academic collaboration",
    "Conference participation",
    "Technology licensing",
    "Visiting scholar programs",
    "Startup incubators",
    "SDO participation"
]
left_text_box = slide.shapes.add_textbox(Inches(left_col_x + 0.3), Inches(y_start + 0.6), Inches(col_width - 0.6), Inches(col_height - 0.8))
tf = left_text_box.text_frame
for item in left_content:
    p = tf.add_paragraph()
    p.text = f"- {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

# Right column: MCF-Leveraged Outcomes
right_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(right_col_x), Inches(y_start),
    Inches(col_width), Inches(col_height)
)
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(40, 50, 70)
right_box.line.color.rgb = WHITE
right_box.line.width = Pt(1.5)

# Right column title
right_title = slide.shapes.add_textbox(Inches(right_col_x + 0.2), Inches(y_start + 0.1), Inches(col_width - 0.4), Inches(0.4))
tf = right_title.text_frame
tf.text = "MCF-Leveraged Outcomes"
p = tf.paragraphs[0]
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# Right column content
right_content = [
    "Knowledge transfer to PLA-linked labs",
    "Talent recruitment for defense R&D",
    "IP acquisition for dual-use systems",
    "Shadow joint programs via proxies",
    "Tech scouting for military needs",
    "Standards positioning for control"
]
right_text_box = slide.shapes.add_textbox(Inches(right_col_x + 0.3), Inches(y_start + 0.6), Inches(col_width - 0.6), Inches(col_height - 0.8))
tf = right_text_box.text_frame
for item in right_content:
    p = tf.add_paragraph()
    p.text = f"- {item}"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

# Bidirectional arrow between columns
arrow_y = y_start + col_height / 2
# Left arrow
left_arrow = slide.shapes.add_connector(
    1,  # Straight connector
    Inches(left_col_x + col_width), Inches(arrow_y - 0.15),
    Inches(right_col_x), Inches(arrow_y - 0.15)
)
left_arrow.line.color.rgb = GOLD
left_arrow.line.width = Pt(2)

# Right arrow
right_arrow = slide.shapes.add_connector(
    1,
    Inches(right_col_x), Inches(arrow_y + 0.15),
    Inches(left_col_x + col_width), Inches(arrow_y + 0.15)
)
right_arrow.line.color.rgb = GOLD
right_arrow.line.width = Pt(2)

# Label in center
label = slide.shapes.add_textbox(Inches(4.2), Inches(arrow_y - 0.3), Inches(1.6), Inches(0.6))
tf = label.text_frame
tf.text = "Knowledge/\nData/Talent\nFlows"
p = tf.paragraphs[0]
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

add_notes(slide, """Slide 13: Gray-Zone Tech Acquisition (MCF Support)

Gray-zone acquisition refers to activities that are technically legal but exploited by MCF/NQPF to acquire dual-use knowledge, data, and talent. The bidirectional flow indicates that legitimate activities create pathways for MCF-leveraged outcomes.

Legitimate Activities:
- Academic collaboration: Joint research with Western universities
- Conference participation: Presenting and networking at academic/industry conferences
- Technology licensing: Acquiring commercial IP through legal channels
- Visiting scholar programs: Short-term research positions at foreign institutions
- Startup incubators: Chinese entrepreneurs in Western accelerators
- SDO participation: Membership in standards development organizations

MCF-Leveraged Outcomes:
- Knowledge transfer to PLA-linked labs: Research findings funneled to military applications
- Talent recruitment for defense R&D: Identifying and recruiting foreign experts
- IP acquisition for dual-use systems: Licensing tech that has military applications
- Shadow joint programs via proxies: Using third-party institutions to obscure PLA involvement
- Tech scouting for military needs: Identifying technologies of interest to PLA modernization
- Standards positioning for control: Shaping global standards to favor Chinese tech dominance

The key insight: most of these activities are legal and involve willing Western participants who may not be aware of the ultimate end-use or beneficiary.

SOURCES:
- ASPI reports on ZGC and CITTC overseas activities
- CSET analysis of talent recruitment programs
- BIS Entity List case studies
- Open-source academic affiliation databases

Placeholder data; replace with specific cases from project corpus when available.
""")

print("  [13/16] Gray-Zone Tech Acquisition - two columns with bidirectional flow")

# ========== SLIDE 14: Illicit & Clandestine Acquisition ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Illicit & Clandestine Acquisition (MCF-Relevant)"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Threat Tree - root node at center, branches outward
center_x = 5.0
center_y = 2.5

# Root node
root = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(center_x - 1.2), Inches(center_y),
    Inches(2.4), Inches(0.6)
)
root.fill.solid()
root.fill.fore_color.rgb = RGBColor(120, 30, 30)
root.line.color.rgb = GOLD
root.line.width = Pt(2)
tf = root.text_frame
tf.text = "Illicit Acquisition\nMethods"
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Branch nodes
branches = [
    ("Cyber Intrusions", 1.5, 1.5, "APT groups targeting\ndefense contractors"),
    ("Insider Theft", 8.5, 1.5, "Recruitment of insiders\nfor IP exfiltration"),
    ("Export Evasion", 1.5, 4.5, "Transshipment via\nthird countries"),
    ("Procurement Masking", 5.0, 5.2, "Shell companies &\nfalse end-users"),
    ("Cloud Exfiltration", 8.5, 4.5, "Unauthorized access\nto shared platforms")
]

for branch_title, x, y, desc in branches:
    # Connector from root to branch
    connector = slide.shapes.add_connector(
        1,
        Inches(center_x), Inches(center_y + 0.3),
        Inches(x + 0.9), Inches(y + 0.3)
    )
    connector.line.color.rgb = GOLD
    connector.line.width = Pt(1.5)

    # Branch node
    branch = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(1.8), Inches(0.6)
    )
    branch.fill.solid()
    branch.fill.fore_color.rgb = RGBColor(70, 40, 40)
    branch.line.color.rgb = WHITE
    branch.line.width = Pt(1)

    tf = branch.text_frame
    tf.text = branch_title
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Description below branch
    desc_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.7), Inches(1.8), Inches(0.5))
    tf = desc_box.text_frame
    tf.text = desc
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

add_notes(slide, """Slide 14: Illicit & Clandestine Acquisition (MCF-Relevant)

This slide maps illicit and clandestine methods used to acquire dual-use technology in support of MCF/NQPF objectives. These activities are illegal under host-country laws and often violate export controls.

Cyber Intrusions: APT (Advanced Persistent Threat) groups linked to PLA or MSS target defense contractors, aerospace firms, and research institutions. Goal is to exfiltrate IP, research data, and technical specifications. Examples: APT1, APT10, APT40.

Insider Theft: Recruitment of insiders (employees, researchers) at Western companies or labs to steal IP or research findings. Incentivized through financial rewards, nationalist appeals, or coercion. Cases include theft of trade secrets from semiconductor firms, aerospace companies, and biotech labs.

Export Evasion: Use of transshipment routes through third countries to evade export controls. Controlled items are routed through jurisdictions with weaker enforcement (e.g., Hong Kong, UAE, Singapore) before reaching China. Shell companies obscure the end-user.

Procurement Masking: Front companies and false end-user certificates are used to acquire controlled dual-use items. The ultimate recipient (PLA lab, defense contractor) is hidden behind layers of corporate entities.

Cloud Exfiltration: Unauthorized access to cloud platforms used by Western firms for R&D collaboration. Data is exfiltrated from shared environments where access controls are weak.

SOURCES:
- FBI press releases on Chinese espionage cases
- DOJ indictments of PLA officers and MSS agents
- BIS export control enforcement actions
- Crowdstrike, FireEye, and Mandiant APT reports
- CSIS reports on supply-chain vulnerabilities

All data here is illustrative; replace with specific case details from project corpus when available.
""")

print("  [14/16] Illicit Acquisition - threat tree with 5 branches")

# ========== SLIDE 15: Capacity Gaps Map ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Capacity Gaps Map: Where & Why"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# 4-column table
table_data = [
    ("Domain", "Current Weakness", "MCF/NQPF Leverage", "Needed Capacity"),
    ("Academia", "Weak affiliation screening", "Talent programs & shadow joint programs", "Research security frameworks & disclosure requirements"),
    ("Industry", "Supply chain opacity", "Procurement masking & transshipment", "End-use verification & sub-tier mapping"),
    ("Space", "Joint R&D ambiguity", "Dual-use data collection & tracking", "Oversight protocols & tech audits"),
    ("Bio", "Open data sharing norms", "PLA-linked genomic & biotech research", "Biosecurity protocols & benefit-sharing governance")
]

col_widths = [1.5, 2.2, 2.5, 2.8]
row_height = 1.0
y_start = 1.3

for row_idx, row in enumerate(table_data):
    for col_idx, cell_text in enumerate(row):
        x_pos = 0.5 + sum(col_widths[:col_idx])
        y_pos = y_start + row_idx * row_height

        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_pos), Inches(y_pos),
            Inches(col_widths[col_idx]), Inches(row_height)
        )

        if row_idx == 0:  # Header row
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD
            cell.line.color.rgb = WHITE
            cell.line.width = Pt(1.5)

            tf = cell.text_frame
            tf.text = cell_text
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = NAVY
            p.alignment = PP_ALIGN.CENTER
        else:  # Data rows
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(40, 50, 70)
            cell.line.color.rgb = WHITE
            cell.line.width = Pt(1)

            tf = cell.text_frame
            tf.text = cell_text
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT
            if col_idx == 0:  # Domain column
                p.font.bold = True
                p.font.color.rgb = GOLD

add_notes(slide, """Slide 15: Capacity Gaps Map - Where & Why

This slide summarizes actionable capacity responses mapped to institutional weaknesses that MCF/NQPF activities exploit.

Academia:
- Current Weakness: Universities lack robust affiliation screening. Dual appointments and undisclosed foreign affiliations are common.
- MCF/NQPF Leverage: Talent programs (1000 Talents, etc.) and shadow joint programs exploit this weakness to access research and recruit talent.
- Needed Capacity: Research security frameworks, mandatory affiliation disclosure, risk-tiering of collaborations, and institutional training.

Industry:
- Current Weakness: Supply chain opacity, especially at sub-tier levels. End-use verification is weak or absent.
- MCF/NQPF Leverage: Procurement masking via shell companies, transshipment through third countries, and false end-user certificates.
- Needed Capacity: End-use verification protocols, sub-tier mapping tools, and enhanced due diligence in procurement.

Space:
- Current Weakness: Joint R&D arrangements with ambiguous military applications. Facility oversight is limited.
- MCF/NQPF Leverage: Dual-use data collection (satellite tracking, space-domain awareness) and potential PLA access.
- Needed Capacity: Oversight protocols for joint facilities, tech audits of data flows, and clear dual-use red lines.

Bio:
- Current Weakness: Open data sharing norms in genomics and biotech. Benefit-sharing governance is underdeveloped.
- MCF/NQPF Leverage: PLA-linked research institutions access genomic databases and biotech IP for military applications.
- Needed Capacity: Biosecurity protocols, controlled access to sensitive databases, and benefit-sharing agreements that consider security implications.

SOURCES:
- CSET research security reports
- RAND biosecurity studies
- CSIS supply-chain analyses
- BIS guidance on end-use verification

All data here is illustrative; replace with specific policy recommendations from project corpus when available.
""")

print("  [15/16] Capacity Gaps Map - 4-column table")

# ========== SLIDE 16: Key Takeaways & References ==========
slide = prs.slides.add_slide(blank_layout)
add_dark_background(slide)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title.text_frame
tf.text = "Key Takeaways & References"
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Two columns
left_col_x = 0.8
right_col_x = 5.5
col_width = 4.2
y_start = 1.3

# Left column: Key Takeaways
left_title = slide.shapes.add_textbox(Inches(left_col_x), Inches(y_start), Inches(col_width), Inches(0.4))
tf = left_title.text_frame
tf.text = "Key Takeaways"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GOLD

takeaways = [
    "Continuity: MCF to NQPF is a rebrand, not a policy shift. Objectives remain unchanged.",
    "Legal Compulsion: Chinese laws mandate cooperation with state intelligence and military needs.",
    "Global Scope: Tech transfer happens via lawful, gray-zone, and illicit channels abroad.",
    "Targeted Capacity: Gaps exist in academia, industry, space, and bio. Capacity building must be domain-specific.",
    "Urgency: Advanced and emerging dual-use tech is at stake. Inaction cedes strategic advantage."
]

takeaways_box = slide.shapes.add_textbox(Inches(left_col_x), Inches(y_start + 0.5), Inches(col_width), Inches(5.0))
tf = takeaways_box.text_frame
for takeaway in takeaways:
    p = tf.add_paragraph()
    p.text = f"- {takeaway}"
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.space_before = Pt(10)
    p.line_spacing = 1.2

# Right column: References
right_title = slide.shapes.add_textbox(Inches(right_col_x), Inches(y_start), Inches(col_width), Inches(0.4))
tf = right_title.text_frame
tf.text = "References"
p = tf.paragraphs[0]
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = GOLD

references = [
    "CSET (Georgetown): Research security, talent recruitment analyses",
    "ASPI: China Defence Universities Tracker, China Tech Map",
    "RAND: Biosecurity, supply-chain vulnerabilities",
    "CSIS: Investment screening, standards participation",
    "BIS: Entity List designations, export control guidance",
    "Reuters: Investigative reporting on Huawei Safe City projects",
    "FBI/DOJ: Press releases and indictments on espionage cases",
    "Open-source: Policy documents, MoU databases, academic affiliations"
]

references_box = slide.shapes.add_textbox(Inches(right_col_x), Inches(y_start + 0.5), Inches(col_width), Inches(5.0))
tf = references_box.text_frame
for reference in references:
    p = tf.add_paragraph()
    p.text = f"- {reference}"
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.space_before = Pt(6)
    p.line_spacing = 1.1

add_notes(slide, """Slide 16: Key Takeaways & References

KEY TAKEAWAYS:

1. Continuity: The transition from MCF to NQPF is cosmetic. The underlying policy objectives—leveraging civilian innovation for military modernization—remain unchanged. This is a rebrand, not a retreat.

2. Legal Compulsion: Chinese laws (National Security Law 2015, National Intelligence Law 2017 Art.7, Data Security Law 2021, Revised State Secrets Law 2024) create a legal framework compelling all Chinese entities—companies, universities, individuals—to cooperate with state intelligence and military needs. This is not optional.

3. Global Scope: Tech transfer happens through three channels:
   - Lawful: academic collaboration, licensing, SDO participation
   - Gray-zone: legitimate activities leveraged for MCF/NQPF outcomes
   - Illicit: cyber intrusions, insider theft, export evasion

4. Targeted Capacity: Institutional capacity gaps exist in academia (research security), industry (supply-chain transparency), space (facility oversight), and bio (data governance). One-size-fits-all approaches fail. Capacity building must be domain-specific and risk-tiered.

5. Urgency: Advanced and emerging dual-use technologies—AI, quantum, semiconductors, biotech, space—are at stake. Inaction or delayed responses cede strategic advantage to MCF/NQPF-aligned entities.

REFERENCES:
The analysis draws on open-source policy documents, academic research, investigative journalism, and enforcement actions. Key sources include:
- CSET (Center for Security and Emerging Technology, Georgetown): Research security frameworks, talent recruitment program analyses
- ASPI (Australian Strategic Policy Institute): China Defence Universities Tracker, China Tech Map, unitracker
- RAND Corporation: Biosecurity studies, supply-chain vulnerabilities
- CSIS (Center for Strategic and International Studies): Investment screening, standards participation analyses
- BIS (Bureau of Industry and Security): Entity List designations, export control guidance
- Reuters: Investigative reporting on Huawei Safe City projects and Chinese tech infrastructure abroad
- FBI/DOJ: Press releases and indictments documenting espionage cases
- Open-source materials: Chinese policy documents, MoU databases, academic affiliation records

All citations are illustrative; replace with specific references from project corpus when available.
""")

print("  [16/16] Key Takeaways & References - two columns")

# Save presentation
prs.save('MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx')
print("\n" + "="*60)
print("SUCCESS: All 16 slides completed!")
print("File: MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")
print("="*60)
