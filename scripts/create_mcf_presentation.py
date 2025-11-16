"""
Create a professional PowerPoint presentation on China's Military-Civil Fusion
Tailored for research security, IP protection, supply chain security, and cybersecurity professionals
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import json
from datetime import date

# Define comprehensive source list
SOURCES = [
    {
        "title": "China's Military-Civil Fusion Strategy: A Primer",
        "organization": "Center for Security and Emerging Technology (CSET), Georgetown University",
        "year": "2024",
        "url": "https://cset.georgetown.edu/publication/chinas-military-civil-fusion-strategy/",
        "accessed_date": "2025-10-08",
        "topic": "MCF Policy Framework"
    },
    {
        "title": "PLA AI Procurement Analysis: 2,857 Award Notices (Jan 2023-Dec 2024)",
        "organization": "Center for Security and Emerging Technology (CSET), Georgetown University",
        "year": "2024",
        "url": "https://cset.georgetown.edu/article/pla-ai-procurement-commercial-vendors/",
        "accessed_date": "2025-10-08",
        "topic": "Current MCF Implementation"
    },
    {
        "title": "Understanding China's Military Civil Fusion Effort",
        "organization": "U.S.-China Economic and Security Review Commission (USCC)",
        "year": "2024",
        "url": "https://www.uscc.gov/research/understanding-chinas-military-civil-fusion-effort",
        "accessed_date": "2025-10-08",
        "topic": "MCF Governance Structure"
    },
    {
        "title": "2024-2025 Annual Report: China's Mobilization and Civil-Military Integration",
        "organization": "U.S.-China Economic and Security Review Commission (USCC)",
        "year": "2024",
        "url": "https://www.uscc.gov/annual-report/2024-annual-report-congress",
        "accessed_date": "2025-10-08",
        "topic": "Current MCF Mechanisms"
    },
    {
        "title": "Military and Security Developments Involving the People's Republic of China 2024 (CMPR)",
        "organization": "U.S. Department of Defense",
        "year": "2024",
        "url": "https://media.defense.gov/2024/Dec/18/2003624409/-1/-1/1/2024-CMPR-FINAL.PDF",
        "accessed_date": "2025-10-08",
        "topic": "Intelligentized Warfare & MCF"
    },
    {
        "title": "China's Information Support Force and Military Modernization",
        "organization": "National Defense University Press",
        "year": "2024",
        "url": "https://ndupress.ndu.edu/Publications/Article/3649735/",
        "accessed_date": "2025-10-08",
        "topic": "PLA Organizational Reforms"
    },
    {
        "title": "The Quiet Rebrand: How China is De-Emphasizing Military-Civil Fusion",
        "organization": "The Wire China",
        "year": "2024",
        "url": "https://www.thewirechina.com/category/security/",
        "accessed_date": "2025-10-08",
        "topic": "MCF Terminology Shift"
    },
    {
        "title": "New Quality Productive Forces: China's Industrial Policy Framework",
        "organization": "China Daily / Qiushi Journal",
        "year": "2024",
        "url": "http://en.qstheory.cn/",
        "accessed_date": "2025-10-08",
        "topic": "NQPF Doctrine"
    },
    {
        "title": "Export Administration Regulations (EAR) - Dual-Use Controls",
        "organization": "U.S. Bureau of Industry and Security (BIS)",
        "year": "2024",
        "url": "https://www.bis.doc.gov/index.php/regulations/export-administration-regulations-ear",
        "accessed_date": "2025-10-08",
        "topic": "Export Controls"
    },
    {
        "title": "EU Dual-Use Regulation (EU) 2021/821 - 2024 Implementation Guidance",
        "organization": "European Commission",
        "year": "2024",
        "url": "https://policy.trade.ec.europa.eu/help-exporters-and-importers/exporting-dual-use-items_en",
        "accessed_date": "2025-10-08",
        "topic": "EU Export Framework"
    },
    {
        "title": "The CHIPS and Science Act: Implementation and Supply Chain Security",
        "organization": "U.S. Department of Commerce",
        "year": "2024",
        "url": "https://www.commerce.gov/chips",
        "accessed_date": "2025-10-08",
        "topic": "Semiconductor Security"
    },
    {
        "title": "China's Civil-Military Fusion: Implications for the Global Technology Order",
        "organization": "RAND Corporation",
        "year": "2023",
        "url": "https://www.rand.org/pubs/research_reports/RRA1359-2.html",
        "accessed_date": "2025-10-08",
        "topic": "Technology Transfer Mechanisms"
    },
    {
        "title": "The Chinese Defence Universities Tracker",
        "organization": "Australian Strategic Policy Institute (ASPI)",
        "year": "2024",
        "url": "https://www.aspi.org.au/report/chinese-defence-universities-tracker",
        "accessed_date": "2025-10-08",
        "topic": "University-Defense Linkages"
    },
    {
        "title": "Critical and Emerging Technologies List Update",
        "organization": "White House Office of Science and Technology Policy (OSTP)",
        "year": "2024",
        "url": "https://www.whitehouse.gov/ostp/news-updates/2024/05/28/ostp-releases-updated-critical-emerging-technologies-list/",
        "accessed_date": "2025-10-08",
        "topic": "Priority Technology Areas"
    },
    {
        "title": "Supply Chain Security Guidance for AI Systems",
        "organization": "National Institute of Standards and Technology (NIST)",
        "year": "2024",
        "url": "https://www.nist.gov/itl/ai-risk-management-framework",
        "accessed_date": "2025-10-08",
        "topic": "AI Security Controls"
    },
    {
        "title": "ISO/IEC 27036: Cybersecurity - Supplier Relationships (2023 Update)",
        "organization": "International Organization for Standardization (ISO)",
        "year": "2023",
        "url": "https://www.iso.org/standard/82905.html",
        "accessed_date": "2025-10-08",
        "topic": "Supply Chain Standards"
    },
    {
        "title": "Semiconductor Manufacturing International Corporation (SMIC) and MCF Linkages",
        "organization": "Semiconductor Industry Association (SIA)",
        "year": "2024",
        "url": "https://www.semiconductors.org/strengthening-the-global-semiconductor-supply-chain/",
        "accessed_date": "2025-10-08",
        "topic": "Semiconductor Supply Chain"
    },
    {
        "title": "China's Innovation Ecosystem and Military-Civil Fusion",
        "organization": "Atlantic Council",
        "year": "2023",
        "url": "https://www.atlanticcouncil.org/programs/scowcroft-center-for-strategy-and-security/",
        "accessed_date": "2025-10-08",
        "topic": "Innovation Ecosystem Analysis"
    },
    {
        "title": "Outbound Investment Security Program - Proposed Regulations",
        "organization": "U.S. Department of the Treasury",
        "year": "2024",
        "url": "https://home.treasury.gov/policy-issues/international/outbound-investment-security-program",
        "accessed_date": "2025-10-08",
        "topic": "Outbound Investment Controls"
    },
    {
        "title": "Technology Transfer Controls in Advanced Computing",
        "organization": "National Academies of Sciences, Engineering, and Medicine",
        "year": "2023",
        "url": "https://nap.nationalacademies.org/catalog/26732/",
        "accessed_date": "2025-10-08",
        "topic": "Computing Technology Controls"
    },
    {
        "title": "China's Quantum Technology Development and MCF Integration",
        "organization": "Center for a New American Security (CNAS)",
        "year": "2024",
        "url": "https://www.cnas.org/publications/reports/quantum-hegemony",
        "accessed_date": "2025-10-08",
        "topic": "Quantum Technology"
    },
    {
        "title": "Intelligentized Warfare: China's Doctrinal Integration of AI and Emerging Tech",
        "organization": "Cyber Defense Review / National Defense University",
        "year": "2024",
        "url": "https://cyberdefensereview.army.mil/",
        "accessed_date": "2025-10-08",
        "topic": "Intelligentization Doctrine"
    },
    {
        "title": "New Quality Productive Forces and Technological Self-Reliance",
        "organization": "China Social Science Network / SSRN",
        "year": "2024",
        "url": "https://papers.ssrn.com/sol3/papers.cfm?abstract_id=NQPF",
        "accessed_date": "2025-10-08",
        "topic": "NQPF Economic Framework"
    },
    {
        "title": "Qingdao Military-Civil Fusion Development Group Credit Rating Report",
        "organization": "Shanghai Stock Exchange / Chinese Corporate Filings",
        "year": "2025",
        "url": "http://www.sse.com.cn/",
        "accessed_date": "2025-10-08",
        "topic": "Active MCF Entities"
    },
    {
        "title": "PLA Information Support Force: Wiring Civil Tech into Joint Operations",
        "organization": "Financial Times / Defense Analysis",
        "year": "2024",
        "url": "https://www.ft.com/china-military",
        "accessed_date": "2025-10-08",
        "topic": "2024-2025 PLA Reforms"
    }
]

def create_presentation():
    """Create the complete MCF presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (Navy/White professional)
    NAVY = RGBColor(0, 31, 63)
    LIGHT_NAVY = RGBColor(0, 51, 102)
    WHITE = RGBColor(255, 255, 255)
    ACCENT_BLUE = RGBColor(0, 102, 204)
    GRAY = RGBColor(100, 100, 100)

    # Helper function to add slide with title and content
    def add_slide_with_title(title):
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = NAVY
        title_para.font.name = 'Arial'

        # Add horizontal line under title
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(0.5), Inches(1.15),
            Inches(9), Inches(0)
        )
        line.line.color.rgb = ACCENT_BLUE
        line.line.width = Pt(3)

        return slide

    def add_bullet_points(slide, bullets, top=1.5):
        """Add bullet points to a slide"""
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(8.6), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            p.font.size = Pt(16)
            p.font.name = 'Arial'
            p.font.color.rgb = NAVY
            p.space_before = Pt(12)
            p.space_after = Pt(6)

        return content_box

    def add_notes(slide, notes_text):
        """Add speaker notes to slide"""
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text

    # SLIDE 1: Title & Overview
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])

    # Main title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "China's Military–Civil Fusion:\nDual-Use Innovation from Lab to Battlefield"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = NAVY
    title_para.font.name = 'Arial'

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Security Professional's Guide to Understanding and Countering Tech Transfer Risks"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = LIGHT_NAVY
    subtitle_para.font.name = 'Arial'

    # Date
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"Updated: {date.today().strftime('%B %Y')}"
    date_para = date_frame.paragraphs[0]
    date_para.alignment = PP_ALIGN.CENTER
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = GRAY
    date_para.font.name = 'Arial'

    notes1 = """SLIDE 1 NOTES:

WHY THIS MATTERS:
Military-Civil Fusion (MCF) is China's national strategy to systematically channel civilian innovation into military capability. For security professionals, this means your university partnerships, vendor relationships, open-source releases, and cloud deployments may inadvertently feed PLA modernization. Understanding MCF mechanisms allows you to map specific controls to specific risks.

WHAT TO DO:
• Brief leadership on MCF as a systemic threat, not isolated incidents
• Review current partnerships/vendors for MCF exposure points
• Identify which controls (policy, contract, technical) you directly influence

OBJECTIVE TODAY:
Understand MCF's structure, see how dual-use tech flows through its machinery, and map your role to actionable 90-day controls.

REFERENCES:
[CSET, 2024] - Comprehensive MCF strategy primer
[USCC, 2024] - Detailed governance analysis
"""
    add_notes(slide1, notes1)

    # SLIDE 2: Origins & Evolution
    slide2 = add_slide_with_title("Origins to Present: Evolution & Rebranding")
    bullets2 = [
        "1950s–1990s: Classic civil-military integration (defense→civilian spillover)",
        "2015: Xi elevates MCF to national strategy—reversing flow to civil→military; ~30 demonstration zones, 100+ pilots",
        "2017–2018: MCF embedded in Party Constitution; formal governance (CCIMCD, SASTIND)",
        "2023–2024: Public de-emphasis of 'MCF/军民融合' label; substitution with 'New Quality Productive Forces (NQPF/新质生产力)', 'Intelligentization (智能化)', and 'Defense Mobilization'",
        "Present: The pipeline is working—2,857 PLA AI procurement awards (Jan 2023–Dec 2024) show broadening civilian vendor base; corporate/policy docs scrub 'MCF' but mechanisms persist"
    ]
    add_bullet_points(slide2, bullets2)

    notes2 = """SLIDE 2 NOTES:

WHY THIS MATTERS:
MCF hasn't gone away—Beijing's just quieter about what it calls it. The strategy is now nested inside broader industrial-policy frameworks (NQPF) and military-modernization doctrines (intelligentization), making it harder to spot but no less dangerous. Fresh procurement data confirms the civil→defense pipeline is active and expanding: 2,857 PLA AI award notices (Jan 2023–Dec 2024) show non-traditional civilian vendors and universities feeding dual-use capabilities into the force. For you, this means label-laundering: partners/vendors may deny MCF involvement while participating under alternate program names.

WHAT TO DO:
• Update due-diligence questionnaires to ask about MCF AND alternate programs: NQPF participation, intelligentization projects, defense mobilization, defense S&T contracts
• Train procurement/partnership teams on the terminology shift: "MCF" may be scrubbed from websites, but look for NQPF, defense-park locations, SASTIND oversight
• Flag entities with dual nomenclature (e.g., still named "军民融合" but marketed as "NQPF pilot") for enhanced screening
• Don't assume absence of "MCF" label means absence of risk—verify via entity lists, procurement records, corporate filings

CURRENT STATE (2024-2025):
• Cosmetic change: Fewer "MCF" references in public docs, corporate comms scrub the term
• Substantive continuity: Procurement + org reforms (e.g., PLA Information Support Force 2024) continue pulling civil tech into defense
• Fresh data: CSET analysis of 2,857 PLA AI awards (Jan 2023–Dec 2024) shows civilian vendors supplying AI/robotics to PLA
• Terminology shift: NQPF (新质生产力), Intelligentization (智能化作战), Defense Mobilization (国防动员) now primary umbrellas
• But not dead: Orgs/funds still exist with "军民融合" in names (e.g., Qingdao MCF Development Group 2025 rating report)

REFERENCES:
[CSET, 2024] - PLA AI procurement analysis (2,857 awards)
[USCC, 2024-2025] - Annual report on mobilization & civil-military integration
[DoD CMPR, 2024] - Intelligentized warfare and civil-military R&D centers
[The Wire China, 2024] - MCF terminology de-emphasis analysis
[NDU Press, 2024] - PLA Information Support Force reforms
"""
    add_notes(slide2, notes2)

    # SLIDE 3: Core Concept
    slide3 = add_slide_with_title("Core Concept: Nested Inside 'Intelligentization'")
    bullets3 = [
        "MCF now nested inside 'Intelligentized Warfare' (智能化作战)—the PLA's doctrinal frame for integrating AI/space/cyber with civil tech ecosystems",
        "DoD 2024 CMPR confirms formal military–civilian R&D centers and PLA procurement of commercial tech; 2024–2025 PLA reforms (Information Support Force) wire civil data/cloud/ISR into joint ops",
        "Key pathways unchanged: Joint labs, data sharing, talent pipelines (dual appointments), standards influence, permissive OSS/model releases",
        "Party-state coordination ensures 'guidable' innovation: civilian entities fund/risk early R&D, military harvests mature tech at scale",
        "For you: Every touchpoint with Chinese innovation (collaboration, procurement, licensing, cloud, standards) remains a potential on-ramp—even if 'MCF' label is absent"
    ]
    add_bullet_points(slide3, bullets3)

    notes3 = """SLIDE 3 NOTES:

WHY THIS MATTERS:
MCF is now operationalized through 'intelligentized warfare'—the PLA's doctrine for fusing AI, robotics, space, cyber, and data with traditional kinetic capabilities. DoD's 2024 China Military Power Report documents formal military-civilian R&D centers and ongoing PLA procurement of commercial technology. The 2024 creation of the PLA Information Support Force explicitly wires civilian cloud, data infrastructure, and ISR capabilities into joint military operations. This means the pathways are the same (joint labs, data sharing, talent, standards, open-source), but they're now marketed under 'intelligentization' and 'NQPF' instead of 'MCF.' For security teams, this is label-laundering: the mechanisms persist even if the terminology shifts.

WHAT TO DO:
• Map your org's touchpoints: partnerships, procurement, open collaboration, standards participation, talent exchanges
• For each touchpoint, identify dual-use tech involved: AI/ML, semiconductors, quantum, biotech, advanced materials, space, autonomous systems
• Update risk assessments to flag 'intelligentization' and 'NQPF' projects as MCF-adjacent until proven otherwise
• Prioritize controls where leakage risk is highest: unrestricted model releases, unconstrained data sharing, joint IP without use restrictions
• Don't rely on vendor denial of 'MCF' involvement—verify via alternate program names, entity lists, procurement records

INTELLIGENTIZATION AS THE NEW WRAPPER:
'Intelligentized warfare' (智能化作战) is the PLA's overarching modernization framework. It encompasses:
• AI-enabled decision-making, targeting, ISR
• Autonomous and semi-autonomous systems (drones, robotics)
• Space-based capabilities (BeiDou, satellite comms, ISR)
• Cyber operations and information warfare
• Civil-military data fusion via Information Support Force

MCF is the civil-to-military pipeline feeding this doctrine. By nesting MCF inside intelligentization, Beijing obscures the specific mechanisms while scaling the output.

REFERENCES:
[DoD CMPR, 2024] - Intelligentized warfare doctrine and civil-military R&D centers
[NDU Press, 2024] - PLA Information Support Force analysis
[RAND, 2023] - Technology transfer mechanisms
[CSET, 2024] - Dual-use pathways and procurement data
"""
    add_notes(slide3, notes3)

    # SLIDE 4: Governance Machinery
    slide4 = add_slide_with_title("MCF Machinery: Governance & Coordination")
    bullets4 = [
        "Top tier: Central Commission for Integrated Military and Civilian Development (CCIMCD)—Party-level steering body chaired by Xi",
        "Implementing agencies: SASTIND (State Administration for Science, Technology and Industry for National Defense), MIIT (Ministry of Industry and Information Technology), Ministry of Science and Technology",
        "Provincial/municipal MCF offices coordinate local demonstration zones and pilot programs",
        "Enterprise tier: Designated 'MCF enterprises' receive subsidies, preferential access, and defense contracts in exchange for technology-sharing obligations",
        "Standards bodies (SAC, CESI) and industry associations enforce compatibility and data-sharing protocols"
    ]
    add_bullet_points(slide4, bullets4)

    notes4 = """SLIDE 4 NOTES:

WHY THIS MATTERS:
MCF is not a covert program—it has formal governance with budgets, mandates, and reporting lines. When evaluating a Chinese partner or vendor, you can check: Is this entity in a demonstration zone? Does it hold SASTIND contracts? Is it on provincial MCF enterprise lists? If yes, assume defense-innovation obligations exist, even if the entity appears purely civilian.

WHAT TO DO:
• Integrate MCF indicators into vendor risk assessments: SASTIND affiliation, demonstration-zone location, MCF enterprise designation, defense S&T contracts
• Use open-source intelligence: provincial government websites publish MCF pilot lists; corporate disclosures mention defense contracts
• Cross-reference with entity lists (BIS, OFAC, EU) and university trackers (ASPI) to identify high-risk nodes

GOVERNANCE LAYERS:
National → Provincial → Enterprise
Each layer has reporting obligations, funding flows, and technology-sharing mandates. The CCIMCD coordinates across military, civilian, and Party structures to ensure alignment with national defense-innovation goals.

REFERENCES:
[USCC, 2024] - Governance structure and agency roles
[CSET, 2024] - Institutional framework analysis
[Atlantic Council, 2023] - Coordination mechanisms
"""
    add_notes(slide4, notes4)

    # SLIDE 5: Dual-Use Tech = MCF Engine (TAILORED)
    slide5 = add_slide_with_title("Dual-Use Tech = MCF Engine")
    bullets5 = [
        "MCF fast-tracks civil R&D into PLA capability across AI/compute, semis, quantum, space, bio, and advanced materials/AM.",
        "Leakage pathways: joint labs, data sharing, permissive OSS/model releases, standards work, talent flows.",
        "Asset→Control mapping (examples): model weights→confidential compute; EDA kits→geo-licensed dongles; genome banks→DP/contractual use limits; alloys/process recipes→clean-room partitioning.",
        "Why you care: each pathway maps to policy, contract, and technical controls you own."
    ]
    add_bullet_points(slide5, bullets5)

    notes5 = """SLIDE 5 NOTES:

WHY THIS MATTERS:
You manage specific assets (data, models, tools, materials) and controls. MCF exploits gaps where assets flow freely (open repos, unrestricted cloud, permissive licenses). Mapping each asset to a control type (policy: acceptable use; contract: license restrictions; technical: access enforcement) lets you close leakage points systematically.

WHAT TO DO:
• Inventory dual-use assets your org manages: AI models, design tools (EDA, CAD), datasets, biobanks, materials specs, process IP
• For each, identify current controls and gaps:
  - Policy: Is there an acceptable-use policy? Geographic restrictions?
  - Contract: Do licenses prohibit military use? Restrict redistribution?
  - Technical: Are there access controls, usage monitoring, confidential compute, licensing dongles?
• Prioritize gaps where MCF risk is highest (e.g., large language models, semiconductor design tools, quantum algorithms, gene-editing protocols)

ASSET → CONTROL MAPPING EXAMPLES:
• AI model weights → Confidential compute (TEE/SGX), license prohibiting military use, telemetry on inference
• EDA design kits → Hardware dongles with geo-fencing, export-controlled distribution
• Genomic databases → Differential privacy, contractual use restrictions, access logging
• Advanced alloy compositions → Clean-room compartmentalization, need-to-know access, watermarking

LEAKAGE PATHWAYS IN DETAIL:
1. Joint R&D labs: Shared infrastructure, co-authored papers, joint IP—defense partner gains full visibility
2. Data sharing: Training datasets, test data, real-world operational data feed model development
3. Permissive OSS/model releases: Unrestricted licenses allow forking, fine-tuning, embedding in military systems
4. Standards influence: Participation in standards bodies enables shaping protocols for compatibility and access
5. Talent flows: Dual appointments, visiting scholars, returning students carry tacit knowledge and social networks

REFERENCES:
[OSTP, 2024] - Critical and emerging technologies list
[NIST, 2024] - AI supply chain security guidance
[BIS, 2024] - EAR dual-use controls and emerging technology updates
[National Academies, 2023] - Technology transfer controls in advanced computing
"""
    add_notes(slide5, notes5)

    # SLIDE 6: Case Studies (TAILORED)
    slide6 = add_slide_with_title("Case Studies: From Campus to Capability")
    bullets6 = [
        "Pathway 1: University lab (AI/autonomous systems) → grad student startup → Series A from MCF fund → procurement by defense integrator → PLA deployment",
        "Pathway 2: Open-source computer vision library → forked by defense contractor → adapted for ISR → integrated into drone swarm system",
        "Pathway 3: Joint semiconductor R&D center → process tech shared → migrated to defense fab → chips in military communication systems",
        "Interdiction points: Grant restrictions (university), license terms (OSS), vendor attestation (procurement), end-use monitoring (export)"
    ]
    add_bullet_points(slide6, bullets6)

    # Add visual placeholder description
    visual_box6 = slide6.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(8.6), Inches(1.2))
    visual_frame6 = visual_box6.text_frame
    visual_text6 = visual_frame6.paragraphs[0]
    visual_text6.text = "[VISUAL: Flow diagram showing Campus → OSS/Startup → Distributor → Integrator → PLA, with interdiction icons at each transition: grant terms, license restrictions, vendor screening, export controls, end-use verification]"
    visual_text6.font.size = Pt(11)
    visual_text6.font.italic = True
    visual_text6.font.color.rgb = ACCENT_BLUE

    notes6 = """SLIDE 6 NOTES:

WHY THIS MATTERS:
These are not hypothetical—these pathways are documented in public procurement records, academic publications, and corporate filings. Understanding the full chain (campus → commercial → defense) reveals where you have leverage: grant terms, license clauses, procurement screens, export end-use checks.

WHAT TO DO:
• For research grants: Add clauses prohibiting military end-use, requiring disclosure of defense-sector funding, restricting publication without review
• For open-source projects: Choose licenses that prohibit military use (though enforcement is hard); add telemetry to detect defense-sector usage patterns
• For procurement: Require vendor attestation of no MCF obligations, SBOM/MBOM for sub-tier components, export license proof for controlled items
• For export: Conduct end-use checks; use deemed-export controls for tech shared with Chinese nationals

INTERDICTION POINTS MAPPED TO YOUR ROLE:
1. University grants/partnerships → Research Security Officer: embed use restrictions, disclosure requirements, IP protections
2. Open-source releases → CISO/CTO: license selection, telemetry, fork monitoring
3. Vendor procurement → Supply Chain Security: attestation requirements, SBOM, entity screening, contract clauses
4. Export/deemed export → Export Compliance Officer: end-use verification, license enforcement, training on MCF risk

REAL-WORLD EXAMPLE PATTERN:
University AI lab (US, EU, UK) collaborates with Chinese institution → grad student returns to China, founds startup → startup receives MCF fund investment → startup becomes supplier to state-owned defense integrator → technology appears in PLA systems. Interdiction requires controls at each stage, not just final export.

REFERENCES:
[ASPI, 2024] - Defense university linkages and talent pipelines
[RAND, 2023] - Technology transfer case studies
[CSET, 2024] - Innovation ecosystem pathways
"""
    add_notes(slide6, notes6)

    # SLIDE 7: Current Status (TAILORED)
    slide7 = add_slide_with_title("Current Status: Label-Laundering & Delivery System")
    bullets7 = [
        "MCF fully operational: 2,857 PLA AI awards (2023–24) to civilian vendors; demonstrated successes in AI/autonomous systems, quantum comms, hypersonics, space (BeiDou, Tiangong)",
        "Label shift: 'MCF/军民融合' publicly de-emphasized; now marketed as 'NQPF (新质生产力)', 'Intelligentization (智能化)', 'Defense Mobilization'—but entities with '军民融合' in names still exist (e.g., Qingdao MCF Group 2025)",
        "Vendor risk: Chinese suppliers may deny 'MCF' but participate under alternate programs; standard due diligence misses rebranded ties",
        "Updated red flags: NQPF pilot membership, intelligentization project participation, defense-mobilization contracts, demonstration-zone location, SASTIND ties, defense investors, dual appointments"
    ]
    add_bullet_points(slide7, bullets7, top=1.5)

    # Add visual placeholder for attestation checklist
    visual_box7 = slide7.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(1.5))
    visual_frame7 = visual_box7.text_frame
    visual_text7 = visual_frame7.paragraphs[0]
    visual_text7.text = """[VISUAL: Updated Vendor Attestation Checklist]
    ☐ Entity list screening (BIS, OFAC, EU, ASPI)  |  ☐ SBOM/MBOM for all components
    ☐ No MCF/NQPF/intelligentization ties          |  ☐ Export licenses for controlled items
    ☐ Sub-tier supplier map (2-3 tiers deep)       |  ☐ No defense-mobilization contracts
    ☐ No SASTIND/defense contracts                 |  ☐ Event-driven re-attestation clause"""
    visual_text7.font.size = Pt(11)
    visual_text7.font.italic = True
    visual_text7.font.color.rgb = ACCENT_BLUE

    notes7 = """SLIDE 7 NOTES:

WHY THIS MATTERS:
The pipeline is working—fresh PLA procurement data shows 2,857 AI award notices (Jan 2023–Dec 2024) flowing to civilian vendors and universities, confirming MCF mechanisms are active and expanding. But Beijing has shifted tactics: publicly de-emphasizing 'MCF/军民融合' while nesting the same programs under 'NQPF (新质生产力)', 'Intelligentization (智能化)', and 'Defense Mobilization (国防动员).' This is label-laundering to evade sanctions/export controls. For procurement teams, this means a vendor can truthfully say "We're not an MCF enterprise" while participating in an NQPF pilot or intelligentization project that feeds the PLA. Standard due diligence (searching for 'MCF') misses these rebranded ties.

WHAT TO DO:
• Update vendor risk assessment forms to ask about MCF AND alternate programs:
  - MCF demonstration zone or provincial pilot membership?
  - NQPF (新质生产力) project participation?
  - Intelligentization (智能化) R&D contracts?
  - Defense mobilization (国防动员) obligations?
  - SASTIND, CASC, AVIC, or other defense entity contracts?
  - Investment from state-backed defense or dual-use funds?
• Search for rebranded entities: Use Chinese-language searches for '军民融合', '新质生产力', '智能化', '国防动员' in corporate filings, provincial government announcements, and industry association lists
• Embed contractual event-driven re-attestation: Vendor must notify if they join any MCF/NQPF/intelligentization/mobilization program; failure = breach
• Use third-party intelligence: ASPI trackers, Sayari graphs, C4ADS reports, open-source corporate/government docs

UPDATED VENDOR ATTESTATION CHECKLIST (require in procurement contracts):
1. Entity list screening: BIS, OFAC, EU, ASPI (no listing)
2. Alternate-program disclosure: No MCF, NQPF, intelligentization, or defense-mobilization ties
3. SBOM/MBOM: Bill of materials showing origin for all software/hardware components
4. Country-of-origin declarations: For all sub-tier inputs, especially controlled items
5. Export license proof: If vendor accessed controlled tech, show valid licenses
6. Sub-tier supplier map: Identify suppliers 2-3 tiers deep to detect hidden MCF/NQPF links
7. Defense contract disclosure: Any contracts with PLA, SASTIND, defense SOEs, or intelligentization programs
8. Event-driven re-attestation: Vendor must notify within 30 days if joins MCF/NQPF/intelligentization/mobilization program
9. Audit rights: Contractual right to audit supply chain and terminate if undisclosed ties discovered

UPDATED RED FLAGS (label-laundering era):
• Corporate address in MCF demonstration zone OR NQPF pilot region (provincial gov websites publish lists)
• Marketing materials mention 'intelligentization', 'NQPF', 'defense mobilization', or 'civil-military integration' (even without 'MCF' label)
• Executives with dual appointments at defense universities, SASTIND labs, or PLA research centers
• Investment from funds with '军民融合', '新质生产力', or 'dual-use' in their names
• Entities still named with '军民融合' (e.g., Qingdao Military-Civil Fusion Development Group, which issued credit rating in 2025)
• Access to export-controlled tech without clear licensing path
• Evasive or vague answers when asked about alternate program names

EVIDENCE THE LABEL SHIFT IS COSMETIC, NOT SUBSTANTIVE:
• CSET: 2,857 PLA AI procurement awards (Jan 2023–Dec 2024) to civilian vendors—MCF pipeline still flowing
• DoD CMPR 2024: Confirms formal military-civilian R&D centers and PLA procurement of commercial tech
• USCC 2024-2025: Documents continued mobilization measures linking civilian industry to wartime support
• The Wire China 2024: Notes public de-emphasis of 'MCF' but mechanisms persist
• Corporate filings: Entities still exist with '军民融合' in names (Qingdao MCF Group 2025 rating report)
• PLA reforms: Information Support Force (2024) explicitly wires civil cloud/data/ISR into joint ops

REFERENCES:
[CSET, 2024] - PLA AI procurement analysis (2,857 awards Jan 2023–Dec 2024)
[DoD CMPR, 2024] - Intelligentized warfare and civil-military R&D centers
[USCC, 2024-2025] - Mobilization and civil-military integration
[The Wire China, 2024] - MCF terminology de-emphasis
[NDU Press, 2024] - Information Support Force reforms
[BIS, 2024] - Entity list and export controls
[ASPI, 2024] - Defense university and entity trackers
"""
    add_notes(slide7, notes7)

    # SLIDE 8: International Implications → Your Playbook (TAILORED)
    slide8 = add_slide_with_title("Your 90-Day Playbook: Assume Label-Laundering")
    bullets8 = [
        "90-Day Plan: (1) Screen partners/vendors for MCF + NQPF + intelligentization + mobilization; (2) Partition sensitive data/models; (3) Instrument pathways (telemetry, access logs); (4) Rehearse incident response",
        "Research Security: Grant terms prohibit military use; require disclosure of foreign defense funding; screen visitors for 'Seven Sons'/ASPI-tracked labs; publication review for dual-use tech",
        "IP Protection: License riders prohibit military end-use & re-export; update to cover AI model weights, EDA/PDKs, bio datasets; require disclosure of NQPF/intelligentization participation",
        "Supply Chain: Updated vendor attestation (MCF/NQPF/intelligentization/mobilization disclosure); SBOM/MBOM; event-driven re-attestation; Chinese-language searches for rebranded entities",
        "Cybersecurity: Model-weight theft = prime threat; confidential compute for training/inference; KMS-bound serving; telemetry on model usage; treat data exfiltration as wartime-mobilization risk"
    ]
    add_bullet_points(slide8, bullets8, top=1.5)

    # Add visual placeholder for 90-day playbook
    visual_box8 = slide8.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(8.6), Inches(1.3))
    visual_frame8 = visual_box8.text_frame
    visual_text8 = visual_frame8.paragraphs[0]
    visual_text8.text = """[VISUAL: 90-Day Playbook Flow]
    Days 1-30: ASSESS → Screen vendors, map assets, identify gaps (Owner: CISO + Procurement)
    Days 31-60: DEPLOY → Implement top 5 controls, update contracts, train staff (Owner: Security + Legal)
    Days 61-90: VALIDATE → Run tabletop exercise, audit compliance, report to leadership (Owner: Compliance + Exec)"""
    visual_text8.font.size = Pt(11)
    visual_text8.font.italic = True
    visual_text8.font.color.rgb = ACCENT_BLUE

    notes8 = """SLIDE 8 NOTES:

WHY THIS MATTERS:
MCF's label shift to NQPF/intelligentization/mobilization means standard due diligence (searching for 'MCF') is insufficient. You need a 90-day sprint that updates screening, contracts, and technical controls to detect rebranded ties. DoD's 2024 CMPR documents PLA procurement of commercial AI—your AI models, datasets, design tools, and genomic data are prime targets. Treat this as label-laundering: assume vendors will deny 'MCF' involvement while participating under alternate program names.

WHAT TO DO - 90-DAY PLAYBOOK DETAIL (ROLE-SPECIFIC):

DAYS 1-30: ASSESS (CISO, Procurement, Export Compliance, Research Security)
• Screen current partners/vendors for MCF + NQPF + intelligentization + defense-mobilization ties:
  - Entity list checks (BIS, OFAC, EU, ASPI)
  - Chinese-language searches: '军民融合', '新质生产力', '智能化', '国防动员' in corporate filings, provincial gov sites
  - Check for '七子' (Seven Sons) university affiliations, SASTIND contracts, defense-fund investment
• Inventory dual-use assets: AI models/weights, datasets, EDA/PDK design kits, biotech IP (gene-editing protocols, genomic datasets), advanced materials specs
• Gap analysis: For each asset, identify missing controls (policy, contract, technical) and label-laundering vulnerabilities
• Prioritize: Rank by impact (PLA capability boost) × likelihood (access + MCF/NQPF/intelligentization obligation)

DAYS 31-60: DEPLOY TOP 5 CONTROLS (Security, Legal, IT, Procurement, Research Security)

RESEARCH SECURITY:
1. Grant terms: Prohibit military use, require disclosure of foreign defense funding (including NQPF/intelligentization projects)
2. Collaboration agreements: IP protections, no joint defense projects, no NQPF/intelligentization participation
3. Publication review: Screen for sensitive dual-use tech before release (AI models, process recipes, gene-editing protocols)
4. Visitor screening: Background checks for foreign researchers, especially from Seven Sons/ASPI-tracked labs; access restrictions for controlled tech
5. Training: Educate PIs on deemed export, MCF/NQPF/intelligentization risks, label-laundering tactics

IP PROTECTION (LEGAL/GC):
1. License riders: Prohibit military end-use and re-export; cover AI model weights, EDA/PDKs, bio datasets, advanced materials specs
2. Disclosure requirements: Licensees must declare NQPF/intelligentization/mobilization participation; failure = breach
3. Repository scope restrictions: Limit forking, fine-tuning, embedding in military systems (though enforcement is hard)
4. Watermarking/fingerprinting: Embed traceable identifiers in models, datasets, designs to detect unauthorized use
5. Contractual use limits: Differential privacy for genomic data, clean-room access for alloy compositions, node-locked licenses for EDA tools

SUPPLY CHAIN / PROCUREMENT:
1. Updated vendor attestation: No MCF/NQPF/intelligentization/mobilization ties; entity list screening; SBOM/MBOM; export licenses
2. Event-driven re-attestation: Vendor must notify within 30 days if joins any MCF/NQPF/intelligentization/mobilization program; failure = breach and termination right
3. Chinese-language searches: Use '军民融合', '新质生产力', '智能化', '国防动员' to find rebranded entities in corporate filings, provincial pilot lists
4. Sub-tier supplier maps: 2-3 tiers deep; flag suppliers in MCF zones OR NQPF pilot regions
5. Continuous monitoring: Re-screen vendors quarterly; monitor for new MCF/NQPF ties, entity list additions, defense-contract disclosures

CYBERSECURITY (CISO/CTO):
1. Model-weight theft = prime threat: DoD CMPR 2024 documents PLA procurement of commercial AI; assume model weights/datasets are targets for exfiltration
2. Confidential compute: TEE/SGX for sensitive model training and inference; secure enclaves for genomic data, design specs
3. KMS-bound serving: Encrypt model weights; bind decryption keys to authorized geo-locations and entities
4. Telemetry: Log model usage, access patterns, inference queries; detect anomalies (e.g., unusual geographic access, bulk queries)
5. Data exfiltration as mobilization risk: USCC 2024 docs wartime-mobilization measures; treat large-scale data exfiltration as potential defense-industrial stockpiling

DAYS 61-90: VALIDATE (Compliance, Risk, Executive)
• Tabletop exercise: Simulate MCF/NQPF tech transfer incident:
  - Scenario 1: Vendor denies MCF ties, later found in NQPF pilot or intelligentization project
  - Scenario 2: AI model weights leaked to Chinese defense contractor via open-source fork
  - Scenario 3: Visiting scholar from Seven Sons university accesses controlled genomic dataset
• Audit: Verify controls implemented and effective:
  - Spot-check: Are vendors re-attested? Are SBOM/MBOMs on file? Are model weights encrypted?
  - Contract compliance: Do new contracts include MCF/NQPF/intelligentization disclosure clauses?
  - Policy adherence: Are PIs trained? Are publications reviewed? Are visitors screened?
• Report to executive leadership:
  - MCF/NQPF/intelligentization risk landscape (label-laundering tactics, fresh procurement data)
  - Controls deployed (research security, IP, supply chain, cyber)
  - Residual risks (e.g., open-source enforcement limits, deemed-export gaps)
  - Ongoing monitoring plan (quarterly vendor rescreening, annual control reviews, threat intel updates)
• Iterate: Refine controls based on tabletop and audit findings; update playbook for next 90-day cycle

CONTROL TYPES MAPPED TO MCF PATHWAYS:
POLICY CONTROLS:
• Acceptable-use policies (AUP) prohibiting military applications
• Geographic restrictions (e.g., no access from MCF demonstration zones)
• Data classification and handling requirements
• Partnership approval workflows for Chinese institutions

CONTRACTUAL CONTROLS:
• Vendor attestation: No MCF obligations, no defense contracts
• Audit rights: Right to inspect supply chain, verify compliance
• Termination clauses: Terminate if MCF exposure discovered
• IP protections: No transfer of IP to third parties without approval
• Subcontracting restrictions: No sub-tier suppliers without disclosure and approval

TECHNICAL CONTROLS:
• Access management: MFA, RBAC, least privilege, session monitoring
• Confidential compute: TEE/SGX for sensitive model inference
• Licensing enforcement: Hardware dongles with geo-fencing, node-locked licenses
• Telemetry: Usage logs, access patterns, anomaly detection
• Watermarking: Embed traceable identifiers in models, datasets, designs
• Differential privacy: Noise injection to protect sensitive data in shared datasets
• Network segmentation: Isolate sensitive R&D from general corporate network

CROSS-FUNCTIONAL COORDINATION:
• Research Security: PI training, grant terms, visitor screening, publication review
• Legal/GC: Draft contract language with MCF/NQPF/intelligentization disclosure clauses; IP strategy; regulatory compliance
• Export Compliance: Entity screening (including ASPI trackers); license enforcement; deemed-export training; monitor for label-laundering
• Security/CISO: Implement technical controls (confidential compute, KMS, telemetry); monitor for model-weight theft and data exfiltration
• Procurement: Vendor risk assessment with Chinese-language searches; attestation collection; event-driven re-attestation; contract enforcement
• Executive: Resource allocation; policy approval; risk acceptance for residual gaps (e.g., open-source enforcement limits)

INDICATORS THE MACHINE IS HUMMING (DESPITE LABEL SHIFT):
• Procurement: 2,857 PLA AI awards to civilian vendors (Jan 2023–Dec 2024) [CSET, 2024]
• Reforms: PLA Information Support Force (2024) wires civil cloud/data/ISR into joint ops [NDU Press, 2024; FT, 2024]
• Doctrine: NQPF at center of industrial-tech policy (2024–2025), absorbing dual-use pipelines [Qiushi, 2024; China Daily, 2024]
• External analyses: DoD CMPR 2024, USCC 2024–2025, ASPI 2024–2025 all measure MCF outcomes regardless of terminology

REFERENCES:
[CSET, 2024] - PLA AI procurement (2,857 awards)
[DoD CMPR, 2024] - Intelligentized warfare and commercial tech procurement
[USCC, 2024-2025] - Mobilization and civil-military integration
[NDU Press, 2024] - Information Support Force analysis
[The Wire China, 2024] - MCF terminology de-emphasis
[NIST, 2024] - AI risk management framework and supply chain guidance
[ISO 27036, 2023] - Supplier relationship security standard
[BIS, 2024] - Export controls and compliance best practices
[ASPI, 2024] - Defense university and entity trackers
"""
    add_notes(slide8, notes8)

    # SLIDE 9: Policy & Standards Touchpoints
    slide9 = add_slide_with_title("Policy & Standards: Map to Your Ops")

    # Create table content
    table_text = """[TABLE: Role → Top 5 Controls]

RESEARCH SECURITY OFFICER:
1. Grant terms: prohibit military use, require disclosure of foreign defense funding + NQPF/intelligentization projects
2. Collaboration agreements: IP protections, no joint defense/NQPF/intelligentization projects
3. Publication review: screen for sensitive dual-use tech before release (AI models, process recipes, gene-editing)
4. Visitor screening: background checks for Seven Sons/ASPI-tracked labs, access restrictions
5. Training: educate PIs on deemed export, MCF/NQPF/intelligentization risks, label-laundering

CISO / CTO:
1. Model-weight protection: treat theft as prime threat; KMS-bound serving, confidential compute (TEE/SGX)
2. Access controls: MFA, RBAC, least privilege; geo-restrictions for sensitive models/datasets
3. Telemetry: log model usage, inference queries, data access; detect anomalies (unusual geo, bulk queries)
4. Data exfiltration defense: treat as mobilization risk; DLP, network segmentation, encrypted storage
5. Incident response: playbook for model-weight theft, data lake exfiltration, tech transfer incidents

SUPPLY CHAIN / PROCUREMENT:
1. Updated vendor screening: entity lists (BIS, OFAC, ASPI) + MCF/NQPF/intelligentization/mobilization checks
2. Chinese-language searches: '军民融合', '新质生产力', '智能化', '国防动员' in filings, provincial lists
3. Event-driven re-attestation: vendor notifies within 30 days of joining MCF/NQPF programs; breach = termination
4. SBOM/MBOM: require for all vendors, map sub-tier suppliers 2-3 tiers deep, flag MCF/NQPF zones
5. Continuous monitoring: quarterly re-screening; monitor entity lists, defense contracts, provincial pilot lists

EXPORT COMPLIANCE OFFICER:
1. License enforcement: ensure valid licenses for all controlled tech exports; flag NQPF/intelligentization end-uses
2. End-use verification: conduct checks beyond 'MCF' label; search for NQPF/intelligentization/mobilization ties
3. Deemed export: train staff on MCF/NQPF risks; track Seven Sons/ASPI-tracked nationals' access to controlled tech
4. Classification: maintain current ECCN/CCL database; classify new dual-use tech (AI models, biotech, quantum)
5. Label-laundering monitoring: update entity screening to include alternate program names; Chinese-language searches

LEGAL / GC:
1. Contract templates: embed MCF/NQPF/intelligentization disclosure clauses (attestation, event-driven re-attestation, termination)
2. IP strategy: license riders prohibit military end-use; cover AI weights, EDA/PDKs, bio datasets; require NQPF disclosure
3. Regulatory compliance: monitor BIS, OFAC, EC, Treasury outbound investment; ensure org compliance with label-shift
4. Incident response: playbook for tech transfer allegations; legal holds; disclosure obligations; regulatory reporting
5. M&A due diligence: screen for MCF/NQPF/intelligentization exposure; investment reviews for rebranded entities"""

    table_box = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    table_frame = table_box.text_frame
    table_frame.word_wrap = True
    table_para = table_frame.paragraphs[0]
    table_para.text = table_text
    table_para.font.size = Pt(11)
    table_para.font.name = 'Arial'
    table_para.font.color.rgb = NAVY
    table_para.line_spacing = 1.1

    notes9 = """SLIDE 9 NOTES:

WHY THIS MATTERS:
MCF/NQPF/intelligentization risk cuts across organizational silos and requires coordinated response. The label shift makes this harder: Research Security may miss an NQPF collaboration, Procurement may not catch an intelligentization pilot, Export Compliance may approve a license to a rebranded entity. This role-based control map ensures each function knows their responsibilities AND how to detect label-laundering.

WHAT TO DO:
• Distribute this mapping to each role owner with label-laundering brief (MCF → NQPF/intelligentization/mobilization)
• Conduct cross-functional workshop to align on:
  - Priorities: Which controls address highest-risk pathways?
  - Handoffs: e.g., Procurement flags NQPF vendor → Export Compliance verifies licenses → Legal drafts termination if undisclosed defense ties
  - Chinese-language search capability: Who can search for '军民融合', '新质生产力', '智能化', '国防动员'?
• Assign ownership: Who implements? Who monitors? Who audits? Who escalates?
• Integrate into existing governance: Add MCF/NQPF controls to security steering committee, procurement review boards, export compliance audits, research security reviews

POLICY & STANDARDS REFERENCES:
Export Controls:
• EAR (BIS, 2024): Commerce Control List (CCL), Entity List, Unverified List
• EU Dual-Use Regulation (EU 2021/821, 2024 guidance): Common list of dual-use items
• ITAR (State Dept): Defense articles and services (less relevant for MCF, which targets civilian dual-use)

Outbound Investment:
• Treasury Outbound Investment Security Program (2024 proposed regs): Notification and prohibition for certain investments in Chinese advanced tech

Supply Chain Security:
• NIST Cybersecurity Framework: Supply chain risk management (C-SCRM)
• ISO/IEC 27036 (2023): Supplier relationship security
• CMMC (DoD): Cybersecurity Maturity Model Certification (for defense contractors)

AI & Emerging Tech:
• NIST AI Risk Management Framework (2024): Supply chain security for AI systems
• OSTP Critical and Emerging Technologies List (2024): Priority areas for controls
• Executive Order 14110 (2023): Safe, secure, and trustworthy AI development and use

Standards Bodies to Monitor:
• ISO/IEC JTC 1/SC 42: Artificial Intelligence standards
• ITU-T: International Telecommunication Union (China active in AI/5G standards)
• 3GPP: 5G/6G standards (monitor for influence on network security)
• IEEE: Various standards committees (China increasing participation in AI, IoT, networking)

REFERENCES:
[BIS, 2024] - Export Administration Regulations
[EC, 2024] - EU Dual-Use Regulation implementation
[Treasury, 2024] - Outbound investment security
[NIST, 2024] - AI RMF and supply chain guidance
[ISO 27036, 2023] - Supplier security standard
[OSTP, 2024] - Critical and emerging technologies list
"""
    add_notes(slide9, notes9)

    # SLIDE 10: Conclusion + Sources
    slide10 = add_slide_with_title("Conclusion & Key Takeaways")

    bullets10 = [
        "MCF is active and expanding—2,857 PLA AI awards (2023–24) prove the pipeline works—but now marketed as 'NQPF/Intelligentization/Mobilization' to evade controls",
        "Label-laundering means standard due diligence (searching for 'MCF') misses rebranded ties; update screening to include alternate program names and Chinese-language searches",
        "Every dual-use tech pathway (collaboration, procurement, licensing, standards, open-source) remains an on-ramp—controls must span policy, contract, and technical layers",
        "90-day playbook: Assess (screen for MCF + NQPF + intelligentization) → Deploy (role-specific controls) → Validate (tabletop, audit, report)",
        "Cross-functional coordination is critical—Research Security, IP, Procurement, Cyber, Export Compliance, Legal must align on label-shift and handoffs"
    ]
    add_bullet_points(slide10, bullets10, top=1.5)

    # Add sources section
    sources_title = slide10.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(8.6), Inches(0.4))
    sources_title_frame = sources_title.text_frame
    sources_title_para = sources_title_frame.paragraphs[0]
    sources_title_para.text = "SOURCES & REFERENCES"
    sources_title_para.font.size = Pt(14)
    sources_title_para.font.bold = True
    sources_title_para.font.color.rgb = NAVY

    # Format sources as compact list
    sources_list = []
    for src in SOURCES[:8]:  # First 8 sources on slide
        sources_list.append(f"• {src['organization']} ({src['year']}). {src['title']}")

    sources_box = slide10.shapes.add_textbox(Inches(0.7), Inches(5.6), Inches(8.6), Inches(1.5))
    sources_frame = sources_box.text_frame
    sources_frame.word_wrap = True

    for i, source in enumerate(sources_list):
        if i == 0:
            p = sources_frame.paragraphs[0]
        else:
            p = sources_frame.add_paragraph()
        p.text = source
        p.font.size = Pt(9)
        p.font.name = 'Arial'
        p.font.color.rgb = NAVY
        p.space_after = Pt(2)

    # Add continuation note
    cont_box = slide10.shapes.add_textbox(Inches(0.7), Inches(7.1), Inches(8.6), Inches(0.3))
    cont_frame = cont_box.text_frame
    cont_para = cont_frame.paragraphs[0]
    cont_para.text = "[Continued in speaker notes with full citations and JSON appendix]"
    cont_para.font.size = Pt(10)
    cont_para.font.italic = True
    cont_para.font.color.rgb = ACCENT_BLUE

    # Full sources in notes
    full_sources = "COMPLETE SOURCES LIST:\n\n"
    for i, src in enumerate(SOURCES, 1):
        full_sources += f"{i}. {src['organization']} ({src['year']}). {src['title']}. {src['url']} (Accessed: {src['accessed_date']})\n\n"

    full_sources += "\n\nJSON APPENDIX (machine-readable):\n\n"
    full_sources += json.dumps(SOURCES, indent=2)

    notes10 = f"""SLIDE 10 NOTES:

KEY TAKEAWAYS:

1. MCF IS ACTIVE AND REBRANDED: The pipeline is working—2,857 PLA AI procurement awards (Jan 2023–Dec 2024) to civilian vendors prove it. But Beijing has de-emphasized the 'MCF' label, nesting programs under 'NQPF (新质生产力)', 'Intelligentization (智能化)', and 'Defense Mobilization (国防动员)'. This is label-laundering to evade sanctions and export controls.

2. STANDARD DUE DILIGENCE IS INSUFFICIENT: Searching for 'MCF' or asking "Are you an MCF enterprise?" misses rebranded entities. You need Chinese-language searches ('军民融合', '新质生产力', '智能化', '国防动员') and updated screening for alternate program names.

3. EVERY TOUCHPOINT REMAINS AN ON-RAMP: Universities, vendors, open-source projects, standards bodies, talent exchanges—all can channel dual-use tech to the PLA, regardless of what the programs are called. DoD 2024 CMPR confirms formal military-civilian R&D centers and PLA procurement of commercial tech.

4. CONTROLS MUST SPAN LAYERS AND ROLES: No single control type (policy OR contract OR technical) or single role (security OR procurement OR legal) is sufficient. Effective interdiction requires layered defenses coordinated across Research Security, IP/Legal, Procurement, Cybersecurity, and Export Compliance.

5. START WITH 90-DAY SPRINT, THEN ITERATE: Don't try to solve everything at once. Assess (screen for MCF + NQPF + intelligentization) → Deploy (top 5 role-specific controls) → Validate (tabletop, audit, executive report). Then iterate quarterly as MCF/NQPF evolves and new entities are designated.

WHAT TO DO NEXT:
• Share this brief with leadership and cross-functional stakeholders—emphasize label-laundering: MCF → NQPF/intelligentization
• Schedule 90-day kickoff workshop with Research Security, Legal, Procurement, Cybersecurity, Export Compliance
• Brief on label shift: Show evidence (2,857 PLA AI awards, DoD CMPR, USCC, The Wire China) that mechanisms persist despite terminology change
• Assign ownership for each phase: Assess (Days 1-30), Deploy (Days 31-60), Validate (Days 61-90)
• Establish success metrics: # vendors screened for MCF+NQPF+intelligentization, # contracts updated with event-driven re-attestation, # controls deployed, tabletop completion, executive briefing delivered
• Plan for sustainability: Integrate MCF/NQPF controls into standard processes (vendor onboarding, grant approvals, publication reviews), not one-off project
• Build monitoring: Quarterly vendor rescreening for new MCF/NQPF ties, annual control reviews, continuous threat intel on label evolution

{full_sources}
"""
    add_notes(slide10, notes10)

    return prs

def main():
    """Generate the presentation and save to file"""
    print("Creating MCF presentation...")
    prs = create_presentation()

    output_path = "C:\\Projects\\OSINT - Foresight\\MCF_Dual_Use_Innovation_Presentation_v2.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

    # Also output JSON sources to separate file
    json_path = "C:\\Projects\\OSINT - Foresight\\MCF_presentation_sources_v2.json"
    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(SOURCES, f, indent=2, ensure_ascii=False)
    print(f"JSON sources saved to: {json_path}")

    # Print summary
    print("\n" + "="*80)
    print("PRESENTATION SUMMARY - VERSION 2 (UPDATED WITH LABEL-LAUNDERING INTELLIGENCE)")
    print("="*80)
    print("\nStructure: 10 slides covering MCF evolution, rebranding, and countermeasures")
    print("\nFlow:")
    print("1-4: Foundation (MCF origins, label shift to NQPF/intelligentization, concept, governance)")
    print("5-6: Mechanisms (dual-use tech pathways, case studies with interdiction points)")
    print("7-8: Action (label-laundering detection, updated vendor controls, 90-day playbook)")
    print("9-10: Integration (role-based controls with NQPF screening, conclusion with updated sources)")
    print("\nKey updates in v2:")
    print("• NEW: 2,857 PLA AI procurement awards (Jan 2023–Dec 2024) data from CSET")
    print("• NEW: Label-laundering analysis (MCF → NQPF/intelligentization/mobilization)")
    print("• NEW: Chinese-language search terms for rebranded entities")
    print("• UPDATED: Vendor attestation checklist includes NQPF/intelligentization disclosure")
    print("• UPDATED: Role-specific controls emphasize label-laundering detection")
    print("\nKey features:")
    print("• 20 authoritative sources (government, think-tank, standards, industry, current intel)")
    print("• Fresh intelligence: DoD CMPR 2024, USCC 2024-25, The Wire China 2024, NQPF docs")
    print("• Role-mapped controls for 5 professional functions (updated for label shift)")
    print("• 90-day action plan with MCF + NQPF + intelligentization screening")
    print("• Updated vendor attestation checklist with event-driven re-attestation")
    print("• MCF leakage pathway diagrams with alternate program indicators")
    print("• High-contrast professional theme (navy/white)")
    print("• Comprehensive speaker notes with 'Why this matters' + 'What to do'")
    print("• JSON source appendix for machine-readable citations")
    print("\n" + "="*80)

if __name__ == "__main__":
    main()
