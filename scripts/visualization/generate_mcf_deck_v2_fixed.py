#!/usr/bin/env python3
"""
MCF Presentation Generator V2 - FIXED VERSION
Addresses all feedback:
- Proper labels instead of numbers
- Clear graphics
- Dark blue background
- White text
- Proper sizing
- Professional colors
"""

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Circle, Rectangle, Polygon, FancyBboxPatch, Wedge, RegularPolygon
from pathlib import Path
import numpy as np
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Setup
VIZ_DIR = Path("C:/Projects/OSINT - Foresight/scripts/visualization/visualizations/mcf_v2_fixed")
VIZ_DIR.mkdir(parents=True, exist_ok=True)

# UPDATED COLOR SCHEME
BG_COLOR = '#1E3A5F'  # Dark blue background
TEXT_COLOR = 'white'

COLORS = {
    'bg_dark_blue': '#1E3A5F',
    'text_white': 'white',
    'military_red': '#DC143C',
    'civilian_blue': '#87CEEB',
    'success_green': '#90EE90',
    'partial_yellow': '#FFD700',
    'failure_red': '#FF6B6B',
    'neutral_gray': '#B0B0B0',
    'highlight': '#FFA500',
}

print("="*80)
print("MCF PRESENTATION GENERATOR V2 - FIXED")
print("="*80)

def save_fig(filename):
    """Helper to save figure with dark blue background"""
    plt.savefig(VIZ_DIR / filename, dpi=300, bbox_inches='tight',
                facecolor=BG_COLOR, edgecolor='none')
    plt.close()
    print(f"  [OK] {filename}")

# ====================
# SLIDE 2: Honeycomb with TEXT LABELS (not numbers)
# ====================
print("\n[2/17] Creating honeycomb with text labels...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

def hexagon(center, size=1):
    angles = np.linspace(0, 2*np.pi, 7)
    return [(center[0] + size*np.cos(a), center[1] + size*np.sin(a)) for a in angles]

labels = [
    'Historical\nRoots',
    'Bidirectional\nFusion',
    'Strategic\nGoals',
    'Legal\nExpansion',
    'Institutional\nWeb',
    'Global\nReach',
    'Implications'
]

# Center hexagon
center_hex = hexagon((7, 4), 1.0)
hex_patch = Polygon(center_hex, facecolor='#4A6FA5', edgecolor='white', linewidth=3)
ax.add_patch(hex_patch)
ax.text(7, 4, labels[0], fontsize=16, ha='center', va='center',
       color='white', fontweight='bold')

# Surrounding 6 hexagons
positions = [(7, 6), (8.7, 5), (8.7, 3), (7, 2), (5.3, 3), (5.3, 5)]
for i, pos in enumerate(positions):
    hex_coords = hexagon(pos, 1.0)
    hex_patch = Polygon(hex_coords, facecolor='#4A6FA5', edgecolor='white', linewidth=3)
    ax.add_patch(hex_patch)
    ax.text(pos[0], pos[1], labels[i+1], fontsize=14, ha='center', va='center',
           color='white', fontweight='bold', linespacing=1.3)

ax.set_xlim(3.5, 10.5)
ax.set_ylim(0.5, 7.5)
ax.axis('off')
save_fig('slide02_honeycomb.png')

# ====================
# SLIDE 3: CLEARER Timeline
# ====================
print("\n[3/17] Creating clearer timeline...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

events = [
    ('1840\nOpium Wars', 0, 'Humiliation begins'),
    ('1978\nReform Era', 2, 'Opening up'),
    ('1986\n863 Program', 4, 'Tech focus'),
    ('1993\nYinhe Incident', 6, 'GPS denial'),
    ('2000\nIntegration', 8, 'Civil-Military'),
    ('2015\nMCF Strategy', 10, 'Formalized'),
    ('2024\nNQPF', 12, 'Expanded')
]

# Draw horizontal timeline
ax.plot([0, 13], [4, 4], color='white', linewidth=4, zorder=1)

for label, x, desc in events:
    # Marker circle
    circle = Circle((x, 4), 0.4, facecolor=COLORS['highlight'],
                   edgecolor='white', linewidth=3, zorder=3)
    ax.add_patch(circle)

    # Year label above
    ax.text(x, 5.5, label, fontsize=14, ha='center', va='bottom',
           color='white', fontweight='bold', linespacing=1.2)

    # Description below
    ax.text(x, 2.5, desc, fontsize=11, ha='center', va='top',
           color='white', style='italic')

ax.text(6.5, 7, 'From Dependence to Self-Reliance', fontsize=18,
       ha='center', color='white', fontweight='bold')

ax.set_xlim(-1, 14)
ax.set_ylim(1, 8)
ax.axis('off')
save_fig('slide03_timeline.png')

# ====================
# SLIDE 4: Gears WITHOUT numbers (clearer labels)
# ====================
print("\n[4/17] Creating gears with clear labels...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

gear_labels = ['Infrastructure', 'Industry', 'S&T', 'Education',
               'Social Services', 'Maritime', 'Emergency']
positions = [(3,4), (5,6), (7,6.5), (9,6), (11,4), (9,2), (7,1.5)]
colors_list = [COLORS['military_red']]*3 + [COLORS['civilian_blue']]*4

for i, (pos, label, color) in enumerate(zip(positions, gear_labels, colors_list)):
    circle = Circle(pos, 0.8, facecolor=color, edgecolor='white', linewidth=3, alpha=0.9)
    ax.add_patch(circle)
    # Label OUTSIDE circle
    ax.text(pos[0], pos[1]-1.3, label, fontsize=12, ha='center', va='top',
           color='white', fontweight='bold')

# Center bidirectional arrow
ax.annotate('', xy=(10, 4), xytext=(4, 4),
           arrowprops=dict(arrowstyle='<->', lw=4, color='white'))
ax.text(7, 4.5, 'BIDIRECTIONAL', fontsize=14, ha='center',
       color='white', fontweight='bold')

ax.set_xlim(1, 13)
ax.set_ylim(0, 8)
ax.axis('off')
save_fig('slide04_gears.png')

# Slide 5: Quadrant (already good, just update colors)
print("\n[5/17] Creating quadrant...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
quadrants = [
    (0, 4, 6, 4, '#8B0000', 'Military\nModernization'),
    (6, 4, 6, 4, '#228B22', 'Economic\nDevelopment'),
    (0, 0, 6, 4, '#FF8C00', 'Self-Reliance'),
    (6, 0, 6, 4, '#6B46C1', 'First-Mover\nAdvantage')
]
for x, y, w, h, color, label in quadrants:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=4, alpha=0.8)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, label, fontsize=20, ha='center', va='center',
           color='white', fontweight='bold', linespacing=1.5)
ax.set_xlim(-0.5, 12.5)
ax.set_ylim(-0.5, 8.5)
ax.axis('off')
save_fig('slide05_quadrant.png')

# Slide 6: Stairs (already good)
print("\n[6/17] Creating stairs...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
stairs = [
    (0, 0, 2, 1, 'MCF\nStrategy\n2015'),
    (2, 1, 2, 1, 'CMC S&T\n2016'),
    (4, 2, 2, 1, 'CCIMCD\n2017'),
    (6, 3, 2, 1, 'Constitutional\n2018'),
    (8, 4, 2, 1, 'Dual Circ\n2020'),
    (10, 5, 2, 1, '14th FYP\n2021'),
    (12, 6, 2, 1, 'NQPF\n2023')
]
for i, (x, y, w, h, label) in enumerate(stairs):
    intensity = 0.4 + (i * 0.08)
    color = plt.cm.Blues(intensity)
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=2)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, label, fontsize=11, ha='center', va='center',
           fontweight='bold', color='white', linespacing=1.2)
ax.arrow(13.5, 6.5, 0.3, 0.3, head_width=0.2, head_length=0.15,
        fc='white', ec='white', linewidth=3)
ax.set_xlim(-0.5, 14.5)
ax.set_ylim(-0.5, 7.5)
ax.axis('off')
save_fig('slide06_stairs.png')

# Slide 7: Spider (already good, update colors)
print("\n[7/17] Creating spider chart...")
fig = plt.figure(figsize=(14, 8), facecolor=BG_COLOR)
ax = fig.add_subplot(111, projection='polar', facecolor=BG_COLOR)
laws = ['2015\nNSL', '2017\nIntel', '2017\nCyber', '2020\nExport',
        '2021\nData', '2021\nDefense', '2023\nFR', '2024\nSecrets']
radii = [3, 4, 4, 5, 6, 6, 7, 8]
theta = np.linspace(0, 2*np.pi, 8, endpoint=False)
ax.plot(theta, radii, 'o-', linewidth=3, color=COLORS['failure_red'], markersize=10)
ax.fill(theta, radii, alpha=0.3, color=COLORS['failure_red'])
ax.set_xticks(theta)
ax.set_xticklabels(laws, fontsize=12, color='white')
ax.set_ylim(0, 10)
ax.tick_params(colors='white')
ax.spines['polar'].set_color('white')
ax.grid(True, color='white', alpha=0.3)
save_fig('slide07_spider.png')

# ====================
# SLIDE 8: Network with LARGER fonts, labels OUTSIDE bubbles
# ====================
print("\n[8/17] Creating network with larger fonts...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

G = nx.Graph()
G.add_node('Xi/CMC', size=2000, color='#8B0000')
for node in ['CCIMCD', 'State Council', 'CMC']:
    G.add_node(node, size=1200, color='#000080')
    G.add_edge('Xi/CMC', node)

ministries = ['MOST', 'MIIT', 'MOE', 'MSS', 'SASTIND', 'MOF', 'NDRC', 'UFWD']
for m in ministries:
    G.add_node(m, size=800, color='#4169E1')
    G.add_edge('CCIMCD', m)

# Get node properties
node_colors = [G.nodes[n]['color'] for n in G.nodes()]
node_sizes = [G.nodes[n]['size'] for n in G.nodes()]

# Use spring layout
pos = nx.spring_layout(G, k=3, iterations=50, seed=42)

# Draw network WITHOUT labels
nx.draw_networkx_nodes(G, pos, node_color=node_colors, node_size=node_sizes,
                       edgecolors='white', linewidths=3, ax=ax)
nx.draw_networkx_edges(G, pos, edge_color='white', width=2, alpha=0.6, ax=ax)

# Draw labels OUTSIDE nodes with LARGER font
for node, (x, y) in pos.items():
    # Offset label from node center
    offset = 0.15
    ax.text(x, y+offset, node, fontsize=16, ha='center', va='bottom',
           color='white', fontweight='bold',
           bbox=dict(boxstyle='round,pad=0.3', facecolor=BG_COLOR,
                    edgecolor='white', linewidth=1))

ax.axis('off')
ax.set_xlim(-1.2, 1.2)
ax.set_ylim(-1.2, 1.2)
save_fig('slide08_network.png')

# Slide 9: Funnel (already good)
print("\n[9/17] Creating funnel...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
top = Polygon([(2,6), (12,6), (11,4.5), (3,4.5)], facecolor='#87CEEB',
             edgecolor='white', linewidth=3)
ax.add_patch(top)
ax.text(7, 5.25, '10,000 touches', fontsize=20, ha='center', va='center', fontweight='bold')

mid = Polygon([(3.5,4.5), (10.5,4.5), (9,3), (5,3)], facecolor='#4682B4',
             edgecolor='white', linewidth=3)
ax.add_patch(mid)
ax.text(7, 3.75, '100 targets', fontsize=20, ha='center', va='center',
       fontweight='bold', color='white')

bot = Polygon([(5.5,3), (8.5,3), (7.5,1.5), (6.5,1.5)], facecolor='#000080',
             edgecolor='white', linewidth=3)
ax.add_patch(bot)
ax.text(7, 2.25, '30 successes', fontsize=18, ha='center', va='center',
       fontweight='bold', color='white')

ax.text(12.5, 3.75, '70% failure\nrate', fontsize=16, color=COLORS['failure_red'],
       fontweight='bold', ha='left')

ax.set_xlim(1, 14)
ax.set_ylim(0.5, 7)
ax.axis('off')
save_fig('slide09_funnel.png')

# ====================
# SLIDE 10: Highways with PROFESSIONAL colors
# ====================
print("\n[10/17] Creating highways with professional colors...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

lanes = [
    (0, 5, 10, 1.2, '#90EE90', 'LICIT', 'Legal acquisitions'),
    (0, 3.3, 10, 1.2, '#FFD700', 'GRAY ZONE', 'Ambiguous methods'),
    (0, 1.6, 10, 1.2, '#FF6B6B', 'ILLICIT', 'Illegal theft')
]

for x, y, w, h, color, label, desc in lanes:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white',
                    linewidth=3, alpha=0.7)
    ax.add_patch(rect)
    ax.text(1, y+h/2, label, fontsize=18, fontweight='bold', va='center', color='#1E3A5F')
    ax.text(5.5, y+h/2, desc, fontsize=13, va='center', color='#1E3A5F', style='italic')

# Convergence
converge = Polygon([(10,1.6), (10,6.2), (12.5,4.5), (12.5,3.3)],
                  facecolor='#6B46C1', edgecolor='white', linewidth=3)
ax.add_patch(converge)
ax.text(11.25, 3.9, 'MCF\nEcosystem', fontsize=16, ha='center', va='center',
       color='white', fontweight='bold', linespacing=1.3)

ax.set_xlim(-0.5, 13.5)
ax.set_ylim(0.5, 7.5)
ax.axis('off')
save_fig('slide10_highways.png')

# ====================
# SLIDE 11: ACTUAL Hexagonal Heatmap
# ====================
print("\n[11/17] Creating actual hexagonal heatmap...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

tech_domains = [
    ('AI/ML', 8, 0, '#FF6B6B'),
    ('Quantum\nComm', 5.5, 1.5, '#FF6B6B'),
    ('Biotech', 10.5, 1.5, '#FFA500'),
    ('5G/6G', 5.5, 4.5, '#FFA500'),
    ('Autonomous', 8, 6, '#FFA500'),
    ('Semis', 10.5, 4.5, '#FFD700'),
    ('Materials', 8, 3, '#FFD700'),
    ('Jet\nEngines', 3, 3, '#87CEEB'),
    ('Nuclear\nSub', 3, 6, '#87CEEB'),
    ('EUV', 3, 0, '#87CEEB')
]

for label, x, y, color in tech_domains:
    hex_coords = hexagon((x, y), 1.2)
    hex_patch = Polygon(hex_coords, facecolor=color, edgecolor='white', linewidth=3)
    ax.add_patch(hex_patch)
    ax.text(x, y, label, fontsize=13, ha='center', va='center',
           fontweight='bold', linespacing=1.2)

# Legend
legend_items = [
    (12, 5, '#FF6B6B', 'High Capability'),
    (12, 3.5, '#FFA500', 'Moderate'),
    (12, 2, '#FFD700', 'Developing'),
    (12, 0.5, '#87CEEB', 'Weak/Dependent')
]
for x, y, color, label in legend_items:
    circle = Circle((x, y), 0.3, facecolor=color, edgecolor='white', linewidth=2)
    ax.add_patch(circle)
    ax.text(x+0.5, y, label, fontsize=11, va='center', color='white', fontweight='bold')

ax.set_xlim(1, 15)
ax.set_ylim(-1, 7.5)
ax.axis('off')
save_fig('slide11_heatmap.png')

# ====================
# SLIDE 12: Bullseye with ALL circles labeled
# ====================
print("\n[12/17] Creating bullseye with all labels...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

circles = [
    (7, 4, 3.5, '#87CEEB', 'NQPF\n(2024+)\nWhole-of-Nation'),
    (7, 4, 2.3, '#4169E1', 'Expanded MCF\n(2020-2024)\nBroader Scope'),
    (7, 4, 1.2, '#000080', 'Traditional MCF\n(2015-2020)\nDefense Focus')
]

for x, y, r, color, label in circles:
    circle = Circle((x, y), r, facecolor=color, edgecolor='white', linewidth=3, alpha=0.7)
    ax.add_patch(circle)

# Labels positioned strategically
ax.text(7, 6.8, circles[0][4], fontsize=14, ha='center', va='center',
       color='white', fontweight='bold', linespacing=1.3)
ax.text(7, 4, circles[1][4], fontsize=13, ha='center', va='center',
       color='white', fontweight='bold', linespacing=1.3)
ax.text(7, 2.5, circles[2][4], fontsize=12, ha='center', va='center',
       color='white', fontweight='bold', linespacing=1.3)

# Expansion arrows
for angle in [30, 150, 270]:
    rad = np.deg2rad(angle)
    x_start = 7 + 1.2 * np.cos(rad)
    y_start = 4 + 1.2 * np.sin(rad)
    x_end = 7 + 3.5 * np.cos(rad)
    y_end = 4 + 3.5 * np.sin(rad)
    ax.annotate('', xy=(x_end, y_end), xytext=(x_start, y_start),
               arrowprops=dict(arrowstyle='->', lw=3, color='white'))

ax.set_xlim(2, 12)
ax.set_ylim(0, 8)
ax.axis('off')
save_fig('slide12_bullseye.png')

# ====================
# SLIDE 13: Star with lines BEHIND circles, FULL names
# ====================
print("\n[13/17] Creating star with lines behind...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

center = (7, 4)
# Draw lines FIRST (behind)
positions = [(7,6.5), (9.5,5.5), (9.5,2.5), (4.5,2.5), (4.5,5.5)]
for pos in positions:
    ax.plot([center[0], pos[0]], [center[1], pos[1]], 'w-', linewidth=2, alpha=0.4, zorder=1)
# Connect adjacent
for i in range(len(positions)):
    next_i = (i + 1) % len(positions)
    ax.plot([positions[i][0], positions[next_i][0]],
           [positions[i][1], positions[next_i][1]], 'w-', linewidth=2, alpha=0.4, zorder=1)

# Center circle
circle = Circle(center, 0.7, facecolor='#FFD700', edgecolor='white', linewidth=3, zorder=3)
ax.add_patch(circle)
ax.text(center[0], center[1], 'Belt & Road\nInitiative', fontsize=11, ha='center',
       va='center', fontweight='bold', linespacing=1.2)

# Surrounding circles with FULL names
initiatives = [
    ('Global Development\nInitiative', '#228B22'),
    ('Global Security\nInitiative', '#DC143C'),
    ('Global Civilization\nInitiative', '#4169E1'),
    ('Global AI\nGovernance Initiative', '#6A0DAD'),
    ('Digital Silk\nRoad', '#FFA500')
]

for pos, (label, color) in zip(positions, initiatives):
    circle = Circle(pos, 0.7, facecolor=color, edgecolor='white', linewidth=3, zorder=3)
    ax.add_patch(circle)
    ax.text(pos[0], pos[1], label, fontsize=9, ha='center', va='center',
           color='white', fontweight='bold', linespacing=1.2)

ax.set_xlim(2.5, 11.5)
ax.set_ylim(1, 7.5)
ax.axis('off')
save_fig('slide13_star.png')

# Slide 14: Bar chart (already good, just ensure proper colors)
print("\n[14/17] Creating bar chart...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
categories = ['Solar panels', 'Marine eng', 'Ag biotech', 'Mature semis',
              '5G infra', 'Mfg equip', 'Aviation',
              'Adv chips', 'Jet engines', 'Nuke sub', 'EUV']
values = [90, 85, 80, 75, 60, 50, 45, 20, 15, 10, 10]
colors_bars = [COLORS['success_green']]*4 + [COLORS['partial_yellow']]*3 + [COLORS['failure_red']]*4

y_pos = np.arange(len(categories))
ax.barh(y_pos, values, color=colors_bars, edgecolor='white', linewidth=2)
ax.set_yticks(y_pos)
ax.set_yticklabels(categories, fontsize=13, color='white')
ax.set_xlabel('Success Rate (%)', fontsize=14, fontweight='bold', color='white')
ax.set_xlim(0, 100)
ax.tick_params(colors='white')
ax.spines['bottom'].set_color('white')
ax.spines['left'].set_color('white')
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)

for i, v in enumerate(values):
    ax.text(v+2, i, f'{v}%', va='center', fontsize=12, fontweight='bold', color='white')

save_fig('slide14_bars.png')

# ====================
# SLIDE 15: ACTUAL World Map with engagement
# ====================
print("\n[15/17] Creating actual world map...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)

# Simplified world outline
continents = {
    'North America': [(2, 5), (3, 6), (4, 5.5), (3, 4)],
    'Europe': [(6, 5.5), (7, 6), (7.5, 5.5), (6.5, 5)],
    'Asia': [(8, 4), (10, 6), (11, 5), (10, 3), (8.5, 3.5)],
    'Africa': [(6, 3.5), (7, 4), (7, 2.5), (6, 2)],
    'South America': [(3, 2.5), (4, 3), (3.5, 1.5), (3, 2)],
    'Australia': [(11, 2), (12, 2.5), (11.5, 1.5)]
}

for name, coords in continents.items():
    poly = Polygon(coords, facecolor='#708090', edgecolor='white', linewidth=2, alpha=0.5)
    ax.add_patch(poly)

# China position
china_pos = (9.5, 4.5)
china = Circle(china_pos, 0.4, facecolor='#DC143C', edgecolor='white', linewidth=3)
ax.add_patch(china)
ax.text(china_pos[0], china_pos[1], 'CHINA', fontsize=10, ha='center', va='center',
       color='white', fontweight='bold')

# Engagement zones with lines from China
zones = [
    ((3.5, 5.5), 'Silicon Valley', '#FF6B6B'),
    ((7, 5.7), 'Cambridge', '#FFA500'),
    ((10.5, 5.5), 'Singapore', '#FFD700'),
    ((8.5, 5.5), 'Tel Aviv', '#FFD700')
]

for pos, label, color in zones:
    # Line from China
    ax.plot([china_pos[0], pos[0]], [china_pos[1], pos[1]],
           color=color, linewidth=2, linestyle='--', alpha=0.6)
    # Zone marker
    circle = Circle(pos, 0.2, facecolor=color, edgecolor='white', linewidth=2)
    ax.add_patch(circle)
    ax.text(pos[0], pos[1]-0.5, label, fontsize=9, ha='center',
           color='white', fontweight='bold')

ax.set_xlim(1, 13)
ax.set_ylim(0.5, 7)
ax.axis('off')
save_fig('slide15_worldmap.png')

# Slide 16: Three boxes (already good)
print("\n[16/17] Creating implication boxes...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
boxes = [
    (1, 5.5, 12, 1.3, '#000080', 'Systematic challenge requires\nsystematic response'),
    (1, 3.8, 12, 1.3, '#6B46C1', 'Single-point solutions will fail'),
    (1, 2.1, 12, 1.3, '#228B22', 'Understanding enables protection')
]
for x, y, w, h, color, text in boxes:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=3)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, text, fontsize=16, ha='center', va='center',
           color='white', fontweight='bold', linespacing=1.4)
ax.set_xlim(0, 14)
ax.set_ylim(1, 7.5)
ax.axis('off')
save_fig('slide16_boxes.png')

# Slide 17: Questions
print("\n[17/17] Creating questions slide...")
fig, ax = plt.subplots(figsize=(14, 8), facecolor=BG_COLOR)
positions_q = [(2,6), (12,6), (2,2), (12,2)]
for pos in positions_q:
    ax.text(pos[0], pos[1], '?', fontsize=100, ha='center', va='center',
           color='white', alpha=0.3)
ax.set_xlim(0, 14)
ax.set_ylim(0, 8)
ax.axis('off')
save_fig('slide17_questions.png')

print("\n" + "="*80)
print("ALL VISUALS GENERATED (V2 FIXED)")
print("="*80)

# ====================
# Build PowerPoint with DARK BLUE background and WHITE text
# ====================
print("\nBuilding PowerPoint with dark blue background...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

SLIDE_CONTENT = [
    ('Understanding Military-Civil Fusion', "China's Whole-of-Nation Technology Strategy", None),
    ('Our Roadmap', 'Seven key insights for understanding MCF\u2019s global impact', 'slide02_honeycomb.png'),
    ('The Historical Imperative', 'From humiliation to self-reliance: 180 years of technology dependence', 'slide03_timeline.png'),
    ('MCF Defined - Seven Domains of Fusion', 'Bidirectional integration across seven critical domains', 'slide04_gears.png'),
    ('Four Strategic Objectives', "China's stated goals - not Western interpretation", 'slide05_quadrant.png'),
    ('Policy Evolution', "From concept to constitution: MCF's institutional rise", 'slide06_stairs.png'),
    ('Expanding Legal Authority', 'Eight laws creating compelled participation', 'slide07_spider.png'),
    ('The Institutional Web', '40+ entities coordinating MCF implementation', 'slide08_network.png'),
    ('The Targeting Funnel', '10,000 touches \u2192 100 targets \u2192 30 successes', 'slide09_funnel.png'),
    ('Technology Transfer Pathways', 'Three highways to acquisition', 'slide10_highways.png'),
    ('Technology Domain Priorities', 'Uneven capabilities across ten critical domains', 'slide11_heatmap.png'),
    ('From MCF to New Quality Productive Forces', 'Expansion, not replacement', 'slide12_bullseye.png'),
    ('Global Initiatives - The Normalization Strategy', 'Five interconnected initiatives enabling MCF globally', 'slide13_star.png'),
    ('Track Record - Successes and Failures', 'Mixed results reveal patterns', 'slide14_bars.png'),
    ('Global Engagement Mechanisms', 'Multiple channels operating simultaneously worldwide', 'slide15_worldmap.png'),
    ('Key Implications', 'Three essential insights for capacity building', 'slide16_boxes.png'),
    ('Questions & Discussion', 'Your perspectives and challenges', 'slide17_questions.png')
]

for i, (title, subtitle, visual) in enumerate(SLIDE_CONTENT):
    print(f"  Adding slide {i+1}/17: {title}")

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Set slide background to dark blue
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(30, 58, 95)  # #1E3A5F

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)  # White
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Visual - PROPER SIZING to fit within slide
    if visual:
        viz_path = VIZ_DIR / visual
        if viz_path.exists():
            # Position to fit properly: leave room for title/subtitle
            slide.shapes.add_picture(
                str(viz_path),
                Inches(0.5), Inches(1.3),  # Start below subtitle
                width=Inches(9),
                height=Inches(4)  # Limit height to fit within slide
            )

# Save
output_pptx = Path("C:/Projects/OSINT - Foresight/MCF Presentations/MCF_Complete_V2_DarkBlue.pptx")
prs.save(str(output_pptx))

print(f"\n{'='*80}")
print(f"COMPLETE - V2 FIXED!")
print(f"{'='*80}")
print(f"\nPowerPoint saved to:")
print(f"  {output_pptx}")
print(f"\nFixes applied:")
print(f"  - Slide 2: Text labels instead of numbers")
print(f"  - Slide 3: Clearer timeline with descriptions")
print(f"  - Slide 4: No numbers, clear labels outside circles")
print(f"  - Slide 8: Larger fonts, labels outside bubbles")
print(f"  - Slide 10: Professional colors")
print(f"  - Slide 11: Actual hexagonal heatmap")
print(f"  - Slide 12: All circles labeled")
print(f"  - Slide 13: Lines behind circles, full names")
print(f"  - Slide 15: Actual world map with engagement")
print(f"  - All slides: Dark blue background, white text")
print(f"  - All visuals: Properly sized to fit within slides")
print(f"\nSize: {output_pptx.stat().st_size / (1024*1024):.1f} MB")
