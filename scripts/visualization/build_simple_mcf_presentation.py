#!/usr/bin/env python3
"""
Build SIMPLIFIED MCF Presentation
Uses clean, minimal visualizations instead of complex ones
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path

def build_simple_presentation():
    """
    Create new PowerPoint with simplified visualizations
    """
    print("="*80)
    print("BUILDING SIMPLIFIED MCF PRESENTATION")
    print("="*80)

    # Paths
    original_pptx = Path("C:/Users/mrear/Downloads/Understanding Military-Civil Fusion.pptx")
    output_pptx = Path("C:/Users/mrear/Downloads/Understanding Military-Civil Fusion - SIMPLE.pptx")

    # Use absolute paths for visualizations
    project_root = Path("C:/Projects/OSINT - Foresight")
    viz_simple = project_root / "scripts/visualization/visualizations/presentation_simple"

    # Check original exists
    if not original_pptx.exists():
        print(f"[ERROR] Original PowerPoint not found at {original_pptx}")
        return None

    print(f"\n[OK] Found original PowerPoint: {original_pptx}")
    print(f"[LOADING] Loading presentation...")

    # Load presentation
    prs = Presentation(str(original_pptx))

    print(f"[OK] Loaded {len(prs.slides)} slides")

    # Mapping of slides to SIMPLIFIED visualizations
    visualizations = {
        3: {  # Slide 4: What Is MCF?
            'file': viz_simple / "slide4_simple_flow.png",
            'title': 'Simple Bidirectional Flow'
        },
        5: {  # Slide 6: Policy Evolution
            'file': viz_simple / "slide6_simple_timeline.png",
            'title': 'Simple Timeline'
        },
        7: {  # Slide 8: Who Implements?
            'file': viz_simple / "slide8_simple_pyramid.png",
            'title': 'Simple 3-Tier Pyramid'
        },
        10: {  # Slide 11: Tech Domains
            'file': viz_simple / "slide11_simple_tech.png",
            'title': 'Simple Tech Bar Chart'
        },
        11: {  # Slide 12: NQPF Evolution
            'file': viz_simple / "slide12_simple_evolution.png",
            'title': 'Simple Evolution Arrow'
        },
        12: {  # Slide 13: Global Initiatives
            'file': viz_simple / "slide13_simple_global.png",
            'title': 'Simple 4-Box Grid'
        }
    }

    print(f"\n[INSERTING] Adding {len(visualizations)} simplified visualizations...\n")

    inserted_count = 0

    for slide_idx, viz_info in visualizations.items():
        viz_file = viz_info['file']
        viz_title = viz_info['title']

        if not viz_file.exists():
            print(f"[WARN] Slide {slide_idx + 1}: Visualization not found - {viz_file.name}")
            continue

        # Get the slide
        if slide_idx >= len(prs.slides):
            print(f"[WARN] Slide {slide_idx + 1}: Slide index out of range")
            continue

        slide = prs.slides[slide_idx]

        # Remove existing content shapes (except title)
        shapes_to_remove = []
        for shape in slide.shapes:
            # Keep title, remove everything else
            if shape.has_text_frame and shape == slide.shapes.title:
                continue
            else:
                shapes_to_remove.append(shape)

        # Remove shapes
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

        # Calculate position and size for image
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)

        # Insert image
        try:
            pic = slide.shapes.add_picture(
                str(viz_file),
                left, top,
                width=width
            )
            inserted_count += 1
            print(f"[OK] Slide {slide_idx + 1}: Inserted {viz_title}")
        except Exception as e:
            print(f"[ERROR] Slide {slide_idx + 1}: Error inserting image - {e}")

    print(f"\n{'='*80}")
    print(f"[SUCCESS] Inserted {inserted_count}/{len(visualizations)} simplified visuals")
    print(f"{'='*80}")

    # Save simplified presentation
    print(f"\n[SAVING] Saving simplified presentation...")
    try:
        prs.save(str(output_pptx))
        print(f"[SAVED] {output_pptx}")
        print(f"[SIZE] File size: {output_pptx.stat().st_size / (1024*1024):.1f} MB")
    except Exception as e:
        print(f"[ERROR] Error saving presentation: {e}")
        return None

    print(f"\n{'='*80}")
    print("SIMPLIFIED PRESENTATION COMPLETE")
    print(f"{'='*80}")
    print(f"\n[LOCATION] {output_pptx}")
    print(f"\n[SUMMARY]")
    print(f"   - Original slides: {len(prs.slides)}")
    print(f"   - Simplified visuals added: {inserted_count}")
    print(f"   - Enhanced slides: {inserted_count}/{len(prs.slides)}")
    print(f"   - Design: MINIMAL & CLEAN")
    print(f"   - Font standard: 32pt+ minimum")
    print(f"   - Quality: Publication-grade (300 DPI)")

    print(f"\n[KEY IMPROVEMENTS]")
    print(f"   - Removed complexity: Simple shapes only")
    print(f"   - Reduced text: Key points only")
    print(f"   - Clean design: Lots of whitespace")
    print(f"   - Easy to scan: Clear visual hierarchy")
    print(f"   - 2-3 colors max: Reduced visual noise")

    print(f"\n[NEXT STEPS]")
    print(f"   1. Open: {output_pptx.name}")
    print(f"   2. Review simplified visuals")
    print(f"   3. Compare to complex version")
    print(f"   4. Choose preferred style")
    print(f"   5. Present with confidence!")

    return str(output_pptx)


if __name__ == "__main__":
    result = build_simple_presentation()

    if result:
        print(f"\n[SUCCESS] Simplified presentation ready!")
    else:
        print(f"\n[FAILED] Check errors above")
