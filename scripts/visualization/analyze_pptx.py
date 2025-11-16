#!/usr/bin/env python3
"""
Extract and analyze PowerPoint presentation content
"""

from pptx import Presentation
import json
from pathlib import Path

def analyze_presentation(pptx_path):
    """Extract all content from PowerPoint file"""

    prs = Presentation(pptx_path)

    presentation_data = {
        'slide_count': len(prs.slides),
        'slides': []
    }

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_number': slide_num,
            'title': '',
            'content': [],
            'notes': '',
            'layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
            'shape_count': len(slide.shapes)
        }

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Check if it's likely a title
                if hasattr(shape, "text_frame"):
                    if shape == slide.shapes.title:
                        slide_data['title'] = shape.text.strip()
                    else:
                        slide_data['content'].append(shape.text.strip())

        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_data['notes'] = notes_slide.notes_text_frame.text.strip()

        presentation_data['slides'].append(slide_data)

    return presentation_data

if __name__ == "__main__":
    import sys
    import codecs

    # Set UTF-8 encoding for stdout
    if sys.stdout.encoding != 'utf-8':
        sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer, 'strict')

    pptx_path = "C:/Users/mrear/Downloads/Understanding Military-Civil Fusion.pptx"

    print("="*80)
    print("ANALYZING POWERPOINT PRESENTATION")
    print("="*80)
    print(f"\nFile: {pptx_path}\n")

    data = analyze_presentation(pptx_path)

    print(f"Total Slides: {data['slide_count']}\n")
    print("="*80)

    for slide in data['slides']:
        print(f"\n--- SLIDE {slide['slide_number']} ---")
        print(f"Layout: {slide['layout_name']}")
        print(f"Shapes: {slide['shape_count']}")

        if slide['title']:
            # Clean title for display
            title = slide['title'].encode('ascii', 'replace').decode('ascii')
            print(f"\nTITLE: {title}")

        if slide['content']:
            print(f"\nCONTENT:")
            for i, content in enumerate(slide['content'], 1):
                # Truncate very long content and clean for display
                content_preview = content[:500] + "..." if len(content) > 500 else content
                content_clean = content_preview.encode('ascii', 'replace').decode('ascii')
                print(f"  {i}. {content_clean}")

        if slide['notes']:
            notes_clean = slide['notes'][:200].encode('ascii', 'replace').decode('ascii')
            print(f"\nNOTES: {notes_clean}...")

        print("\n" + "-"*80)

    # Save to JSON for detailed analysis
    output_path = "C:/Projects/OSINT - Foresight/scripts/visualization/presentation_analysis.json"
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

    print(f"\n[SAVED] Detailed analysis: {output_path}")
    print("="*80)
