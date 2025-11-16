#!/usr/bin/env python3
"""
MCF Presentation - FULLY EDITABLE VERSION
All graphics created using PowerPoint native shapes
User can edit any element directly in PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import math

output_pptx = Path("C:/Projects/OSINT - Foresight/MCF Presentations/MCF_Complete_V3_Editable.pptx")

# Colors
BG_COLOR = RGBColor(30, 58, 95)  # Dark blue
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(220, 20, 60)
BLUE = RGBColor(135, 206, 235)
GREEN = RGBColor(144, 238, 144)
YELLOW = RGBColor(255, 215, 0)
ORANGE = RGBColor(255, 165, 0)
PURPLE = RGBColor(107, 70, 193)
NAVY = RGBColor(0, 0, 128)

print("="*80)
print("MCF PRESENTATION - FULLY EDITABLE")
print("Creating all graphics using PowerPoint shapes")
print("="*80)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_subtitle(slide, title_text, subtitle_text):
    """Add title and subtitle to slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.3))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 1: Title Slide
# ====================
print("\n[1/17] Creating Slide 1...")
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide1, BG_COLOR)

# Main title
title1 = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
tf1 = title1.text_frame
tf1.text = "Understanding Military-Civil Fusion"
p1 = tf1.paragraphs[0]
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle1 = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.8))
tf_sub1 = subtitle1.text_frame
tf_sub1.text = "China's Whole-of-Nation Technology Strategy"
p_sub1 = tf_sub1.paragraphs[0]
p_sub1.font.size = Pt(28)
p_sub1.font.color.rgb = WHITE
p_sub1.alignment = PP_ALIGN.CENTER

# Add decorative circles (editable)
circle1 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(4.5), Inches(0.8), Inches(0.8))
circle1.fill.solid()
circle1.fill.fore_color.rgb = BLUE
circle1.line.color.rgb = WHITE

circle2 = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(4.5), Inches(0.8), Inches(0.8))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RED
circle2.line.color.rgb = WHITE

# ====================
# SLIDE 2: Honeycomb with TEXT
# ====================
print("\n[2/17] Creating Slide 2 - Editable honeycomb...")
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide2, BG_COLOR)
add_title_subtitle(slide2, "Our Roadmap", "Seven key insights for understanding MCF's global impact")

labels = [
    'Historical\nRoots',
    'Bidirectional\nFusion',
    'Strategic\nGoals',
    'Legal\nExpansion',
    'Institutional\nWeb',
    'Global\nReach',
    'Implications'
]

# Create hexagons as rounded rectangles (easier to edit than true hexagons)
hex_size = Inches(1.2)
positions = [
    (Inches(4.4), Inches(2.5)),  # Center
    (Inches(4.4), Inches(1.5)),  # Top
    (Inches(5.5), Inches(2)),    # Top right
    (Inches(5.5), Inches(3)),    # Bottom right
    (Inches(4.4), Inches(3.5)),  # Bottom
    (Inches(3.3), Inches(3)),    # Bottom left
    (Inches(3.3), Inches(2))     # Top left
]

for i, (left, top) in enumerate(positions):
    # Rounded rectangle for easier editing
    shape = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, hex_size, Inches(0.8)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(74, 111, 165)
    shape.line.color.rgb = WHITE
    shape.line.width = Pt(3)

    # Add text
    text_frame = shape.text_frame
    text_frame.text = labels[i]
    text_frame.word_wrap = True
    para = text_frame.paragraphs[0]
    para.font.size = Pt(13)
    para.font.bold = True
    para.font.color.rgb = WHITE
    para.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# ====================
# SLIDE 3: Timeline
# ====================
print("\n[3/17] Creating Slide 3 - Editable timeline...")
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide3, BG_COLOR)
add_title_subtitle(slide3, "The Historical Imperative",
                   "From humiliation to self-reliance: 180 years of technology dependence")

events = [
    ('1840\nOpium Wars', Inches(1), 'Humiliation begins'),
    ('1978\nReform Era', Inches(2), 'Opening up'),
    ('1986\n863 Program', Inches(3.2), 'Tech focus'),
    ('1993\nYinhe', Inches(4.4), 'GPS denial'),
    ('2000\nIntegration', Inches(5.6), 'Civil-Military'),
    ('2015\nMCF', Inches(6.8), 'Formalized'),
    ('2024\nNQPF', Inches(8), 'Expanded')
]

# Timeline line (editable connector)
line = slide3.shapes.add_connector(1, Inches(1), Inches(3), Inches(9), Inches(3))
line.line.color.rgb = WHITE
line.line.width = Pt(4)

for label, x, desc in events:
    # Circle marker
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.7), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(2)

    # Year label above
    year_box = slide3.shapes.add_textbox(x - Inches(0.3), Inches(2), Inches(1.2), Inches(0.5))
    tf = year_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Description below
    desc_box = slide3.shapes.add_textbox(x - Inches(0.3), Inches(3.5), Inches(1.2), Inches(0.4))
    tf_desc = desc_box.text_frame
    tf_desc.text = desc
    p_desc = tf_desc.paragraphs[0]
    p_desc.font.size = Pt(9)
    p_desc.font.italic = True
    p_desc.font.color.rgb = WHITE
    p_desc.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 4: Gears (as circles with labels)
# ====================
print("\n[4/17] Creating Slide 4 - Editable gears...")
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide4, BG_COLOR)
add_title_subtitle(slide4, "MCF Defined - Seven Domains of Fusion",
                   "Bidirectional integration across seven critical domains")

gear_labels = ['Infrastructure', 'Industry', 'S&T', 'Education',
               'Social Services', 'Maritime', 'Emergency']
gear_positions = [
    (Inches(2), Inches(2.5)),
    (Inches(3), Inches(1.8)),
    (Inches(4), Inches(1.5)),
    (Inches(5), Inches(1.8)),
    (Inches(6), Inches(2.5)),
    (Inches(5), Inches(3.2)),
    (Inches(3), Inches(3.2))
]
colors_list = [RED]*3 + [BLUE]*4

for i, (left, top) in enumerate(gear_positions):
    # Circle
    circle = slide4.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors_list[i]
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(3)

    # Label below
    label_box = slide4.shapes.add_textbox(left - Inches(0.3), top + Inches(0.9),
                                          Inches(1.4), Inches(0.4))
    tf = label_box.text_frame
    tf.text = gear_labels[i]
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Bidirectional arrow
arrow_left = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.5), Inches(2.7),
                                     Inches(2), Inches(0.4))
arrow_left.fill.solid()
arrow_left.fill.fore_color.rgb = WHITE
arrow_left.line.color.rgb = WHITE

arrow_right = slide4.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.5), Inches(2.7),
                                      Inches(2), Inches(0.4))
arrow_right.fill.solid()
arrow_right.fill.fore_color.rgb = WHITE
arrow_right.line.color.rgb = WHITE

# Label
arrow_label = slide4.shapes.add_textbox(Inches(3.5), Inches(3.3), Inches(3), Inches(0.3))
tf_arrow = arrow_label.text_frame
tf_arrow.text = "BIDIRECTIONAL"
p_arrow = tf_arrow.paragraphs[0]
p_arrow.font.size = Pt(14)
p_arrow.font.bold = True
p_arrow.font.color.rgb = WHITE
p_arrow.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 5: Quadrant
# ====================
print("\n[5/17] Creating Slide 5 - Editable quadrant...")
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide5, BG_COLOR)
add_title_subtitle(slide5, "Four Strategic Objectives",
                   "China's stated goals - not Western interpretation")

quadrants = [
    (Inches(1), Inches(1.5), RGBColor(139, 0, 0), 'Military\nModernization'),
    (Inches(5.5), Inches(1.5), RGBColor(34, 139, 34), 'Economic\nDevelopment'),
    (Inches(1), Inches(3.5), RGBColor(255, 140, 0), 'Self-Reliance'),
    (Inches(5.5), Inches(3.5), RGBColor(107, 70, 193), 'First-Mover\nAdvantage')
]

for left, top, color, label in quadrants:
    rect = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Inches(4), Inches(1.8))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.fill.transparency = 0.2
    rect.line.color.rgb = WHITE
    rect.line.width = Pt(4)

    tf = rect.text_frame
    tf.text = label
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 6: Stairs
# ====================
print("\n[6/17] Creating Slide 6 - Editable stairs...")
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide6, BG_COLOR)
add_title_subtitle(slide6, "Policy Evolution", "From concept to constitution: MCF's institutional rise")

stairs_data = [
    (Inches(0.5), Inches(4.5), Inches(1.2), Inches(0.7), 'MCF\n2015'),
    (Inches(1.8), Inches(4), Inches(1.2), Inches(0.7), 'CMC\n2016'),
    (Inches(3.1), Inches(3.5), Inches(1.2), Inches(0.7), 'CCIMCD\n2017'),
    (Inches(4.4), Inches(3), Inches(1.2), Inches(0.7), 'Const\n2018'),
    (Inches(5.7), Inches(2.5), Inches(1.2), Inches(0.7), 'Dual\n2020'),
    (Inches(7), Inches(2), Inches(1.2), Inches(0.7), '14FYP\n2021'),
    (Inches(8.3), Inches(1.5), Inches(1.2), Inches(0.7), 'NQPF\n2023')
]

for i, (left, top, width, height, text) in enumerate(stairs_data):
    intensity = 0.4 + (i * 0.08)
    color_val = int(255 * intensity)
    rect = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0, 0, color_val)
    rect.line.color.rgb = WHITE
    rect.line.width = Pt(2)

    tf = rect.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 7: Spider Chart
# ====================
print("\n[7/17] Creating Slide 7 - Editable spider chart...")
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide7, BG_COLOR)
add_title_subtitle(slide7, "Expanding Legal Authority", "Eight laws creating compelled participation")

# Center of spider
center_x, center_y = Inches(5), Inches(3)
radius = Inches(1.5)

# 8 axes
laws = ['National\nSecurity', 'Intelligence\nLaw', 'Cyber\nsecurity', 'Data\nSecurity',
        'Encryption', 'Export\nControl', 'Nat Def\nMobilization', 'Mil Facilities\nProtection']

for i in range(8):
    angle = i * (2 * math.pi / 8) - math.pi/2
    end_x = center_x + radius * math.cos(angle)
    end_y = center_y + radius * math.sin(angle)

    # Draw line
    line = slide7.shapes.add_connector(1, center_x, center_y, end_x, end_y)
    line.line.color.rgb = WHITE
    line.line.width = Pt(2)

    # Add label at end
    label_x = center_x + (radius + Inches(0.3)) * math.cos(angle) - Inches(0.4)
    label_y = center_y + (radius + Inches(0.3)) * math.sin(angle) - Inches(0.2)
    label_box = slide7.shapes.add_textbox(label_x, label_y, Inches(0.8), Inches(0.4))
    tf = label_box.text_frame
    tf.text = laws[i]
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

# Center circle
center_circle = slide7.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.2), center_y - Inches(0.2),
                                        Inches(0.4), Inches(0.4))
center_circle.fill.solid()
center_circle.fill.fore_color.rgb = RED
center_circle.line.color.rgb = WHITE

# ====================
# SLIDE 8: Network
# ====================
print("\n[8/17] Creating Slide 8 - Editable network...")
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide8, BG_COLOR)
add_title_subtitle(slide8, "The Institutional Web", "40+ entities coordinating MCF implementation")

# Simple network layout
network_nodes = [
    ('Xi/CMC', Inches(5), Inches(2.8), Inches(0.5), RED),
    ('SASTIND', Inches(3.5), Inches(2.2), Inches(0.4), ORANGE),
    ('MOST', Inches(6.5), Inches(2.2), Inches(0.4), ORANGE),
    ('NDRC', Inches(3.5), Inches(3.5), Inches(0.4), BLUE),
    ('MIIT', Inches(6.5), Inches(3.5), Inches(0.4), BLUE),
    ('PLA', Inches(5), Inches(4.2), Inches(0.4), NAVY)
]

# Draw connections FIRST (behind nodes)
connections = [(0,1), (0,2), (0,3), (0,4), (0,5)]
for start, end in connections:
    line = slide8.shapes.add_connector(1,
                                       network_nodes[start][1] + network_nodes[start][2]/2,
                                       network_nodes[start][2] + network_nodes[start][2]/2,
                                       network_nodes[end][1] + network_nodes[end][2]/2,
                                       network_nodes[end][2] + network_nodes[end][2]/2)
    line.line.color.rgb = WHITE
    line.line.width = Pt(2)

# Draw nodes with labels OUTSIDE
for label, x, y, size, color in network_nodes:
    circle = slide8.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(3)

    # Label below circle
    label_box = slide8.shapes.add_textbox(x - Inches(0.3), y + size + Inches(0.05),
                                          size + Inches(0.6), Inches(0.3))
    tf = label_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 9: Funnel
# ====================
print("\n[9/17] Creating Slide 9 - Editable funnel...")
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide9, BG_COLOR)
add_title_subtitle(slide9, "The Targeting Funnel", "10,000 touches → 100 targets → 30 successes")

# Three-level funnel as rectangles with decreasing width
funnel_levels = [
    (Inches(1.5), Inches(1.8), Inches(7), Inches(0.8), '10,000 Touches', RGBColor(100, 150, 200)),
    (Inches(2.5), Inches(2.8), Inches(5), Inches(0.8), '100 Targets', RGBColor(70, 120, 180)),
    (Inches(3.5), Inches(3.8), Inches(3), Inches(0.8), '30 Successes', RGBColor(40, 90, 160))
]

for left, top, width, height, text, color in funnel_levels:
    rect = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = WHITE
    rect.line.width = Pt(3)

    tf = rect.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 10: Three Highways
# ====================
print("\n[10/17] Creating Slide 10 - Editable highways...")
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide10, BG_COLOR)
add_title_subtitle(slide10, "Technology Transfer Pathways", "Three highways to acquisition")

# Three lanes
highways = [
    (Inches(1), Inches(1.8), 'LICIT\nLicensing, Investment', GREEN),
    (Inches(1), Inches(2.8), 'GRAY ZONE\nAcademic, Dual-use', YELLOW),
    (Inches(1), Inches(3.8), 'ILLICIT\nTheft, Espionage', RGBColor(255, 107, 107))
]

for left, top, text, color in highways:
    rect = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                    Inches(6.5), Inches(0.7))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.fill.transparency = 0.3
    rect.line.color.rgb = WHITE
    rect.line.width = Pt(3)

    tf = rect.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

# Convergence point
converge = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2.5),
                                    Inches(1.8), Inches(1.3))
converge.fill.solid()
converge.fill.fore_color.rgb = PURPLE
converge.line.color.rgb = WHITE
converge.line.width = Pt(3)
tf = converge.text_frame
tf.text = 'MCF\nEcosystem'
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 11: Hexagonal Heatmap
# ====================
print("\n[11/17] Creating Slide 11 - Editable heatmap...")
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide11, BG_COLOR)
add_title_subtitle(slide11, "Technology Domain Priorities", "Uneven capabilities across ten critical domains")

# Simplified heatmap as colored rounded rectangles in honeycomb-like arrangement
tech_domains = [
    ('AI/ML', Inches(2), Inches(2), RGBColor(255, 107, 107)),
    ('Quantum', Inches(3.5), Inches(2), RGBColor(255, 107, 107)),
    ('Biotech', Inches(5), Inches(2), RGBColor(255, 165, 0)),
    ('5G/6G', Inches(6.5), Inches(2), RGBColor(255, 165, 0)),
    ('Autonomous', Inches(8), Inches(2), RGBColor(255, 215, 0)),
    ('Semicon', Inches(2.75), Inches(3.2), RGBColor(255, 215, 0)),
    ('Materials', Inches(4.25), Inches(3.2), RGBColor(255, 215, 0)),
    ('Jet Engines', Inches(5.75), Inches(3.2), RGBColor(135, 206, 235)),
    ('Nuclear Sub', Inches(7.25), Inches(3.2), RGBColor(135, 206, 235)),
    ('EUV Litho', Inches(4.5), Inches(4.3), RGBColor(135, 206, 235))
]

for label, left, top, color in tech_domains:
    hex_shape = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                         Inches(1.3), Inches(0.7))
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = color
    hex_shape.line.color.rgb = WHITE
    hex_shape.line.width = Pt(3)

    tf = hex_shape.text_frame
    tf.text = label
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

# Legend
legend_items = [
    ('High Capability', RGBColor(255, 107, 107)),
    ('Moderate', RGBColor(255, 165, 0)),
    ('Developing', RGBColor(255, 215, 0)),
    ('Weak/Dependent', RGBColor(135, 206, 235))
]
for i, (label, color) in enumerate(legend_items):
    box = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2 + i*0.5),
                                   Inches(0.3), Inches(0.3))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = WHITE

    text_box = slide11.shapes.add_textbox(Inches(0.9), Inches(2 + i*0.5), Inches(1.5), Inches(0.3))
    tf = text_box.text_frame
    tf.text = label
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.font.bold = True

# ====================
# SLIDE 12: Bullseye
# ====================
print("\n[12/17] Creating Slide 12 - Editable bullseye...")
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide12, BG_COLOR)
add_title_subtitle(slide12, "From MCF to New Quality Productive Forces", "Expansion, not replacement")

# Three concentric circles
circles_data = [
    (Inches(3.5), Inches(2), Inches(3), Inches(3), RGBColor(173, 216, 230), 'NQPF (2024+)\nWhole-of-Nation'),
    (Inches(4), Inches(2.5), Inches(2), Inches(2), RGBColor(100, 149, 237), 'Expanded MCF\n(2020-2024)'),
    (Inches(4.5), Inches(3), Inches(1), Inches(1), RGBColor(0, 0, 128), 'Traditional MCF\n(2015-2020)')
]

for left, top, width, height, color, label in circles_data:
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.fill.transparency = 0.3
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(4)

    # Add label inside
    tf = circle.text_frame
    tf.text = label
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 13: Star Pattern
# ====================
print("\n[13/17] Creating Slide 13 - Editable star...")
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide13, BG_COLOR)
add_title_subtitle(slide13, "Global Initiatives", "Five interconnected initiatives enabling MCF globally")

# Center and 5 surrounding nodes
center_x, center_y = Inches(5), Inches(3)
radius = Inches(1.4)

initiatives = [
    'Global Development\nInitiative',
    'Global Security\nInitiative',
    'Global Civilization\nInitiative',
    'Global AI Governance\nInitiative',
    'Digital Silk Road'
]

# Draw lines FIRST (behind circles)
for i in range(5):
    angle = i * (2 * math.pi / 5) - math.pi/2
    end_x = center_x + radius * math.cos(angle)
    end_y = center_y + radius * math.sin(angle)

    line = slide13.shapes.add_connector(1, center_x, center_y, end_x, end_y)
    line.line.color.rgb = WHITE
    line.line.width = Pt(2)

# Center circle
center_circle = slide13.shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.35), center_y - Inches(0.35),
                                         Inches(0.7), Inches(0.7))
center_circle.fill.solid()
center_circle.fill.fore_color.rgb = YELLOW
center_circle.line.color.rgb = WHITE
center_circle.line.width = Pt(3)
tf = center_circle.text_frame
tf.text = 'Belt & Road\nInitiative'
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.font.size = Pt(9)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.alignment = PP_ALIGN.CENTER

# Surrounding circles
for i, label in enumerate(initiatives):
    angle = i * (2 * math.pi / 5) - math.pi/2
    x = center_x + radius * math.cos(angle) - Inches(0.3)
    y = center_y + radius * math.sin(angle) - Inches(0.3)

    circle = slide13.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.color.rgb = WHITE
    circle.line.width = Pt(3)

    # Label outside
    label_offset = Inches(0.5)
    label_x = center_x + (radius + label_offset) * math.cos(angle) - Inches(0.5)
    label_y = center_y + (radius + label_offset) * math.sin(angle) - Inches(0.2)
    label_box = slide13.shapes.add_textbox(label_x, label_y, Inches(1), Inches(0.4))
    tf_label = label_box.text_frame
    tf_label.text = label
    tf_label.word_wrap = True
    p_label = tf_label.paragraphs[0]
    p_label.font.size = Pt(9)
    p_label.font.color.rgb = WHITE
    p_label.font.bold = True
    p_label.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 14: Bar Chart
# ====================
print("\n[14/17] Creating Slide 14 - Editable bars...")
slide14 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide14, BG_COLOR)
add_title_subtitle(slide14, "Track Record", "Mixed results reveal patterns")

# Horizontal bars
bars_data = [
    ('Solar Panels', Inches(6), GREEN, Inches(2)),
    ('Marine Engineering', Inches(5.5), GREEN, Inches(2.5)),
    ('5G Networks', Inches(4.5), YELLOW, Inches(3)),
    ('EV Manufacturing', Inches(4), YELLOW, Inches(3.5)),
    ('Advanced Chips', Inches(2.5), RGBColor(255, 107, 107), Inches(4)),
    ('Jet Engines', Inches(2), RGBColor(255, 107, 107), Inches(4.5))
]

for label, width, color, top in bars_data:
    # Bar
    bar = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), top,
                                   width, Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.fill.transparency = 0.2
    bar.line.color.rgb = WHITE
    bar.line.width = Pt(2)

    # Label
    label_box = slide14.shapes.add_textbox(Inches(0.3), top, Inches(2), Inches(0.35))
    tf = label_box.text_frame
    tf.text = label
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

# ====================
# SLIDE 15: World Map
# ====================
print("\n[15/17] Creating Slide 15 - Editable world map...")
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide15, BG_COLOR)
add_title_subtitle(slide15, "Global Engagement Mechanisms", "Multiple channels operating simultaneously worldwide")

# Simplified world representation - China and engagement zones
# China (center)
china_circle = slide15.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.5), Inches(2.8), Inches(0.5), Inches(0.5))
china_circle.fill.solid()
china_circle.fill.fore_color.rgb = RED
china_circle.line.color.rgb = WHITE
china_circle.line.width = Pt(3)
label_china = slide15.shapes.add_textbox(Inches(4.3), Inches(3.4), Inches(0.9), Inches(0.3))
tf_china = label_china.text_frame
tf_china.text = 'China'
p_china = tf_china.paragraphs[0]
p_china.font.size = Pt(10)
p_china.font.bold = True
p_china.font.color.rgb = WHITE
p_china.alignment = PP_ALIGN.CENTER

# Engagement zones
zones = [
    ('Silicon\nValley', Inches(1.5), Inches(2), RGBColor(255, 107, 107)),
    ('Cambridge', Inches(3), Inches(1.5), RGBColor(255, 165, 0)),
    ('Singapore', Inches(7), Inches(3.5), YELLOW),
    ('Tel Aviv', Inches(6), Inches(2), YELLOW)
]

for label, x, y, color in zones:
    # Zone circle
    zone = slide15.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.4), Inches(0.4))
    zone.fill.solid()
    zone.fill.fore_color.rgb = color
    zone.line.color.rgb = WHITE
    zone.line.width = Pt(2)

    # Connection line from China
    line = slide15.shapes.add_connector(1, Inches(4.75), Inches(3.05),
                                       x + Inches(0.2), y + Inches(0.2))
    line.line.color.rgb = WHITE
    line.line.width = Pt(1)
    line.line.dash_style = 2  # Dashed

    # Label
    label_box = slide15.shapes.add_textbox(x - Inches(0.2), y + Inches(0.5), Inches(0.8), Inches(0.3))
    tf_zone = label_box.text_frame
    tf_zone.text = label
    p_zone = tf_zone.paragraphs[0]
    p_zone.font.size = Pt(9)
    p_zone.font.color.rgb = WHITE
    p_zone.font.bold = True
    p_zone.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 16: Three Boxes
# ====================
print("\n[16/17] Creating Slide 16 - Editable boxes...")
slide16 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide16, BG_COLOR)
add_title_subtitle(slide16, "Key Implications", "Three essential insights for capacity building")

# Three stacked boxes
boxes = [
    (Inches(1.5), Inches(1.8), NAVY, 'Systematic challenge requires\nsystematic response'),
    (Inches(1.5), Inches(2.7), PURPLE, 'Single-point solutions will fail'),
    (Inches(1.5), Inches(3.6), GREEN, 'Understanding enables protection')
]

for left, top, color, text in boxes:
    box = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Inches(7), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.fill.transparency = 0.2
    box.line.color.rgb = WHITE
    box.line.width = Pt(4)

    tf = box.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# ====================
# SLIDE 17: Questions
# ====================
print("\n[17/17] Creating Slide 17 - Questions...")
slide17 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide17, BG_COLOR)
add_title_subtitle(slide17, "Questions & Discussion", "Your perspectives and challenges")

# Four question marks in corners
positions = [(Inches(1.5), Inches(2)), (Inches(7.5), Inches(2)),
             (Inches(1.5), Inches(4)), (Inches(7.5), Inches(4))]

for x, y in positions:
    q_box = slide17.shapes.add_textbox(x, y, Inches(0.8), Inches(0.8))
    tf = q_box.text_frame
    tf.text = '?'
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Save
prs.save(str(output_pptx))

print(f"\n{'='*80}")
print(f"COMPLETE - FULLY EDITABLE PRESENTATION!")
print(f"{'='*80}")
print(f"\nSaved to: {output_pptx}")
print(f"Size: {output_pptx.stat().st_size / (1024*1024):.1f} MB")
print(f"\nAll graphics are PowerPoint shapes - fully editable!")
print(f"  - Click any shape to edit")
print(f"  - Change colors, sizes, text")
print(f"  - Move, resize, reposition")
print(f"  - Group/ungroup elements")
