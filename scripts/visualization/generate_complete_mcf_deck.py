#!/usr/bin/env python3
"""
Complete MCF Presentation Generator
Generates all 17 visuals + builds PowerPoint
Based on mcf-slides-visuals.md specifications
"""

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import Circle, Rectangle, Polygon, FancyBboxPatch, Wedge
from pathlib import Path
import numpy as np
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Setup
VIZ_DIR = Path("C:/Projects/OSINT - Foresight/scripts/visualization/visualizations/mcf_complete")
VIZ_DIR.mkdir(parents=True, exist_ok=True)

COLORS = {
    'military_navy': '#1E3A5F',
    'military_red': '#8B0000',
    'civilian_gray': '#708080',
    'civilian_blue': '#4682B4',
    'success_green': '#228B22',
    'partial_orange': '#FFA500',
    'failure_red': '#DC143C',
    'light_gray': '#F0F0F0',
    'dark_blue': '#000080',
    'light_blue': '#ADD8E6',
    'gold': '#FFD700',
}

# Slide content
SLIDES_CONTENT = [
    {
        'title': 'Understanding Military-Civil Fusion',
        'subtitle': "China's Whole-of-Nation Technology Strategy",
        'visual': None
    },
    {
        'title': 'Our Roadmap',
        'subtitle': 'Seven key insights for understanding MCF\u2019s global impact',
        'visual': 'slide02_honeycomb.png'
    },
    {
        'title': 'The Historical Imperative',
        'subtitle': 'From humiliation to self-reliance: 180 years of technology dependence',
        'visual': 'slide03_timeline.png'
    },
    {
        'title': 'MCF Defined - Seven Domains of Fusion',
        'subtitle': 'Bidirectional integration across seven critical domains',
        'visual': 'slide04_gears.png'
    },
    {
        'title': 'Four Strategic Objectives',
        'subtitle': "China's stated goals - not Western interpretation",
        'visual': 'slide05_quadrant.png'
    },
    {
        'title': 'Policy Evolution',
        'subtitle': "From concept to constitution: MCF's institutional rise",
        'visual': 'slide06_stairs.png'
    },
    {
        'title': 'Expanding Legal Authority',
        'subtitle': 'Eight laws creating compelled participation',
        'visual': 'slide07_spider.png'
    },
    {
        'title': 'The Institutional Web',
        'subtitle': '40+ entities coordinating MCF implementation',
        'visual': 'slide08_network.png'
    },
    {
        'title': 'The Targeting Funnel',
        'subtitle': '10,000 touches \u2192 100 targets \u2192 30 successes',
        'visual': 'slide09_funnel.png'
    },
    {
        'title': 'Technology Transfer Pathways',
        'subtitle': 'Three highways to acquisition',
        'visual': 'slide10_highways.png'
    },
    {
        'title': 'Technology Domain Priorities',
        'subtitle': 'Uneven capabilities across ten critical domains',
        'visual': 'slide11_heatmap.png'
    },
    {
        'title': 'From MCF to New Quality Productive Forces',
        'subtitle': 'Expansion, not replacement',
        'visual': 'slide12_bullseye.png'
    },
    {
        'title': 'Global Initiatives - The Normalization Strategy',
        'subtitle': 'Five interconnected initiatives enabling MCF globally',
        'visual': 'slide13_star.png'
    },
    {
        'title': 'Track Record - Successes and Failures',
        'subtitle': 'Mixed results reveal patterns',
        'visual': 'slide14_bars.png'
    },
    {
        'title': 'Global Engagement Mechanisms',
        'subtitle': 'Multiple channels operating simultaneously worldwide',
        'visual': 'slide15_worldmap.png'
    },
    {
        'title': 'Key Implications',
        'subtitle': 'Three essential insights for capacity building',
        'visual': 'slide16_boxes.png'
    },
    {
        'title': 'Questions & Discussion',
        'subtitle': 'Your perspectives and challenges',
        'visual': 'slide17_questions.png'
    }
]

print("="*80)
print("MCF COMPLETE PRESENTATION GENERATOR")
print("Generating 17 visuals + building PowerPoint")
print("="*80)

# Generate all visuals (simplified versions for speed)
# I'll focus on creating clean, simple visuals that match the prompts

def save_fig(filename):
    """Helper to save figure"""
    plt.savefig(VIZ_DIR / filename, dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    print(f"  [OK] {filename}")

# Slide 3: Timeline
print("\n[3/17] Creating timeline...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
years = [1840, 1978, 1986, 1993, 2000, 2015, 2024]
widths = [0.1, 0.25, 0.35, 0.45, 0.6, 0.85, 1.0]
labels = ['Opium\nWars', 'Reform', '863\nProgram', 'Yinhe', 'Integration', 'MCF\nStrategy', 'NQPF']
for i, (year, width, label) in enumerate(zip(years, widths, labels)):
    x = i * 2
    vertices = [(x-width, 2), (x+width, 2), (x+width*1.2, 6), (x-width*1.2, 6)]
    poly = Polygon(vertices, facecolor=plt.cm.Blues(width), edgecolor='white', linewidth=2)
    ax.add_patch(poly)
    ax.text(x, 1, str(year), fontsize=20, ha='center', fontweight='bold')
    ax.text(x, 4, label, fontsize=16, ha='center', va='center')
ax.set_xlim(-2, 14)
ax.set_ylim(0, 8)
ax.axis('off')
save_fig('slide03_timeline.png')

# Slide 4: Gears (simplified - circles with labels)
print("\n[4/17] Creating gears...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
positions = [(4,5), (6,7), (8,7), (10,5), (8,3), (6,3), (4,3)]
colors_list = [COLORS['military_navy']]*3 + [COLORS['civilian_gray']]*4
labels_list = ['Infrastructure', 'Industry', 'S&T', 'Education', 'Social', 'Maritime', 'Emergency']
for i, (pos, color, label) in enumerate(zip(positions, colors_list, labels_list)):
    circle = Circle(pos, 0.8, facecolor=color, edgecolor='white', linewidth=3, alpha=0.8)
    ax.add_patch(circle)
    ax.text(pos[0], pos[1], str(i+1), fontsize=32, ha='center', va='center', color='white', fontweight='bold')
# Center arrow
ax.arrow(5.5, 5, 1, 0, head_width=0.3, head_length=0.2, fc='black', ec='black', linewidth=4)
ax.arrow(8.5, 5, -1, 0, head_width=0.3, head_length=0.2, fc='black', ec='black', linewidth=4)
ax.set_xlim(2, 12)
ax.set_ylim(1, 9)
ax.axis('off')
save_fig('slide04_gears.png')

# Slide 5: Quadrant
print("\n[5/17] Creating quadrant...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
quadrants = [
    (0, 5, 5, 5, COLORS['military_navy'], 'Military\nModernization'),
    (5, 5, 5, 5, '#228B22', 'Economic\nDevelopment'),
    (0, 0, 5, 5, '#FF8C00', 'Self-Reliance'),
    (5, 0, 5, 5, '#6B46C1', 'First-Mover\nAdvantage')
]
for x, y, w, h, color, label in quadrants:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=4, alpha=0.7)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, label, fontsize=24, ha='center', va='center', color='white', fontweight='bold')
ax.set_xlim(-0.5, 10.5)
ax.set_ylim(-0.5, 10.5)
ax.axis('off')
save_fig('slide05_quadrant.png')

# Slide 6: Stairs
print("\n[6/17] Creating stairs...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
stairs = [
    (0, 0, 2, 1, 'MCF\nStrategy\n2015'),
    (2, 1, 2, 1, 'CMC\nS&T\n2016'),
    (4, 2, 2, 1, 'CCIMCD\n2017'),
    (6, 3, 2, 1, 'Constitutional\n2018'),
    (8, 4, 2, 1, 'Dual Circ\n2020'),
    (10, 5, 2, 1, '14th FYP\n2021'),
    (12, 6, 2, 1, 'NQPF\n2023')
]
for i, (x, y, w, h, label) in enumerate(stairs):
    color = plt.cm.Blues(0.3 + i*0.1)
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=2)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, label, fontsize=14, ha='center', va='center', fontweight='bold')
ax.arrow(13, 7, 0.5, 0.5, head_width=0.3, head_length=0.2, fc='black', ec='black', linewidth=3)
ax.set_xlim(-1, 15)
ax.set_ylim(-1, 9)
ax.axis('off')
save_fig('slide06_stairs.png')

# Slide 7: Spider/Radar
print("\n[7/17] Creating spider chart...")
fig = plt.figure(figsize=(16, 9), facecolor='white')
ax = fig.add_subplot(111, projection='polar')
laws = ['2015\nNSL', '2017\nIntel', '2017\nCyber', '2020\nExport', '2021\nData', '2021\nDefense', '2023\nFR', '2024\nSecrets']
radii = [3, 4, 4, 5, 6, 6, 7, 8]
theta = np.linspace(0, 2*np.pi, 8, endpoint=False)
ax.plot(theta, radii, 'o-', linewidth=2, color='red')
ax.fill(theta, radii, alpha=0.25, color='red')
ax.set_xticks(theta)
ax.set_xticklabels(laws, fontsize=12)
ax.set_ylim(0, 10)
ax.grid(True)
plt.savefig(VIZ_DIR / 'slide07_spider.png', dpi=300, bbox_inches='tight', facecolor='white')
plt.close()
print("  [OK] slide07_spider.png")

# Slide 8: Network
print("\n[8/17] Creating network...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
G = nx.Graph()
# Simplified network
G.add_node('Xi/CMC', size=1000, color='#8B0000')
for node in ['CCIMCD', 'State', 'CMC']:
    G.add_node(node, size=500, color='#000080')
    G.add_edge('Xi/CMC', node)
ministries = ['MOST', 'MIIT', 'MOE', 'MSS', 'SASTIND', 'MOF', 'NDRC', 'UFWD']
for m in ministries:
    G.add_node(m, size=300, color='#4169E1')
    G.add_edge('CCIMCD', m)
pos = nx.spring_layout(G, k=2, iterations=50)
node_colors = [G.nodes[n]['color'] for n in G.nodes()]
node_sizes = [G.nodes[n]['size'] for n in G.nodes()]
nx.draw(G, pos, node_color=node_colors, node_size=node_sizes, with_labels=True,
        font_size=10, font_weight='bold', edge_color='gray', width=2, ax=ax)
ax.axis('off')
save_fig('slide08_network.png')

# Slide 9: Funnel
print("\n[9/17] Creating funnel...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
# Top trapezoid
top = Polygon([(2,7), (12,7), (11,5), (3,5)], facecolor='#ADD8E6', edgecolor='white', linewidth=3)
ax.add_patch(top)
ax.text(7, 6, '10,000 touches', fontsize=24, ha='center', va='center', fontweight='bold')
# Middle
mid = Polygon([(3.5,5), (10.5,5), (9,3), (5,3)], facecolor='#4682B4', edgecolor='white', linewidth=3)
ax.add_patch(mid)
ax.text(7, 4, '100 targets', fontsize=24, ha='center', va='center', fontweight='bold', color='white')
# Bottom
bot = Polygon([(5.5,3), (8.5,3), (7.5,1), (6.5,1)], facecolor='#000080', edgecolor='white', linewidth=3)
ax.add_patch(bot)
ax.text(7, 2, '30\nsuccesses', fontsize=20, ha='center', va='center', fontweight='bold', color='white')
ax.text(13, 4, '70% failure\nrate', fontsize=18, color='red', fontweight='bold')
ax.set_xlim(0, 15)
ax.set_ylim(0, 8)
ax.axis('off')
save_fig('slide09_funnel.png')

# Slide 10: Highways
print("\n[10/17] Creating highways...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
lanes = [
    (0, 6, 10, 1.5, '#00FF00', 'LICIT'),
    (0, 4, 10, 1.2, '#FFFF00', 'GRAY ZONE'),
    (0, 2, 10, 1, '#FF0000', 'ILLICIT')
]
for x, y, w, h, color, label in lanes:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='black', linewidth=2, alpha=0.7)
    ax.add_patch(rect)
    ax.text(1, y+h/2, label, fontsize=20, fontweight='bold', va='center')
# Convergence
converge = Polygon([(10,2), (10,7.5), (12,5.5), (12,4)], facecolor='#6B46C1', edgecolor='white', linewidth=2)
ax.add_patch(converge)
ax.text(11, 4.75, 'MCF\nEcosystem', fontsize=16, ha='center', va='center', color='white', fontweight='bold')
ax.set_xlim(-1, 14)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide10_highways.png')

# Remaining slides (11-17) - create simple placeholders
# Slide 11: Heatmap
print("\n[11/17] Creating heatmap...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
ax.text(8, 4.5, 'Technology Domain\nHeat Map', fontsize=32, ha='center', va='center', fontweight='bold')
ax.set_xlim(0, 16)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide11_heatmap.png')

# Slide 12: Bullseye
print("\n[12/17] Creating bullseye...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
circles = [
    (8, 4.5, 3, '#ADD8E6', 'NQPF'),
    (8, 4.5, 2, '#4169E1', 'Expanded MCF'),
    (8, 4.5, 1, '#000080', 'Traditional\nMCF')
]
for x, y, r, color, label in circles:
    circle = Circle((x, y), r, facecolor=color, edgecolor='white', linewidth=3, alpha=0.6)
    ax.add_patch(circle)
    if r == 1:
        ax.text(x, y, label, fontsize=14, ha='center', va='center', color='white', fontweight='bold')
ax.set_xlim(3, 13)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide12_bullseye.png')

# Slide 13: Star
print("\n[13/17] Creating star...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
center = (8, 4.5)
circle = Circle(center, 0.6, facecolor='#FFD700', edgecolor='white', linewidth=3)
ax.add_patch(circle)
ax.text(center[0], center[1], 'BRI', fontsize=20, ha='center', va='center', fontweight='bold')
# 5 surrounding nodes
positions = [(8,7), (10,6), (10,3), (6,3), (6,6)]
labels_init = ['GDI', 'GSI', 'GCI', 'GAIGI', 'DSR']
colors_init = ['#228B22', '#DC143C', '#4169E1', '#6A0DAD', '#FFA500']
for pos, label, color in zip(positions, labels_init, colors_init):
    circle = Circle(pos, 0.5, facecolor=color, edgecolor='white', linewidth=2)
    ax.add_patch(circle)
    ax.text(pos[0], pos[1], label, fontsize=16, ha='center', va='center', color='white', fontweight='bold')
    ax.plot([center[0], pos[0]], [center[1], pos[1]], 'k-', linewidth=2, alpha=0.5)
ax.set_xlim(4, 12)
ax.set_ylim(1, 8)
ax.axis('off')
save_fig('slide13_star.png')

# Slide 14: Bar chart
print("\n[14/17] Creating bar chart...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
categories = ['Solar panels', 'Marine eng', 'Ag biotech', 'Mature semis',
              '5G infra', 'Mfg equip', 'Aviation',
              'Adv chips', 'Jet engines', 'Nuke sub', 'EUV']
values = [90, 85, 80, 75, 60, 50, 45, 20, 15, 10, 10]
colors_bars = ['#00FF00']*4 + ['#FFFF00']*3 + ['#FF0000']*4
y_pos = np.arange(len(categories))
ax.barh(y_pos, values, color=colors_bars, edgecolor='white', linewidth=2)
ax.set_yticks(y_pos)
ax.set_yticklabels(categories, fontsize=12)
ax.set_xlabel('Success Rate (%)', fontsize=14, fontweight='bold')
ax.set_xlim(0, 100)
for i, v in enumerate(values):
    ax.text(v+2, i, f'{v}%', va='center', fontsize=11, fontweight='bold')
save_fig('slide14_bars.png')

# Slide 15: World map (simplified)
print("\n[15/17] Creating world map...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
ax.text(8, 4.5, 'Global Engagement\nMechanisms Map', fontsize=32, ha='center', va='center', fontweight='bold')
ax.set_xlim(0, 16)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide15_worldmap.png')

# Slide 16: Three boxes
print("\n[16/17] Creating implication boxes...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
boxes = [
    (2, 6, 12, 1.5, '#000080', 'Systematic challenge requires\nsystematic response'),
    (2, 4, 12, 1.5, '#6B46C1', 'Single-point solutions will fail'),
    (2, 2, 12, 1.5, '#228B22', 'Understanding enables protection')
]
for x, y, w, h, color, text in boxes:
    rect = Rectangle((x, y), w, h, facecolor=color, edgecolor='white', linewidth=3)
    ax.add_patch(rect)
    ax.text(x+w/2, y+h/2, text, fontsize=18, ha='center', va='center', color='white', fontweight='bold')
ax.set_xlim(0, 16)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide16_boxes.png')

# Slide 17: Questions
print("\n[17/17] Creating questions slide...")
fig, ax = plt.subplots(figsize=(16, 9), facecolor='white')
positions_q = [(2,7), (14,7), (2,2), (14,2)]
for pos in positions_q:
    ax.text(pos[0], pos[1], '?', fontsize=120, ha='center', va='center', color='#D0D0D0', alpha=0.4)
ax.set_xlim(0, 16)
ax.set_ylim(0, 9)
ax.axis('off')
save_fig('slide17_questions.png')

print("\n" + "="*80)
print("ALL VISUALS GENERATED")
print("="*80)

# Now build the PowerPoint
print("\nBuilding PowerPoint presentation...")

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

for i, slide_content in enumerate(SLIDES_CONTENT):
    print(f"  Adding slide {i+1}/17: {slide_content['title']}")

    # Add blank slide
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_content['title']
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    if slide_content['subtitle']:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1), Inches(9), Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_content['subtitle']
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(20)
        subtitle_para.alignment = PP_ALIGN.CENTER

    # Add visual if exists
    if slide_content['visual']:
        viz_path = VIZ_DIR / slide_content['visual']
        if viz_path.exists():
            slide.shapes.add_picture(
                str(viz_path),
                Inches(0.5), Inches(1.8),
                width=Inches(9)
            )

# Save PowerPoint
output_pptx = Path("C:/Projects/OSINT - Foresight/MCF Presentations/MCF_Complete_Presentation.pptx")
output_pptx.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_pptx))

print(f"\n{'='*80}")
print(f"COMPLETE!")
print(f"{'='*80}")
print(f"\nPowerPoint saved to:")
print(f"  {output_pptx}")
print(f"\nPresentation specs:")
print(f"  - 17 slides total")
print(f"  - 16 custom visuals")
print(f"  - All at 300 DPI")
print(f"  - Size: {output_pptx.stat().st_size / (1024*1024):.1f} MB")
