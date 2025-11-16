#!/usr/bin/env python3
"""
MCF PowerPoint Generator
Creates PowerPoint presentation with embedded MCF visualizations
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import os


def create_title_slide(prs, title, subtitle):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(231, 76, 60)  # Red

    return slide


def create_content_slide(prs, title, image_path, notes=""):
    """Create slide with title and image"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title

    # Style title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)  # Dark gray
    title_para.alignment = PP_ALIGN.CENTER

    # Add image
    if os.path.exists(image_path):
        img_left = Inches(0.5)
        img_top = Inches(1.3)
        img_width = Inches(9)

        try:
            pic = slide.shapes.add_picture(
                image_path,
                img_left,
                img_top,
                width=img_width
            )
            print(f"  [ADDED] {Path(image_path).name}")
        except Exception as e:
            print(f"  [ERROR] Could not add {Path(image_path).name}: {e}")
    else:
        print(f"  [MISSING] {image_path}")

    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    return slide


def create_section_slide(prs, section_title):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add section title in center
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = section_title

    # Style
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(231, 76, 60)  # Red
    title_para.alignment = PP_ALIGN.CENTER

    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(236, 240, 241)  # Light gray

    return slide


def create_mcf_visualization_presentation(output_dir="visualizations"):
    """
    Create comprehensive PowerPoint presentation with all MCF visualizations
    """
    output_path = Path(output_dir)
    viz_path = output_path

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("=" * 80)
    print("CREATING MCF VISUALIZATION POWERPOINT")
    print("=" * 80)
    print()

    # Title slide
    print("Creating title slide...")
    create_title_slide(
        prs,
        "China's Military-Civil Fusion (MCF)",
        "Institutional Architecture and Technology Flows\nVisualization Compendium"
    )
    print()

    # SECTION 1: Institutional Architecture
    print("Section 1: Institutional Architecture")
    create_section_slide(prs, "Institutional Architecture")

    # Slide: Full Institutional Architecture
    create_content_slide(
        prs,
        "MCF Institutional Architecture - Complete Network",
        str(viz_path / "mcf_institutional_architecture.png"),
        notes="Complete network of 18 institutions across 5 tiers, showing 21 key relationships. "
              "Node sizes represent power level, colors represent organization type. "
              "Hierarchical layout shows command structure from central leadership to implementation entities."
    )

    # Slide: Simplified Architecture
    create_content_slide(
        prs,
        "MCF Core Coordination Structure",
        str(viz_path / "mcf_simplified_architecture.png"),
        notes="Simplified view focusing on Tier 1-2 central authorities and key ministries. "
              "Shows the core coordination mechanisms between central leadership and key government agencies."
    )

    # Slide: Governance Hierarchy
    create_content_slide(
        prs,
        "MCF Governance Hierarchy",
        str(viz_path / "mcf_governance_hierarchy.png"),
        notes="Tree structure showing Xi Jinping at apex, chairing both Central MCF Commission and "
              "Central Military Commission. Shows dual civilian-military command structure with "
              "coordinated oversight of ministries, agencies, research institutions, and implementation entities."
    )

    # Slide: Institutional Layers
    create_content_slide(
        prs,
        "MCF Institutional Layers (Tier 1-5)",
        str(viz_path / "mcf_institutional_layers.png"),
        notes="Five-tier structure: (1) Central Leadership, (2) Key Ministries & Commissions, "
              "(3) Key Agencies, (4) Provincial Implementation, (5) Implementation Entities. "
              "Shows how directives flow from central coordination through provincial and "
              "institutional layers to operational implementation."
    )

    # Slide: Command Flow
    create_content_slide(
        prs,
        "MCF Command and Control Flows",
        str(viz_path / "mcf_command_flow.png"),
        notes="Relationship types: Commands (military chain), Directs (government oversight), "
              "Coordinates (inter-agency cooperation), Funds (research support), Controls (SOE management). "
              "Dashed lines show technology transfer and procurement relationships."
    )

    print()

    # SECTION 2: Technology Flows
    print("Section 2: Technology Flows")
    create_section_slide(prs, "Technology Flows")

    # Slide: Technology Flow
    create_content_slide(
        prs,
        "Technology Flow: Foreign Acquisition → Application",
        str(viz_path / "mcf_technology_flow.png"),
        notes="Three-stage flow: (1) Foreign acquisition from universities and companies, "
              "(2) Domestic processing through CAS, university labs, and SOEs, "
              "(3) Application in military and civilian domains. Flow values represent relative "
              "intensity of technology transfer."
    )

    # Slide: Four Global Initiatives
    create_content_slide(
        prs,
        "China's Four Global Initiatives Flow",
        str(viz_path / "mcf_four_initiatives_flow.png"),
        notes="Shows coordination of Belt and Road Initiative (BRI), Global Development Initiative (GDI), "
              "Global Security Initiative (GSI), and Global Civilization Initiative (GCI). "
              "Demonstrates how these initiatives flow through implementation mechanisms "
              "(infrastructure, finance, technology, security, culture) to achieve strategic objectives "
              "(economic influence, technology acquisition, security presence, soft power)."
    )

    print()

    # SECTION 3: Strategic Sectors
    print("Section 3: Strategic Priorities")
    create_section_slide(prs, "Strategic Priority Sectors")

    # Slide: Made in China 2025 Sectors - note: might not have PNG, use HTML note
    sector_png = viz_path / "mcf_sector_priority_flow.png"
    if sector_png.exists():
        create_content_slide(
            prs,
            "Made in China 2025: Sectors → Research → Application",
            str(sector_png),
            notes="Eight priority sectors from Made in China 2025 flow through three key research "
                  "institutions (CAS, University Defense Labs, SOEs) to military, civilian, and "
                  "dual-use applications. Shows systematic conversion of sectoral priorities into "
                  "research and development outputs."
        )
    else:
        # Create text slide noting HTML version
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        title = slide.shapes.title
        title.text = "Made in China 2025: Sectors → Research → Application"

        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Interactive Sankey diagram available in HTML format:\nmcf_sector_priority_flow.html"

        p = tf.add_paragraph()
        p.text = "\nEight priority sectors:"
        p.level = 0

        sectors = [
            "Information Technology",
            "Robotics & Automation",
            "Aerospace & Aviation",
            "Maritime Equipment",
            "New Energy Vehicles",
            "Power Equipment",
            "New Materials",
            "Biopharmaceuticals"
        ]

        for sector in sectors:
            p = tf.add_paragraph()
            p.text = sector
            p.level = 1

        print(f"  [NOTE] Created text slide for sector priority flow (HTML only)")

    print()

    # Final slide - Summary
    print("Creating summary slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    title.text = "Visualization Summary"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "This presentation includes:"

    viz_types = [
        "NetworkX network graphs (institutional architecture)",
        "Plotly Sankey diagrams (technology and initiative flows)",
        "Graphviz hierarchy trees (governance structure)",
        "",
        "All visualizations available in multiple formats:",
        "  • PNG (for presentations)",
        "  • SVG (for editing/scaling)",
        "  • PDF (for printing)",
        "  • HTML (interactive, Sankey diagrams only)",
        "",
        "Source data: scripts/visualization/mcf_network_data.py"
    ]

    for item in viz_types:
        p = tf.add_paragraph()
        p.text = item
        if item.startswith("  •"):
            p.level = 1

    print()

    # Save presentation
    output_file = output_path / "MCF_Visualization_Compendium.pptx"
    prs.save(str(output_file))

    print("=" * 80)
    print(f"[SAVED] PowerPoint: {output_file}")
    print("=" * 80)
    print()
    print(f"Total slides: {len(prs.slides)}")
    print()

    return str(output_file)


def create_quick_reference_deck(output_dir="visualizations"):
    """
    Create quick reference deck with just the key visualizations
    """
    output_path = Path(output_dir)
    viz_path = output_path

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating Quick Reference Deck...")

    # Title
    create_title_slide(
        prs,
        "MCF Quick Reference",
        "Key Institutional Architecture Visualizations"
    )

    # Key slides only
    create_content_slide(
        prs,
        "MCF Governance Hierarchy",
        str(viz_path / "mcf_governance_hierarchy.png")
    )

    create_content_slide(
        prs,
        "MCF Institutional Layers",
        str(viz_path / "mcf_institutional_layers.png")
    )

    create_content_slide(
        prs,
        "Technology Flow",
        str(viz_path / "mcf_technology_flow.png")
    )

    create_content_slide(
        prs,
        "Four Global Initiatives",
        str(viz_path / "mcf_four_initiatives_flow.png")
    )

    # Save
    output_file = output_path / "MCF_Quick_Reference.pptx"
    prs.save(str(output_file))

    print(f"[SAVED] Quick Reference: {output_file}")
    print()

    return str(output_file)


if __name__ == "__main__":
    # Create comprehensive presentation
    full_deck = create_mcf_visualization_presentation()

    # Create quick reference
    quick_deck = create_quick_reference_deck()

    print("=" * 80)
    print("POWERPOINT GENERATION COMPLETE")
    print("=" * 80)
    print()
    print("Created presentations:")
    print(f"  1. Full Compendium: {Path(full_deck).name}")
    print(f"  2. Quick Reference: {Path(quick_deck).name}")
    print()
    print("All visualizations embedded as high-resolution images")
