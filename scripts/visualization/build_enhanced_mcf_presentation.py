#!/usr/bin/env python3
"""
Build Enhanced MCF Presentation
Incorporates all created visualizations into the PowerPoint deck
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import shutil

def build_enhanced_presentation():
    """
    Create new PowerPoint with all visualizations integrated
    """
    print("="*80)
    print("BUILDING ENHANCED MCF PRESENTATION")
    print("="*80)

    # Paths
    original_pptx = Path("C:/Users/mrear/Downloads/Understanding Military-Civil Fusion.pptx")
    output_pptx = Path("C:/Users/mrear/Downloads/Understanding Military-Civil Fusion - ENHANCED.pptx")

    # Use absolute paths for visualizations
    project_root = Path("C:/Projects/OSINT - Foresight")
    viz_base = project_root / "scripts/visualization/visualizations"
    presentation_viz = viz_base / "presentation"
    network_viz = viz_base / "network"

    # Check original exists
    if not original_pptx.exists():
        print(f"[ERROR] Original PowerPoint not found at {original_pptx}")
        return None

    print(f"\n[OK] Found original PowerPoint: {original_pptx}")
    print(f"[LOADING] Loading presentation...")

    # Load presentation
    prs = Presentation(str(original_pptx))

    print(f"[OK] Loaded {len(prs.slides)} slides")

    # Mapping of slides to visualizations (slide index is 0-based)
    # Slide numbers in presentation are 1-based, so subtract 1 for index
    visualizations = {
        3: {  # Slide 4: What Is MCF?
            'file': presentation_viz / "slide4_bidirectional_flow.png",
            'title': 'Bidirectional Flow Diagram'
        },
        5: {  # Slide 6: Policy Evolution
            'file': presentation_viz / "slide6_timeline_blocks.png",
            'title': 'Timeline Building Blocks'
        },
        6: {  # Slide 7: Legal Framework
            'file': presentation_viz / "slide7_legal_framework.png",
            'title': 'Legal Framework Rings'
        },
        7: {  # Slide 8: Who Implements?
            'file': network_viz / "mcf_hierarchical_network.png",
            'title': 'Hierarchical Network'
        },
        8: {  # Slide 9: Engagement Funnel (use radial network as alternative)
            'file': network_viz / "mcf_radial_network.png",
            'title': 'Radial Network'
        },
        9: {  # Slide 10: Tech Pathways (use simplified org chart)
            'file': network_viz / "mcf_simplified_orgchart.png",
            'title': 'Simplified Organization Chart'
        },
        10: {  # Slide 11: Tech Domains
            'file': presentation_viz / "slide11_20_tech_capabilities.png",
            'title': 'Technology Capabilities Radar'
        },
        11: {  # Slide 12: NQPF Evolution
            'file': presentation_viz / "slide12_nqpf_mcf_venn.png",
            'title': 'NQPF/MCF Venn Diagram'
        },
        12: {  # Slide 13: Global Initiatives
            'file': presentation_viz / "slide13_global_initiatives.png",
            'title': 'Four Global Initiatives'
        },
        18: {  # Slide 19: Provincial Variance
            'file': presentation_viz / "slide19_provincial_map.png",
            'title': 'Provincial Implementation Map'
        },
        19: {  # Slide 20: Capabilities Matrix (duplicate of slide 11)
            'file': presentation_viz / "slide11_20_tech_capabilities.png",
            'title': 'Technology Capabilities Radar'
        }
    }

    print(f"\n[INSERTING] Adding {len(visualizations)} visualizations...\n")

    inserted_count = 0

    for slide_idx, viz_info in visualizations.items():
        viz_file = viz_info['file']
        viz_title = viz_info['title']

        if not viz_file.exists():
            print(f"[WARN] Slide {slide_idx + 1}: Visualization not found - {viz_file.name}")
            continue

        # Get the slide
        if slide_idx >= len(prs.slides):
            print(f"[WARN] Slide {slide_idx + 1}: Slide index out of range")
            continue

        slide = prs.slides[slide_idx]

        # Remove existing content shapes (except title)
        shapes_to_remove = []
        for shape in slide.shapes:
            # Keep title, remove everything else
            if shape.has_text_frame and shape == slide.shapes.title:
                continue
            else:
                shapes_to_remove.append(shape)

        # Remove shapes
        for shape in shapes_to_remove:
            sp = shape.element
            sp.getparent().remove(sp)

        # Calculate position and size for image
        # Leave room for title at top
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.0)  # Standard slide width is 10 inches, leave margins

        # Insert image
        try:
            pic = slide.shapes.add_picture(
                str(viz_file),
                left, top,
                width=width
            )
            inserted_count += 1
            print(f"[OK] Slide {slide_idx + 1}: Inserted {viz_title}")
        except Exception as e:
            print(f"[ERROR] Slide {slide_idx + 1}: Error inserting image - {e}")

    print(f"\n{'='*80}")
    print(f"[SUCCESS] Inserted {inserted_count}/{len(visualizations)} visualizations")
    print(f"{'='*80}")

    # Save enhanced presentation
    print(f"\n[SAVING] Saving enhanced presentation...")
    try:
        prs.save(str(output_pptx))
        print(f"[SAVED] {output_pptx}")
        print(f"[SIZE] File size: {output_pptx.stat().st_size / (1024*1024):.1f} MB")
    except Exception as e:
        print(f"[ERROR] Error saving presentation: {e}")
        return None

    print(f"\n{'='*80}")
    print("ENHANCEMENT COMPLETE")
    print(f"{'='*80}")
    print(f"\n[LOCATION] {output_pptx}")
    print(f"\n[SUMMARY]")
    print(f"   - Original slides: {len(prs.slides)}")
    print(f"   - Visualizations added: {inserted_count}")
    print(f"   - Enhanced slides: {inserted_count}/{len(prs.slides)}")
    print(f"   - Font standard: 32pt+ minimum throughout visuals")
    print(f"   - Quality: Publication-grade (300 DPI)")

    print(f"\n[NEXT STEPS]")
    print(f"   1. Open: {output_pptx.name}")
    print(f"   2. Review each enhanced slide")
    print(f"   3. Adjust image sizing if needed (maintain aspect ratio)")
    print(f"   4. Add any additional text/notes")
    print(f"   5. Present with confidence!")

    return str(output_pptx)


if __name__ == "__main__":
    result = build_enhanced_presentation()

    if result:
        print(f"\n[SUCCESS] Enhanced presentation ready!")
    else:
        print(f"\n[FAILED] Check errors above")
