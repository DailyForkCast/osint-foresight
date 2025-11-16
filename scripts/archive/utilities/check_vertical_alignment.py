from pptx import Presentation
from pptx.enum.text import MSO_VERTICAL_ANCHOR

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')
slide = prs.slides[8]  # Slide 9

print("=" * 80)
print("VERTICAL ALIGNMENT CHECK - SLIDE 9")
print("=" * 80)
print()

# Check title
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    text = shape.text.strip()
    if "BRI + FIVE INITIATIVES" in text:
        print("TITLE BOX:")
        print(f"  Text: {text}")
        print(f"  Position: Left={shape.left.inches:.3f}\", Top={shape.top.inches:.3f}\"")
        print(f"  Size: Width={shape.width.inches:.3f}\", Height={shape.height.inches:.3f}\"")

        # Check vertical anchor
        try:
            tf = shape.text_frame
            print(f"  Vertical anchor: {tf.vertical_anchor}")
            print(f"  Word wrap: {tf.word_wrap}")
        except:
            pass

        # Check font and paragraph
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                print(f"  Font size: {run.font.size.pt if run.font.size else 'N/A'}pt")
                break
            break
        print()

# Check initiative boxes
initiatives = ["DSR", "GSI", "GDI", "GCI", "GAGI"]
for init in initiatives:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text.strip()
        if text == init:
            print(f"{init} ACRONYM BOX:")
            print(f"  Position: Top={shape.top.inches:.3f}\", Height={shape.height.inches:.3f}\"")
            print(f"  Bottom: {(shape.top.inches + shape.height.inches):.3f}\"")
            try:
                print(f"  Vertical anchor: {shape.text_frame.vertical_anchor}")
            except:
                pass
            print()

# Check initiative full names
initiative_names = ["Digital Silk Road", "Global Security Initiative",
                    "Global Development Initiative", "Global Civilization Initiative",
                    "Global AI Governance Initiative"]
for name in initiative_names:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        if name in shape.text:
            print(f"{name[:20]}... BOX:")
            print(f"  Position: Top={shape.top.inches:.3f}\", Height={shape.height.inches:.3f}\"")
            print(f"  Bottom: {(shape.top.inches + shape.height.inches):.3f}\"")
            try:
                print(f"  Vertical anchor: {shape.text_frame.vertical_anchor}")
            except:
                pass
            print()

# Check descriptions
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    text = shape.text.strip()
    if "Data infrastructure" in text or "Security frameworks" in text or \
       "Development norms" in text or "Cultural influence" in text or \
       "AI standards" in text:
        print(f"DESCRIPTION BOX (first 30 chars: {text[:30]}):")
        print(f"  Position: Top={shape.top.inches:.3f}\", Height={shape.height.inches:.3f}\"")
        print(f"  Bottom: {(shape.top.inches + shape.height.inches):.3f}\"")
        try:
            print(f"  Vertical anchor: {shape.text_frame.vertical_anchor}")
        except:
            pass
        print()

print("=" * 80)
