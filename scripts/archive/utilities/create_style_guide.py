from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from collections import defaultdict, Counter

def rgb_to_hex(rgb):
    """Convert RGB tuple to hex color code"""
    if rgb is None:
        return None
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

def get_color_hex(color):
    """Extract hex color from color object"""
    if color is None:
        return None
    try:
        if color.type == 1:  # RGB color
            return rgb_to_hex((color.rgb[0], color.rgb[1], color.rgb[2]))
    except:
        pass
    return None

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')

# Collect style information
fonts = Counter()
font_sizes = Counter()
colors = Counter()
text_colors = Counter()
shape_fills = Counter()
layouts = []
alignments = Counter()
text_positions = {'titles': [], 'body': [], 'other': []}
shape_types = Counter()

# Analyze slides 3-8 (indices 2-7)
for i in range(2, 8):
    slide = prs.slides[i]
    slide_info = {
        'slide_number': i + 1,
        'shapes_count': len(slide.shapes),
        'text_blocks': []
    }

    for shape in slide.shapes:
        shape_types[str(shape.shape_type)] += 1

        # Analyze text
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            is_title = 'title' in shape.name.lower() if hasattr(shape, 'name') and shape.name else False

            # Position info
            pos = {
                'left': round(shape.left.inches, 2),
                'top': round(shape.top.inches, 2),
                'width': round(shape.width.inches, 2),
                'height': round(shape.height.inches, 2)
            }

            if is_title:
                text_positions['titles'].append(pos)
            elif len(text) > 100:
                text_positions['body'].append(pos)
            else:
                text_positions['other'].append(pos)

            for para in shape.text_frame.paragraphs:
                if para.alignment:
                    alignments[str(para.alignment)] += 1

                for run in para.runs:
                    if run.text.strip():
                        # Font info
                        if run.font.name:
                            fonts[run.font.name] += len(run.text)
                        if run.font.size:
                            font_sizes[f"{run.font.size.pt}pt"] += 1

                        # Text color
                        color = get_color_hex(run.font.color)
                        if color:
                            text_colors[color] += 1

        # Analyze fills
        if hasattr(shape, 'fill'):
            try:
                if shape.fill.type == 1:  # Solid fill
                    color = get_color_hex(shape.fill.fore_color)
                    if color:
                        shape_fills[color] += 1
                        colors[color] += 1
            except:
                pass

        # Analyze lines
        if hasattr(shape, 'line'):
            try:
                if hasattr(shape.line, 'color'):
                    color = get_color_hex(shape.line.color)
                    if color:
                        colors[color] += 1
            except:
                pass

    layouts.append(slide_info)

# Generate style guide
style_guide = []

style_guide.append("=" * 80)
style_guide.append("POWERPOINT STYLE GUIDE - MCF v2.1 (Slides 3-8)")
style_guide.append("=" * 80)
style_guide.append("")

# Slide dimensions
style_guide.append("1. SLIDE DIMENSIONS")
style_guide.append("-" * 80)
style_guide.append(f"   Dimensions: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
style_guide.append(f"   Aspect Ratio: 16:9 (widescreen)")
style_guide.append("")

# Typography
style_guide.append("2. TYPOGRAPHY")
style_guide.append("-" * 80)
style_guide.append("   Primary Fonts (by usage):")
for font, count in fonts.most_common(5):
    percentage = (count / sum(fonts.values())) * 100
    style_guide.append(f"      • {font}: {percentage:.1f}% of text content")
style_guide.append("")

style_guide.append("   Font Sizes (by frequency):")
for size, count in sorted(font_sizes.items(), key=lambda x: float(x[0].replace('pt', '')), reverse=True):
    style_guide.append(f"      • {size}: used {count} times")
style_guide.append("")

# Colors
style_guide.append("3. COLOR PALETTE")
style_guide.append("-" * 80)
style_guide.append("   Primary Colors (all elements):")
for color, count in colors.most_common(10):
    style_guide.append(f"      • {color} (used {count} times)")
style_guide.append("")

style_guide.append("   Text Colors:")
for color, count in text_colors.most_common(5):
    style_guide.append(f"      • {color} (used {count} times)")
style_guide.append("")

style_guide.append("   Shape/Background Fills:")
for color, count in shape_fills.most_common(5):
    style_guide.append(f"      • {color} (used {count} times)")
style_guide.append("")

# Layout patterns
style_guide.append("4. LAYOUT PATTERNS")
style_guide.append("-" * 80)

if text_positions['titles']:
    avg_title = {
        'left': sum(p['left'] for p in text_positions['titles']) / len(text_positions['titles']),
        'top': sum(p['top'] for p in text_positions['titles']) / len(text_positions['titles']),
        'width': sum(p['width'] for p in text_positions['titles']) / len(text_positions['titles']),
        'height': sum(p['height'] for p in text_positions['titles']) / len(text_positions['titles']),
    }
    style_guide.append(f"   Title Position (average across {len(text_positions['titles'])} titles):")
    style_guide.append(f"      • Left: {avg_title['left']:.2f}\"")
    style_guide.append(f"      • Top: {avg_title['top']:.2f}\"")
    style_guide.append(f"      • Width: {avg_title['width']:.2f}\"")
    style_guide.append(f"      • Height: {avg_title['height']:.2f}\"")
    style_guide.append("")

if text_positions['body']:
    avg_body = {
        'left': sum(p['left'] for p in text_positions['body']) / len(text_positions['body']),
        'top': sum(p['top'] for p in text_positions['body']) / len(text_positions['body']),
        'width': sum(p['width'] for p in text_positions['body']) / len(text_positions['body']),
        'height': sum(p['height'] for p in text_positions['body']) / len(text_positions['body']),
    }
    style_guide.append(f"   Body Text Position (average across {len(text_positions['body'])} blocks):")
    style_guide.append(f"      • Left: {avg_body['left']:.2f}\"")
    style_guide.append(f"      • Top: {avg_body['top']:.2f}\"")
    style_guide.append(f"      • Width: {avg_body['width']:.2f}\"")
    style_guide.append(f"      • Height: {avg_body['height']:.2f}\"")
    style_guide.append("")

# Text alignment
style_guide.append("5. TEXT ALIGNMENT")
style_guide.append("-" * 80)
for align, count in alignments.most_common():
    style_guide.append(f"   • {align}: {count} instances")
style_guide.append("")

# Shape types
style_guide.append("6. SHAPE TYPES USED")
style_guide.append("-" * 80)
for shape_type, count in shape_types.most_common():
    style_guide.append(f"   • {shape_type}: {count} instances")
style_guide.append("")

# Per-slide summary
style_guide.append("7. PER-SLIDE BREAKDOWN")
style_guide.append("-" * 80)
for slide_info in layouts:
    style_guide.append(f"   Slide {slide_info['slide_number']}: {slide_info['shapes_count']} shapes")
style_guide.append("")

style_guide.append("=" * 80)
style_guide.append("END OF STYLE GUIDE")
style_guide.append("=" * 80)

# Save style guide
with open('C:/Projects/OSINT-Foresight/MCF_STYLE_GUIDE.txt', 'w', encoding='utf-8') as f:
    f.write('\n'.join(style_guide))

print('\n'.join(style_guide))
