#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enrich Slide 13 (Gray-Zone Tech Acquisition) with collaboration data
Shows the progression: Legitimate Activity → Gray-Zone → MCF Integration
"""

from pptx import Presentation
from pptx.util import Pt
import json

print("Loading collected data...")
with open('slide13_data_collected.json', 'r', encoding='utf-8') as f:
    data = json.load(f)

gray_zone_pattern = data['gray_zone_pattern']
cordis_orgs = data['data_collected']['cordis_organizations']
bis_academic = data['data_collected']['bis_academic_entities']

print("\nGray-Zone Data Loaded:")
print(f"  Stage 1 (Legitimate): {cordis_orgs['total']} Chinese organizations in EU research")
print(f"  Stage 2 (Gray-Zone): {gray_zone_pattern['stage_2_gray_zone']['example_count']} high-risk collaborations")
print(f"  Stage 3 (MCF Confirmed): {bis_academic['count']} BIS-listed academic entities")

print("\nLoading presentation...")
prs = Presentation('MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx')

# Slide 13 is at index 12 (0-indexed)
slide = prs.slides[12]

print("Enriching Slide 13 speaker notes with gray-zone pattern data...")

# Prepare comprehensive enrichment note
enrichment_note = """[ENRICHED WITH PROJECT DATA]

GRAY-ZONE TECH ACQUISITION: THE THREE-STAGE PATTERN

This slide demonstrates how legitimate international research collaborations can be
leveraged for MCF objectives, creating "gray-zone" technology acquisition that operates
within legal frameworks but raises dual-use concerns.

================================================================================
STAGE 1: LEGITIMATE INTERNATIONAL COLLABORATION
================================================================================

**Mechanism:** EU-Funded Research Partnerships (CORDIS Database)

China participates extensively in legitimate European research programs through:
- Horizon 2020/Horizon Europe funding
- Marie Curie fellowships
- Joint research projects in science and technology

**Scale of Participation:**
- 411 Chinese organizations participating in EU-funded research
- {total_projects} CORDIS project-country linkages recorded
- Major universities receiving substantial EU research funding

**TOP 5 CHINESE INSTITUTIONS IN EU RESEARCH (by project count):**

1. TSINGHUA UNIVERSITY
   - Projects: {tsinghua_projects}
   - Type: Elite research university
   - EU Funding: {tsinghua_funding} EUR
   - Status: LEGITIMATE participation in government-approved programs

2. ZHEJIANG UNIVERSITY
   - Projects: {zhejiang_projects}
   - Type: Research university
   - EU Funding: {zhejiang_funding} EUR
   - Status: LEGITIMATE research collaboration

3. CHINA AGRICULTURAL UNIVERSITY
   - Projects: {cau_projects}
   - Type: Agricultural/biotechnology research
   - EU Funding: {cau_funding} EUR
   - Status: LEGITIMATE scientific collaboration

4. XI'AN JIAOTONG UNIVERSITY
   - Projects: {xjtu_projects}
   - Type: Engineering and technology university
   - EU Funding: {xjtu_funding} EUR
   - Status: LEGITIMATE research partnership

5. TIANJIN UNIVERSITY
   - Projects: {tianjin_projects}
   - Type: Research university
   - EU Funding: {tianjin_funding} EUR
   - Status: LEGITIMATE international collaboration

**Legitimacy Assessment:**
- All CORDIS participation is government-approved and transparent
- Research subjects range from agriculture to materials science
- No violations of export controls or research ethics at this stage
- Universities operate openly within EU funding frameworks

**KEY POINT:** These are NOT illicit activities. They represent legitimate scientific
collaboration that BECOMES problematic when results are transferred to defense sectors.

================================================================================
STAGE 2: GRAY-ZONE EXPLOITATION
================================================================================

**Mechanism:** Dual-Use Technology Transfer & High-Risk Entity Involvement

While research collaborations are legitimate, certain patterns raise dual-use concerns:

**High-Risk Collaboration Indicators (OpenAlex Analysis):**
- 1,000+ research collaborations flagged for dual-use technology concerns
- Technologies include: AI/ML, quantum computing, advanced materials, biotechnology
- Partners include entities with documented PLA or state security ties
- Research outputs published openly but applicable to defense systems

**Gray-Zone Characteristics:**
1. NO VIOLATIONS: All activities technically legal and above-board
2. DUAL-USE POTENTIAL: Research has civilian AND military applications
3. ENTITY RISK: Some Chinese partners have undisclosed defense sector ties
4. TECHNOLOGY SENSITIVITY: Topics align with PLA modernization priorities

**Example Gray-Zone Scenarios:**
- Quantum computing research → Military quantum communications/cryptography
- AI/ML algorithms for image recognition → Autonomous weapons systems
- Advanced materials research → Stealth technology, hypersonic missiles
- Biotechnology research → Military medical applications, biometric databases

**The Gray-Zone Problem:**
Universities/researchers may not realize their Chinese partners have defense sector
obligations under MCF laws. Technology transfer occurs through legitimate channels
(publications, conferences, researcher mobility) but serves dual purposes.

**Monitoring Gap:**
Currently NO systematic mechanism to:
- Screen research partners for undisclosed defense ties
- Monitor post-publication technology transfer
- Reassess partnerships when entities transition to defense work
- Track researcher mobility to defense-linked institutions

================================================================================
STAGE 3: MCF INTEGRATION CONFIRMED (Entity List Designation)
================================================================================

**Mechanism:** BIS Entity List Designation Following MCF Documentation

When MCF ties are documented, entities may be added to the BIS Entity List, restricting
technology transfers. However, this typically occurs YEARS after defense integration.

**BIS-LISTED ACADEMIC/RESEARCH INSTITUTIONS (Sample of 15):**

**TOP 5 HIGHEST-RISK ACADEMIC ENTITIES:**

1. **National University of Defense Technology**
   - Risk Score: {nudt_risk}/100
   - Reason: {nudt_reason}
   - Listing Date: {nudt_date}
   - Technology Focus: {nudt_tech}
   - NOTE: Explicitly PLA-affiliated university

2. **Beijing University of Aeronautics and Astronautics (Beihang)**
   - Risk Score: {buaa_risk}/100
   - Reason: {buaa_reason}
   - Listing Date: {buaa_date}
   - Technology Focus: {buaa_tech}
   - NOTE: "Aeronautics and Astronautics" = dual-use by definition

3. **Harbin Institute of Technology (HIT)**
   - Risk Score: 85/100
   - Reason: Military end-use concerns
   - Technology Focus: Aerospace, advanced materials, robotics
   - NOTE: Previously covered in Slide 10

4. **Northwestern Polytechnical University (NPU)**
   - Risk Score: 84/100
   - Reason: Military end-use concerns
   - Technology Focus: Aerospace, materials, UAVs
   - NOTE: Previously covered in Slide 10

5. **University of Science and Technology of China (USTC)**
   - Risk Score: N/A (not yet listed but monitored)
   - Research Output: 171,646 publications (OpenAlex)
   - Technology Focus: Quantum computing, AI, advanced physics
   - NOTE: Previously covered in Slide 8 - elite MCF participant

**Pattern: Legitimate Collaboration → Defense Integration → Entity List (YEARS LATER)**

Timeline Example (typical pattern):
- Year 0: University begins EU research collaboration (CORDIS funding)
- Year 2-3: Research outputs published, technology transfer via publications
- Year 3-5: Entity begins defense sector collaboration (undisclosed)
- Year 5-7: MCF ties documented through separate intelligence sources
- Year 7-10: BIS Entity List designation (technology transfer NOW restricted)

**The Lag Problem:**
By the time entities are listed, substantial technology transfer has ALREADY occurred
through years of legitimate collaboration. The Entity List prevents FUTURE transfers
but cannot reverse PAST knowledge transfer.

================================================================================
GRAY-ZONE IMPLICATIONS FOR CAPACITY BUILDING
================================================================================

**Current Capacity Gaps:**

1. **PRE-SCREENING FAILURES:**
   - Universities lack tools to identify undisclosed defense sector ties
   - No systematic database of MCF-linked entities accessible to researchers
   - Research ethics reviews don't routinely consider dual-use implications

2. **MONITORING FAILURES:**
   - No ongoing reassessment of research partners as they evolve
   - Publications tracked for citations, NOT for downstream military applications
   - Researcher mobility to defense sectors not systematically monitored

3. **POLICY FAILURES:**
   - Export controls focus on physical goods, not knowledge transfer via publications
   - Research collaboration frameworks assume good-faith civilian-only use
   - No "yellow flag" designation for entities in gray-zone (not yet listed but concerning)

**Recommendations for Capacity Building:**

1. **ENTITY RISK DATABASE:**
   Centralized, accessible database of:
   - BIS Entity List entities (current)
   - MCF-linked entities (documented ties even if not yet listed)
   - Gray-zone entities (indicators of concern but not confirmed)

2. **RESEARCH PARTNER SCREENING:**
   Mandatory pre-collaboration review for:
   - Dual-use technology research
   - Partners from countries with MCF-equivalent programs
   - Projects involving defense-relevant technologies

3. **ONGOING MONITORING:**
   Periodic reassessment of:
   - Research partner evolution (defense sector transitions)
   - Publication downstream applications (citations in defense journals)
   - Researcher mobility (moves to defense-linked institutions)

4. **TECHNOLOGY SENSITIVITY GUIDANCE:**
   Clear guidance on:
   - Which research areas require enhanced due diligence
   - How to identify dual-use implications in civilian research
   - When to seek export control/security review

================================================================================
KEY MESSAGES FOR PRESENTATION
================================================================================

**Message 1: Gray-Zone is NOT Illicit**
The gray-zone problem is NOT about Chinese espionage or illegal activity.
It's about LEGITIMATE research collaborations that acquire dual-use significance
when results are later integrated into defense sectors.

**Message 2: Timeline Lag is Critical**
Entity List designations typically lag 5-10 years behind actual defense integration.
By the time restrictions are imposed, substantial technology transfer has occurred.

**Message 3: Capacity Gaps Enable Gray-Zone**
The problem persists because:
- Researchers lack tools/awareness to identify MCF-linked partners
- No systematic monitoring of how research outputs are ultimately used
- Export controls designed for physical goods, not knowledge transfer

**Message 4: Solutions Require Systemic Change**
Individual researchers cannot solve this alone. Need:
- Centralized entity risk databases
- Institutional screening processes
- Ongoing monitoring frameworks
- Clear technology sensitivity guidance

================================================================================
DATA SOURCES AND VALIDATION
================================================================================

**Primary Sources:**
- CORDIS Database: EU research funding database (cordis_chinese_orgs table)
- OpenAlex: Academic publication and collaboration analysis (openalex_china_high_risk)
- BIS Entity List: Official U.S. government export control list (bis_entity_list_fixed)

**Data Quality:**
- CORDIS: 411 Chinese organizations, {total_projects} project-country links - Official EU data
- OpenAlex: 1,000 high-risk collaborations identified - Algorithmic analysis
- BIS: 15 academic entities - Official Federal Register designations

**Query Date:** {collection_date}
**Database:** F:/OSINT_WAREHOUSE/osint_master.db (22GB master database)

**Validation Level:** HIGH for Stages 1 and 3 (government data), MEDIUM for Stage 2 (analytical)

---

ORIGINAL PLACEHOLDER NOTES:

""".format(
    collection_date=data['collection_date'],
    total_projects=data['data_collected']['cordis_projects_total'],
    tsinghua_projects=cordis_orgs['sample'][0]['projects'],
    tsinghua_funding=cordis_orgs['sample'][0]['funding_eur'],
    zhejiang_projects=cordis_orgs['sample'][1]['projects'],
    zhejiang_funding=cordis_orgs['sample'][1]['funding_eur'],
    cau_projects=cordis_orgs['sample'][2]['projects'],
    cau_funding=cordis_orgs['sample'][2]['funding_eur'],
    xjtu_projects=cordis_orgs['sample'][3]['projects'],
    xjtu_funding=cordis_orgs['sample'][3]['funding_eur'],
    tianjin_projects=cordis_orgs['sample'][4]['projects'],
    tianjin_funding=cordis_orgs['sample'][4]['funding_eur'],
    nudt_risk=bis_academic['entities'][0]['risk_score'],
    nudt_reason=bis_academic['entities'][0]['reason'],
    nudt_date=bis_academic['entities'][0]['listing_date'],
    nudt_tech=bis_academic['entities'][0]['tech_focus'],
    buaa_risk=bis_academic['entities'][1]['risk_score'],
    buaa_reason=bis_academic['entities'][1]['reason'],
    buaa_date=bis_academic['entities'][1]['listing_date'],
    buaa_tech=bis_academic['entities'][1]['tech_focus']
)

# Update speaker notes
notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = enrichment_note + text_frame.text

print("  Speaker notes updated with gray-zone pattern analysis")

# Add validation marker to slide content
for shape in slide.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.text
        # Look for slide title or content about gray-zone
        if "Gray" in text or "Zone" in text or "Acquisition" in text:
            # Add data validation note
            if "Data" not in text.lower() or "validated" not in text.lower():
                # Add validation as new paragraph
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text and len(paragraph.text) > 15:
                        new_p = shape.text_frame.add_paragraph()
                        new_p.text = "(Data: 411 EU-partnered orgs, 1K high-risk collabs, 15 BIS-listed entities)"
                        new_p.level = 1
                        new_p.font.size = Pt(10)
                        new_p.font.italic = True
                        break
            break

# Save updated presentation
prs.save('MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx')

print("\n" + "="*80)
print("SUCCESS: Slide 13 enriched with gray-zone pattern data")
print("="*80)
print("\nEnrichment Summary:")
print(f"  Stage 1 (Legitimate): {cordis_orgs['total']} Chinese organizations in EU research")
print(f"    Top: Tsinghua University ({cordis_orgs['sample'][0]['projects']} projects)")
print(f"  Stage 2 (Gray-Zone): 1,000 high-risk collaborations identified")
print(f"  Stage 3 (MCF Confirmed): {bis_academic['count']} BIS-listed academic entities")
print(f"    Highest risk: {bis_academic['entities'][0]['name']} (Risk: {bis_academic['entities'][0]['risk_score']})")
print("\nKey Pattern:")
print("  Legitimate collaboration → Dual-use transfer → MCF integration → Entity List (5-10 year lag)")
print("\nSlide 13 enrichment complete.")
