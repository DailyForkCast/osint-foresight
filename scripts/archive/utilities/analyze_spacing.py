from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')

print("=" * 80)
print("DETAILED SPACING ANALYSIS - MCF v2.1 (Slides 3-8)")
print("=" * 80)
print()

# Analyze spacing properties
all_margins = {'left': [], 'right': [], 'top': [], 'bottom': []}
all_line_spacing = []
all_space_before = []
all_space_after = []
text_box_relationships = []

for slide_idx in range(2, 8):  # Slides 3-8
    slide = prs.slides[slide_idx]
    slide_num = slide_idx + 1

    print(f"SLIDE {slide_num}")
    print("-" * 80)

    text_boxes = []

    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            text_frame = shape.text_frame
            text = shape.text.strip()
            is_title = 'title' in shape.name.lower() if hasattr(shape, 'name') and shape.name else False

            # Text box info
            box_info = {
                'name': shape.name if hasattr(shape, 'name') else 'Unknown',
                'is_title': is_title,
                'text_preview': text[:50] + '...' if len(text) > 50 else text,
                'top': shape.top.inches,
                'bottom': shape.top.inches + shape.height.inches,
                'height': shape.height.inches,
            }

            # Text frame margins (internal padding)
            try:
                margins = {
                    'left': text_frame.margin_left.inches if text_frame.margin_left else 0,
                    'right': text_frame.margin_right.inches if text_frame.margin_right else 0,
                    'top': text_frame.margin_top.inches if text_frame.margin_top else 0,
                    'bottom': text_frame.margin_bottom.inches if text_frame.margin_bottom else 0,
                }
                box_info['margins'] = margins

                # Collect margins for statistics
                for key, value in margins.items():
                    if value is not None and value >= 0:
                        all_margins[key].append(value)

                print(f"\n  {box_info['name']} {'[TITLE]' if is_title else '[BODY]'}")
                print(f"    Text: \"{box_info['text_preview']}\"")
                print(f"    Internal Padding:")
                print(f"      Top:    {margins['top']:.3f}\"")
                print(f"      Bottom: {margins['bottom']:.3f}\"")
                print(f"      Left:   {margins['left']:.3f}\"")
                print(f"      Right:  {margins['right']:.3f}\"")

            except Exception as e:
                print(f"\n  {box_info['name']} - Could not read margins: {e}")

            # Paragraph spacing
            for para_idx, para in enumerate(text_frame.paragraphs):
                if para.text.strip():
                    try:
                        # Line spacing
                        if para.line_spacing is not None:
                            if hasattr(para.line_spacing, 'pt'):
                                line_spacing_val = para.line_spacing.pt
                            else:
                                line_spacing_val = para.line_spacing
                            all_line_spacing.append(line_spacing_val)
                            print(f"      Para {para_idx + 1} line spacing: {line_spacing_val}")

                        # Space before/after
                        if para.space_before is not None and hasattr(para.space_before, 'pt'):
                            space_before = para.space_before.pt
                            all_space_before.append(space_before)
                            print(f"      Para {para_idx + 1} space before: {space_before}pt")

                        if para.space_after is not None and hasattr(para.space_after, 'pt'):
                            space_after = para.space_after.pt
                            all_space_after.append(space_after)
                            print(f"      Para {para_idx + 1} space after: {space_after}pt")
                    except Exception as e:
                        pass

            text_boxes.append(box_info)

    # Analyze distance between text boxes (title to body spacing)
    text_boxes.sort(key=lambda x: x['top'])
    for i in range(len(text_boxes) - 1):
        current = text_boxes[i]
        next_box = text_boxes[i + 1]

        gap = next_box['top'] - current['bottom']

        if current['is_title'] and not next_box['is_title']:
            text_box_relationships.append({
                'type': 'title_to_body',
                'gap': gap,
                'slide': slide_num
            })
            print(f"\n  >>> TITLE-TO-BODY GAP: {gap:.3f}\" (Slide {slide_num})")
        elif not current['is_title'] and not next_box['is_title']:
            text_box_relationships.append({
                'type': 'body_to_body',
                'gap': gap,
                'slide': slide_num
            })

    print()

# Summary statistics
print("\n" + "=" * 80)
print("SUMMARY STATISTICS")
print("=" * 80)

print("\n1. TEXT FRAME INTERNAL MARGINS (Padding inside text boxes)")
print("-" * 80)
for side in ['top', 'bottom', 'left', 'right']:
    if all_margins[side]:
        avg = sum(all_margins[side]) / len(all_margins[side])
        min_val = min(all_margins[side])
        max_val = max(all_margins[side])
        print(f"  {side.capitalize():8s}: Avg={avg:.3f}\"  Min={min_val:.3f}\"  Max={max_val:.3f}\"  (n={len(all_margins[side])})")

print("\n2. LINE SPACING")
print("-" * 80)
if all_line_spacing:
    print(f"  Found {len(all_line_spacing)} explicit line spacing values")
    for val in set(all_line_spacing):
        count = all_line_spacing.count(val)
        print(f"    {val}: {count} instances")
else:
    print("  Most paragraphs use default line spacing (1.0 or auto)")

print("\n3. PARAGRAPH SPACING")
print("-" * 80)
if all_space_before:
    avg = sum(all_space_before) / len(all_space_before)
    print(f"  Space Before: Avg={avg:.1f}pt  (n={len(all_space_before)})")
    print(f"    Values: {sorted(set(all_space_before))}")
else:
    print("  Space Before: Using defaults (typically 0pt)")

if all_space_after:
    avg = sum(all_space_after) / len(all_space_after)
    print(f"  Space After:  Avg={avg:.1f}pt  (n={len(all_space_after)})")
    print(f"    Values: {sorted(set(all_space_after))}")
else:
    print("  Space After: Using defaults (typically 0pt)")

print("\n4. TEXT BOX RELATIONSHIPS")
print("-" * 80)
title_to_body_gaps = [r['gap'] for r in text_box_relationships if r['type'] == 'title_to_body']
if title_to_body_gaps:
    avg = sum(title_to_body_gaps) / len(title_to_body_gaps)
    print(f"  Title-to-Body Gap:")
    print(f"    Average: {avg:.3f}\"")
    print(f"    Min: {min(title_to_body_gaps):.3f}\"")
    print(f"    Max: {max(title_to_body_gaps):.3f}\"")
    print(f"    Instances: {len(title_to_body_gaps)}")

body_to_body_gaps = [r['gap'] for r in text_box_relationships if r['type'] == 'body_to_body']
if body_to_body_gaps:
    positive_gaps = [g for g in body_to_body_gaps if g > 0]
    if positive_gaps:
        avg = sum(positive_gaps) / len(positive_gaps)
        print(f"\n  Body-to-Body Gap (when separate boxes):")
        print(f"    Average: {avg:.3f}\"")
        print(f"    Min: {min(positive_gaps):.3f}\"")
        print(f"    Max: {max(positive_gaps):.3f}\"")
        print(f"    Instances: {len(positive_gaps)}")

print("\n" + "=" * 80)
print("RECOMMENDATIONS")
print("=" * 80)

# Generate recommendations
print("\nBased on this analysis:")
print()

if all_margins['top']:
    avg_top = sum(all_margins['top']) / len(all_margins['top'])
    print(f"1. TEXT BOX TOP PADDING: Use {avg_top:.3f}\" ({avg_top * 72:.0f}px)")

if all_margins['bottom']:
    avg_bottom = sum(all_margins['bottom']) / len(all_margins['bottom'])
    print(f"2. TEXT BOX BOTTOM PADDING: Use {avg_bottom:.3f}\" ({avg_bottom * 72:.0f}px)")
    if avg_bottom < 0.05:
        print(f"   ⚠️  WARNING: Bottom padding is very small ({avg_bottom:.3f}\")!")
        print(f"       This may cause text to appear too close to the box edge.")
        print(f"       RECOMMENDATION: Increase to at least 0.05\" (4px)")

if all_margins['left']:
    avg_left = sum(all_margins['left']) / len(all_margins['left'])
    print(f"3. TEXT BOX LEFT PADDING: Use {avg_left:.3f}\" ({avg_left * 72:.0f}px)")

if all_margins['right']:
    avg_right = sum(all_margins['right']) / len(all_margins['right'])
    print(f"4. TEXT BOX RIGHT PADDING: Use {avg_right:.3f}\" ({avg_right * 72:.0f}px)")

if title_to_body_gaps:
    avg_gap = sum(title_to_body_gaps) / len(title_to_body_gaps)
    print(f"\n5. TITLE-TO-BODY SPACING: Use {avg_gap:.3f}\" ({avg_gap * 72:.0f}px)")
    if avg_gap > 0.5:
        print(f"   ⚠️  WARNING: Gap is large ({avg_gap:.3f}\")!")
        print(f"       This may feel disconnected. Consider reducing to 0.2-0.3\"")
    elif avg_gap < 0.1:
        print(f"   ⚠️  WARNING: Gap is very small ({avg_gap:.3f}\")!")
        print(f"       This may feel cramped. Consider increasing to 0.2-0.3\"")

print("\n" + "=" * 80)
