from pptx import Presentation
from pptx.util import Pt, Inches

def rgb_to_hex(rgb):
    """Convert RGB tuple to hex"""
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')
slide = prs.slides[8]  # Slide 9

print("=" * 80)
print("SLIDE 9 - DETAILED TEXT FORMATTING ANALYSIS")
print("=" * 80)
print()

# Focus on specific text boxes
target_boxes = [
    "Domain Overlap",
    "Technology Transfer",
    "DSR", "GSI", "GDI", "GCI", "GAGI",
    "Text 30", "Text 33"  # Based on earlier analysis
]

for shape in slide.shapes:
    if not shape.has_text_frame:
        continue

    text = shape.text.strip()
    if not text:
        continue

    shape_name = shape.name if hasattr(shape, 'name') else 'Unknown'

    # Check if this is one of our target boxes
    is_target = any(target in shape_name or target in text[:50]
                    for target in target_boxes)

    if is_target or "Text 30" in shape_name or "Text 33" in shape_name or \
       "Smart Cities" in text or "Contracts" in text:

        print(f"\n{'=' * 80}")
        print(f"SHAPE: {shape_name}")
        print(f"Position: L={shape.left.inches:.2f}\", T={shape.top.inches:.2f}\", "
              f"W={shape.width.inches:.2f}\", H={shape.height.inches:.2f}\"")
        print(f"Text preview: {text[:80]}")
        print(f"{'=' * 80}")

        # Analyze text frame settings
        tf = shape.text_frame
        print(f"\nText Frame Settings:")
        print(f"  Word wrap: {tf.word_wrap}")
        print(f"  Auto size: {tf.auto_size}")

        # Analyze each paragraph and run
        for para_idx, para in enumerate(tf.paragraphs):
            if not para.text.strip():
                continue

            print(f"\n  Paragraph {para_idx + 1}:")
            print(f"    Text: \"{para.text.strip()}\"")
            print(f"    Alignment: {para.alignment}")
            print(f"    Level: {para.level}")

            # Analyze each run (formatted text segment)
            for run_idx, run in enumerate(para.runs):
                if not run.text.strip():
                    continue

                font = run.font

                # Get color
                color_info = "N/A"
                try:
                    if font.color.type == 1:  # RGB
                        rgb = font.color.rgb
                        color_info = f"{rgb_to_hex((rgb[0], rgb[1], rgb[2]))}"
                except:
                    pass

                print(f"      Run {run_idx + 1}: \"{run.text[:40]}\"")
                print(f"        Font: {font.name if font.name else 'N/A'}")
                print(f"        Size: {font.size.pt if font.size else 'N/A'}pt")
                print(f"        Color: {color_info}")
                print(f"        Bold: {font.bold}")
                print(f"        Italic: {font.italic}")

        # Check if text overflows
        print(f"\n  Text Length Analysis:")
        total_chars = len(text)
        print(f"    Total characters: {total_chars}")
        print(f"    Box height: {shape.height.inches:.3f}\"")
        if total_chars > 100:
            print(f"    WARNING: Long text in relatively small box")

print("\n" + "=" * 80)
print("ANALYSIS COMPLETE")
print("=" * 80)
