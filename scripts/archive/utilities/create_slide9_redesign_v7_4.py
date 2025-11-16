from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor

# Create new presentation with correct dimensions
prs = Presentation()
prs.slide_width = Inches(10.0)
prs.slide_height = Inches(5.62)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Color palette
COLORS = {
    'white': RGBColor(255, 255, 255),
    'orange': RGBColor(243, 156, 18),
    'red': RGBColor(200, 16, 46),
    'blue': RGBColor(46, 107, 168),
    'green': RGBColor(39, 174, 96),
    'alt_orange': RGBColor(230, 126, 34),
    'gray': RGBColor(149, 165, 166),
    'dark_slate': RGBColor(44, 62, 80)
}

# Set slide background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLORS['dark_slate']

def create_text_box(slide, left, top, width, height, text, font_size,
                    color=COLORS['white'], bold=False, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT,
                    vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
                    tight_spacing=False):
    """Create text box with proper settings"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    text_frame = textbox.text_frame
    text_frame.margin_top = Inches(0.08)
    text_frame.margin_bottom = Inches(0.10)
    text_frame.margin_left = Inches(0.10)
    text_frame.margin_right = Inches(0.10)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = vertical_anchor

    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = alignment

    # Tight spacing for long names that wrap
    if tight_spacing:
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.15  # Tighter than standard 1.3
    else:
        paragraph.space_before = Pt(6)
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = 1.3

    run = paragraph.runs[0]
    run.font.name = 'Arial'
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold

    return textbox

def create_label_description_box(slide, left, top, width, height, items, font_size, text_color, bg_color=None):
    """Create label:description text box"""
    if bg_color:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))

    text_frame = textbox.text_frame
    text_frame.margin_top = Inches(0.08)
    text_frame.margin_bottom = Inches(0.10)
    text_frame.margin_left = Inches(0.10)
    text_frame.margin_right = Inches(0.10)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    text_frame.clear()

    for idx, (label, description) in enumerate(items):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        p.line_spacing = 1.3

        run_label = p.add_run()
        run_label.text = label + ":"
        run_label.font.name = 'Arial'
        run_label.font.size = Pt(font_size)
        run_label.font.color.rgb = text_color
        run_label.font.bold = True

        run_desc = p.add_run()
        run_desc.text = " " + description
        run_desc.font.name = 'Arial'
        run_desc.font.size = Pt(font_size)
        run_desc.font.color.rgb = text_color
        run_desc.font.bold = False

    return textbox

def create_colored_box(slide, left, top, width, height, fill_color):
    """Create colored rectangle"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# SLIDE 9 V7.4 - FIXED STRATEGIC ARCHITECTURE CENTERING

# Title background
create_colored_box(slide, 0, 0, 10, 0.6, COLORS['red'])

# Title - vertically centered (MIDDLE alignment)
create_text_box(slide, 0.2, 0, 9.6, 0.6,
                "BRI + FIVE INITIATIVES: GLOBAL MCF INFRASTRUCTURE",
                19.5, COLORS['white'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

# Section header - MIDDLE alignment (not TOP!)
create_text_box(slide, 0.3, 0.7, 4.5, 0.25,
                "Strategic Architecture",
                14, COLORS['red'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

# BRI section - BOTH ELEMENTS MIDDLE ALIGNED
bri_y = 1.05
bri_width = 9.4
bri_height = 0.45

# Orange background for entire BRI section
create_colored_box(slide, 0.3, bri_y, bri_width, bri_height, COLORS['orange'])

# BRI Title - MIDDLE aligned
create_text_box(slide, 0.3, bri_y, bri_width, 0.22,
                "Belt and Road Initiative (BRI)",
                11, COLORS['white'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.LEFT,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

# BRI Description - MIDDLE aligned
create_text_box(slide, 0.3, bri_y + 0.22, bri_width, 0.23,
                "Physical + digital infrastructure creating dependency and access points",
                9, COLORS['white'], alignment=PP_PARAGRAPH_ALIGNMENT.LEFT,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

# FIVE INITIATIVES - V7.1 COMPACT LAYOUT WITH SPECIAL FORMATTING
initiatives_y = 1.67
initiative_width = 1.72
initiative_spacing = 1.88

# V7 dimensions - compact
name_height = 0.35
desc_height = 0.50

initiatives = [
    {'acronym': 'DSR', 'name': 'Digital Silk Road',
     'desc': 'Data infrastructure, smart cities, e-commerce', 'color': COLORS['blue'],
     'tight_spacing': False},  # Short name, normal spacing
    {'acronym': 'GSI', 'name': 'Global Security Initiative',
     'desc': 'Security frameworks, peacekeeping', 'color': COLORS['green'],
     'tight_spacing': False},  # Short name, normal spacing
    {'acronym': 'GDI', 'name': 'Global Development Initiative',
     'desc': 'Development norms, infrastructure', 'color': COLORS['orange'],
     'tight_spacing': True},  # LONG NAME - tight spacing
    {'acronym': 'GCI', 'name': 'Global Civilization Initiative',
     'desc': 'Cultural influence, Confucius Institutes', 'color': COLORS['red'],
     'tight_spacing': False},  # Short name, normal spacing
    {'acronym': 'GAGI', 'name': 'Global AI Governance Initiative',
     'desc': 'AI standards, data sovereignty', 'color': COLORS['alt_orange'],
     'tight_spacing': True}  # LONG NAME - tight spacing
]

# STEP 1: Create all background boxes first (back layer)
for i, init in enumerate(initiatives):
    x_pos = 0.4 + (i * initiative_spacing)

    # Colored header box
    create_colored_box(slide, x_pos, initiatives_y, initiative_width, 0.15, init['color'])

    # Description white background box
    create_colored_box(slide, x_pos, initiatives_y + 0.56, initiative_width, desc_height, COLORS['white'])

# STEP 2: Create all text boxes (front layer)
for i, init in enumerate(initiatives):
    x_pos = 0.4 + (i * initiative_spacing)

    # Acronym - in colored box with MIDDLE alignment
    create_text_box(slide, x_pos, initiatives_y, initiative_width, 0.15,
                    init['acronym'],
                    11, COLORS['white'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                    vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

    # Full name - with tight_spacing flag for long names, TOP alignment
    create_text_box(slide, x_pos, initiatives_y + 0.16, initiative_width, name_height,
                    init['name'],
                    9, COLORS['white'], alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                    vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
                    tight_spacing=init['tight_spacing'])

    # Description text - TOP alignment
    create_text_box(slide, x_pos, initiatives_y + 0.56, initiative_width, desc_height,
                    init['desc'],
                    9, init['color'], alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                    vertical_anchor=MSO_VERTICAL_ANCHOR.TOP)

# Bottom sections
# initiatives end at: 1.67 + 0.56 + 0.50 = 2.73
# Leave 0.30" gap: 2.73 + 0.30 = 3.03
bottom_y = 3.03

# Left box
create_colored_box(slide, 0.3, bottom_y, 3.2, 0.3, COLORS['blue'])
create_text_box(slide, 0.35, bottom_y, 3.1, 0.3,
                "Domain Overlap Zones",
                11, COLORS['white'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

overlap_items = [
    ("Smart Cities", "DSR + GDI integration"),
    ("Maritime", "GSI + BRI ports"),
    ("Standards", "DSR + GSI + GAGI")
]
create_label_description_box(slide, 0.3, bottom_y + 0.32, 3.2, 1.5,
                             overlap_items, 10, COLORS['dark_slate'], COLORS['white'])

# Right box
create_colored_box(slide, 3.7, bottom_y, 3.2, 0.3, COLORS['green'])
create_text_box(slide, 3.75, bottom_y, 3.1, 0.3,
                "Technology Transfer Mechanisms",
                11, COLORS['white'], bold=True, alignment=PP_PARAGRAPH_ALIGNMENT.CENTER,
                vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)

transfer_items = [
    ("Contracts", "Specify Chinese tech (Vendor lock-in)"),
    ("Lock-in", "Proprietary standards create dependencies"),
    ("Training", "Create technical dependencies")
]
create_label_description_box(slide, 3.7, bottom_y + 0.32, 3.2, 1.5,
                             transfer_items, 10, COLORS['dark_slate'], COLORS['white'])

# Save
output_path = 'C:/Projects/OSINT-Foresight/MCF_Slide9_Redesign_v7_4_FINAL.pptx'
prs.save(output_path)

print("=" * 80)
print("SLIDE 9 V7.4 - FIXED STRATEGIC ARCHITECTURE CENTERING")
print("=" * 80)
print()
print("FIXES from V7.3:")
print("  [X] 'Strategic Architecture': TOP -> MIDDLE vertical alignment")
print("  [X] Text now properly centered in its box")
print()
print("ALL HEADERS AND SINGLE-LINE ELEMENTS NOW USE MIDDLE:")
print("  - Main title: 'BRI + FIVE INITIATIVES...' (MIDDLE)")
print("  - Section header: 'Strategic Architecture' (MIDDLE)")
print("  - BRI title: 'Belt and Road Initiative (BRI)' (MIDDLE)")
print("  - BRI description: 'Physical + digital infrastructure...' (MIDDLE)")
print("  - Initiative acronyms: DSR, GSI, GDI, GCI, GAGI (MIDDLE)")
print("  - Bottom section titles: 'Domain Overlap Zones', etc. (MIDDLE)")
print()
print("MULTI-LINE CONTENT USES TOP:")
print("  - Initiative names: May wrap to 2 lines (TOP)")
print("  - Initiative descriptions: May wrap to 2 lines (TOP)")
print("  - Bottom section content: Label:description pairs (TOP)")
print()
print("Simple Rule:")
print("  Single-line headers/titles/labels = MIDDLE")
print("  Multi-line flowing content = TOP")
print()
print(f"Output: {output_path}")
print("=" * 80)
