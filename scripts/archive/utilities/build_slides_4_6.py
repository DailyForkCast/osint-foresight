# -*- coding: utf-8 -*-
"""Build slides 4-6 with full detail for MCF presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import math

# Load existing presentation
prs = Presentation("C:/Projects/OSINT - Foresight/MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")

# Define colors
NAVY = RGBColor(15, 25, 45)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
DARK_GRAY = RGBColor(60, 60, 70)
LIGHT_GRAY = RGBColor(150, 150, 160)

def add_dark_background(slide):
    """Add navy background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def add_notes(slide, notes_text):
    """Add speaker notes to slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

print("Building Slides 4-6 with full detail...")
print("=" * 70)

# ===== SLIDE 4: CHINA'S MOTIVATIONS & LEGAL FOUNDATIONS =====
print("Creating Slide 4: China's Motivations - Concentric Circles")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "China's Motivations & Legal Foundations"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Concentric circles (center of slide)
center_x = 5
center_y = 3.75

# Outer circle: Economic Modernization
outer_circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(center_x - 2.5), Inches(center_y - 2),
    Inches(5), Inches(4)
)
outer_circle.fill.solid()
outer_circle.fill.fore_color.rgb = RGBColor(80, 60, 60)
outer_circle.line.color.rgb = GOLD
outer_circle.line.width = Pt(2)

# Label for outer
outer_label = slide.shapes.add_textbox(Inches(center_x - 1.2), Inches(center_y + 1.6), Inches(2.4), Inches(0.4))
tf = outer_label.text_frame
tf.text = "Economic Modernization"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Middle circle: Tech Leadership
mid_circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(center_x - 1.8), Inches(center_y - 1.4),
    Inches(3.6), Inches(2.8)
)
mid_circle.fill.solid()
mid_circle.fill.fore_color.rgb = RGBColor(100, 70, 70)
mid_circle.line.color.rgb = GOLD
mid_circle.line.width = Pt(2)

# Label for middle
mid_label = slide.shapes.add_textbox(Inches(center_x - 1), Inches(center_y + 0.8), Inches(2), Inches(0.4))
tf = mid_label.text_frame
tf.text = "Tech Leadership"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Inner circle: National Power/PLA
inner_circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(center_x - 1.2), Inches(center_y - 0.9),
    Inches(2.4), Inches(1.8)
)
inner_circle.fill.solid()
inner_circle.fill.fore_color.rgb = RGBColor(180, 50, 50)
inner_circle.line.color.rgb = GOLD
inner_circle.line.width = Pt(3)

# Label for core
core_label = slide.shapes.add_textbox(Inches(center_x - 1), Inches(center_y - 0.3), Inches(2), Inches(0.6))
tf = core_label.text_frame
tf.text = "National Power\nPLA Modernization"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Legal badges on the right
legal_frameworks = [
    ("National Security\nLaw (2015)", 1.2),
    ("National Intelligence\nLaw (2017, Art.7)", 2.2),
    ("Data Security\nLaw (2021)", 3.2),
    ("Revised State Secrets\nLaw (2024)", 4.2)
]

for law_text, y_pos in legal_frameworks:
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8), Inches(y_pos),
        Inches(1.8), Inches(0.8)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_GRAY
    badge.line.color.rgb = GOLD
    badge.line.width = Pt(2)

    tf = badge.text_frame
    tf.text = law_text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_notes(slide, """China's motivations for MCF/NQPF are concentric and mutually reinforcing:

CORE: National Power and PLA Modernization
The innermost objective is enhancing comprehensive national power, with particular emphasis on military capabilities. The PLA's modernization goals drive technology requirements across all domains.

MIDDLE RING: Technology Leadership
China seeks technological self-sufficiency and dominance in strategic sectors. This includes dual-use technologies that serve both civilian and military applications.

OUTER RING: Economic Modernization
The broadest objective is elevating China's economic competitiveness through innovation, industrial upgrading, and integration into global technology supply chains.

LEGAL FRAMEWORKS COMPELLING COOPERATION:

National Security Law (2015): Establishes broad state authority over all entities and individuals in matters affecting national security. All organizations must support national security work.

National Intelligence Law (2017), Article 7: Explicitly requires "all organizations and citizens" to support, assist, and cooperate with intelligence work. This creates legal compulsion for technology transfer.

Data Security Law (2021): Mandates that data processing activities must not harm national security. Gives the state authority to access data held by any entity for national security purposes.

Revised State Secrets Law (2024): Expands the definition of state secrets and strengthens penalties. Relevant to technology information that could be classified retroactively.

IMPLICATIONS: These laws mean that any Chinese entity—whether ostensibly private, academic, or commercial—can be compelled to cooperate with state intelligence and military objectives. Western partners engaging with Chinese entities must understand that legal compulsion overrides contractual obligations or institutional autonomy.""")

print("[OK] Slide 4 complete")

# ===== SLIDE 5: MECHANISM INSIDE CHINA =====
print("Creating Slide 5: Mechanism Inside China - Flow Diagram")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "MCF Mechanism Inside China"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Flow diagram boxes
flow_boxes = [
    ("Research Base", "Universities\nNational Labs\nCAS Institutes", 1.5, 2),
    ("Commercial Bridge", "Tech Companies\nStartups\nIncubators", 3.5, 2),
    ("Industrial Hubs", "Manufacturing\nSOEs\nPrivate Firms", 5.5, 2),
    ("Defense Integration", "PLA\nDefense Industry\nMilitary Academies", 7.5, 2),
]

for i, (title, content, x_pos, y_pos) in enumerate(flow_boxes):
    # Main box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_pos), Inches(y_pos),
        Inches(1.6), Inches(1.4)
    )

    # Color gradient from research (blue) to defense (red)
    if i == 0:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(70, 100, 150)
    elif i == 1:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(100, 100, 100)
    elif i == 2:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(120, 90, 70)
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(160, 60, 60)

    box.line.color.rgb = GOLD
    box.line.width = Pt(2)

    # Title
    title_tf = box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Content
    p2 = title_tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(10)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    # Arrow to next (except last)
    if i < len(flow_boxes) - 1:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(x_pos + 1.7), Inches(y_pos + 0.6),
            Inches(0.7), Inches(0.3)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.width = Pt(0)

# Feedback loop arrow (curved back from Defense to Research)
feedback_label = slide.shapes.add_textbox(Inches(7.5), Inches(3.6), Inches(1.6), Inches(0.4))
tf = feedback_label.text_frame
tf.text = "Feedback Loop"
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = GOLD
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Curved line back (simplified as straight with annotation)
feedback_line = slide.shapes.add_connector(
    1, Inches(8), Inches(3.8), Inches(2.5), Inches(3.8)
)
feedback_line.line.color.rgb = GOLD
feedback_line.line.width = Pt(2)
feedback_line.line.dash_style = 2  # Dashed

# Arrow at end
arrow_end = slide.shapes.add_shape(
    MSO_SHAPE.LEFT_ARROW,
    Inches(2.2), Inches(3.65),
    Inches(0.4), Inches(0.3)
)
arrow_end.fill.solid()
arrow_end.fill.fore_color.rgb = GOLD
arrow_end.line.width = Pt(0)

# Bottom annotation
annotation = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1))
tf = annotation.text_frame
tf.text = "Domestic Loop: Civilian innovation → Commercial application → Industrial scale → Military integration → Requirements feedback"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(slide, """The MCF mechanism inside China operates as a continuous loop:

RESEARCH BASE:
Universities, national laboratories, and Chinese Academy of Sciences (CAS) institutes conduct foundational research. These institutions receive military R&D funding while maintaining civilian academic facades. Key talent programs (Thousand Talents, Changjiang Scholars) recruit foreign expertise.

COMMERCIAL BRIDGE:
Technology companies and startups serve as the primary mechanism for converting research into products. Companies like Huawei, ZTE, DJI, and SenseTime straddle civilian and military applications. Government-backed incubators and venture funds provide capital with implicit dual-use expectations.

INDUSTRIAL HUBS:
Manufacturing at scale occurs through state-owned enterprises (SOEs) and nominally private firms. Industrial policy tools (Made in China 2025, National IC Fund) direct production capacity toward strategic technologies. Supply chain integration ensures military access to civilian manufacturing.

DEFENSE INTEGRATION:
The PLA and defense industry directly incorporate commercial innovations. Defense conglomerates (CASIC, AVIC, CETC, Norinco) have civilian subsidiaries that participate in commercial markets while maintaining military production lines. Military academies conduct joint research with civilian universities.

FEEDBACK LOOP:
Military requirements flow back to research institutions, creating a continuous cycle. The PLA identifies capability gaps, research institutes propose solutions, companies productize innovations, industry scales production, and defense integrates new capabilities. This feedback mechanism accelerates the civil-military innovation cycle.

The entire system is coordinated through party committees, joint civil-military R&D programs, and government procurement that deliberately blurs civilian-military boundaries.""")

print("[OK] Slide 5 complete")

# ===== SLIDE 6: TERMINOLOGY SHIFT =====
print("Creating Slide 6: Terminology Shift - Dual Line Chart")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Terminology Shift: MCF → NQPF (Optics ≠ Operations)"
p = tf.paragraphs[0]
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = WHITE

# Create chart data
chart_data = CategoryChartData()
chart_data.categories = ['2021', '2022', '2023', '2024', '2025']
chart_data.add_series('MCF Keywords', (120, 110, 70, 40, 25))
chart_data.add_series('NQPF Keywords', (5, 10, 45, 80, 95))

# Add chart
x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

# Format chart
chart.has_legend = True
chart.legend.position = 2  # Right
chart.legend.include_in_layout = False

# Format series
mcf_series = chart.series[0]
mcf_series.format.line.color.rgb = RGBColor(200, 100, 100)
mcf_series.format.line.width = Pt(3)

nqpf_series = chart.series[1]
nqpf_series.format.line.color.rgb = GOLD
nqpf_series.format.line.width = Pt(3)

# Add vertical marker line for 2023 transition
marker_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.2), Inches(1.5), Inches(0.4))
tf = marker_box.text_frame
tf.text = "Sep 2023:\nNQPF Introduced"
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = GOLD
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Annotation
annotation = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(8), Inches(0.8))
tf = annotation.text_frame
tf.text = "Language shift does not indicate policy change—legal frameworks and objectives remain consistent.\nPlaceholder data shown; replace with actual keyword frequency analysis when available."
tf.paragraphs[0].font.size = Pt(11)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(slide, """The terminology shift from Military-Civil Fusion (MCF) to New Quality Productive Forces (NQPF) represents a rhetorical adjustment, not a substantive policy change.

PLACEHOLDER DATA NOTE: The chart shows illustrative keyword frequency trends. This should be replaced with actual analysis from:
- Chinese policy documents (State Council, NDRC, MOST publications)
- Official media (People's Daily, Xinhua, CCTV transcripts)
- Academic literature (CNKI database queries)
- Party Congress and NPC reports

CONTINUITY OF INTENT:
Despite the language shift, the core objectives remain:
1. Dual-use technology development serving both economic and military goals
2. Legal compulsion for all entities to support national security objectives
3. Systematic acquisition of foreign technology and expertise
4. Integration of civilian innovation into military capabilities

WHY THE SHIFT MATTERS:
The rebranding to "New Quality Productive Forces" may:
- Reduce international scrutiny by using less militarized language
- Make Chinese technology partnerships appear more benign to foreign partners
- Complicate export control and screening processes that explicitly reference "MCF"
- Create challenges for policy communication ("What is NQPF?" vs. well-understood "MCF")

OPERATIONAL REALITY:
The institutional mechanisms, legal frameworks, and acquisition pathways established under MCF continue under NQPF. The Commission for Integrated Military-Civil Development still operates, talent programs persist, and legal obligations for technology transfer remain in force.

ANALYTICAL IMPLICATION: Analysts and practitioners must look beyond terminology to assess actual dual-use risks. Focusing solely on whether a Chinese entity explicitly mentions "MCF" or "NQPF" will miss the underlying compulsion structure.""")

print("[OK] Slide 6 complete")

# Save progress
print("\nSaving slides 4-6...")
prs.save("C:/Projects/OSINT - Foresight/MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")
print("Progress saved. Slides 1-6 now complete.")
print("=" * 70)
