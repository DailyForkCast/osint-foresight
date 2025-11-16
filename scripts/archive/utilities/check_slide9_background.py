from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')

# Slide 9 (index 8)
slide = prs.slides[8]

print("=" * 80)
print("SLIDE 9 BACKGROUND ANALYSIS")
print("=" * 80)
print()

# Check slide background
try:
    bg = slide.background
    fill = bg.fill

    print(f"Background fill type: {fill.type}")

    if fill.type == MSO_FILL_TYPE.SOLID:
        print("Background: Solid fill")
        try:
            color = fill.fore_color
            if color.type == 1:  # RGB
                rgb = color.rgb
                print(f"  RGB: ({rgb[0]}, {rgb[1]}, {rgb[2]})")
                print(f"  Hex: #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
        except Exception as e:
            print(f"  Could not extract color: {e}")
    elif fill.type == MSO_FILL_TYPE.PATTERNED:
        print("Background: Patterned")
    elif fill.type == MSO_FILL_TYPE.GRADIENT:
        print("Background: Gradient")
    elif fill.type == MSO_FILL_TYPE.PICTURE:
        print("Background: Picture")
    elif fill.type == MSO_FILL_TYPE.TEXTURED:
        print("Background: Textured")
    else:
        print(f"Background: Other type ({fill.type})")

except Exception as e:
    print(f"Could not read background: {e}")
    print("\nChecking for background shapes instead...")

# Check for background rectangles
print("\n" + "=" * 80)
print("CHECKING FOR BACKGROUND SHAPES")
print("=" * 80)

for idx, shape in enumerate(slide.shapes):
    # Look for large rectangles that might be backgrounds
    if hasattr(shape, 'width') and hasattr(shape, 'height'):
        if shape.width.inches >= 9.5 and shape.height.inches >= 5.0:
            print(f"\nLarge shape found (possible background): {shape.name}")
            print(f"  Position: {shape.left.inches:.2f}\", {shape.top.inches:.2f}\"")
            print(f"  Size: {shape.width.inches:.2f}\" x {shape.height.inches:.2f}\"")
            print(f"  Shape type: {shape.shape_type}")

            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type == MSO_FILL_TYPE.SOLID:
                        color = shape.fill.fore_color
                        if color.type == 1:
                            rgb = color.rgb
                            print(f"  Fill color RGB: ({rgb[0]}, {rgb[1]}, {rgb[2]})")
                            print(f"  Fill color Hex: #{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}")
                except:
                    pass

print("\n" + "=" * 80)
