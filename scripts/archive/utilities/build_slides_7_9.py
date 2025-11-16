# -*- coding: utf-8 -*-
"""Build slides 7-9 with full detail for MCF presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import math

# Load existing presentation
prs = Presentation("C:/Projects/OSINT - Foresight/MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")

# Define colors
NAVY = RGBColor(15, 25, 45)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
DARK_GRAY = RGBColor(60, 60, 70)
LIGHT_GRAY = RGBColor(150, 150, 160)

def add_dark_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

print("Building Slides 7-9 with full detail...")
print("=" * 70)

# ===== SLIDE 7: DUAL-USE DOMAINS =====
print("Creating Slide 7: Dual-Use Domains - 6-Hex Cluster")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Dual-Use Technology Domains → Capacity Gaps"
p = tf.paragraphs[0]
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = WHITE

# Center node
center_x, center_y = 5, 3.75
center_node = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    Inches(center_x - 0.6), Inches(center_y - 0.5),
    Inches(1.2), Inches(1)
)
center_node.fill.solid()
center_node.fill.fore_color.rgb = RGBColor(180, 50, 50)
center_node.line.color.rgb = GOLD
center_node.line.width = Pt(3)

center_tf = center_node.text_frame
center_tf.text = "Dual-use\n→ PLA/NQPF"
center_tf.paragraphs[0].font.size = Pt(12)
center_tf.paragraphs[0].font.bold = True
center_tf.paragraphs[0].font.color.rgb = WHITE
center_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
center_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# 6 hexagons around center
hex_domains = [
    ("AI/ML", "Model-weight\ncontrols", 0),
    ("Semiconductors", "Sub-tier\nmapping", 60),
    ("Quantum", "Research\nscreening", 120),
    ("Space", "Facility\noversight", 180),
    ("Biotech", "Genomic\ngovernance", 240),
    ("Advanced\nMaterials", "Process-IP\nprotection", 300)
]

for title, gap, angle_deg in hex_domains:
    # Calculate position
    radius = 2.2
    angle_rad = math.radians(angle_deg)
    hex_x = center_x + radius * math.cos(angle_rad)
    hex_y = center_y + radius * math.sin(angle_rad)

    # Create hexagon (using rectangle as placeholder for complexity)
    hex_box = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(hex_x - 0.7), Inches(hex_y - 0.55),
        Inches(1.4), Inches(1.1)
    )
    hex_box.fill.solid()
    hex_box.fill.fore_color.rgb = DARK_GRAY
    hex_box.line.color.rgb = GOLD
    hex_box.line.width = Pt(2)

    # Text
    hex_tf = hex_box.text_frame
    p = hex_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = hex_tf.add_paragraph()
    p2.text = gap
    p2.font.size = Pt(9)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    # Connector line to center
    connector = slide.shapes.add_connector(
        1, Inches(hex_x), Inches(hex_y),
        Inches(center_x), Inches(center_y)
    )
    connector.line.color.rgb = RGBColor(100, 100, 120)
    connector.line.width = Pt(1)

add_notes(slide, """Six critical dual-use technology domains and their associated capacity gaps:

AI/ML (Model-Weight Controls):
Capacity Gap: Lack of technical mechanisms to prevent dual-use AI model proliferation. Model weights can be copied and fine-tuned for military applications (autonomous weapons, surveillance, command-control systems). Current export controls focus on chips, not trained models.

Semiconductors (Sub-Tier Mapping):
Capacity Gap: Opaque supply chains beyond first-tier suppliers. Chinese companies can access advanced chips through distributors, shell companies, or third-country transshipment. Most entities lack tools to map sub-tier sources and ultimate end-users.

Quantum (Research Screening):
Capacity Gap: Academic institutions inadequately screen quantum research collaborations. Many lack awareness that quantum sensing, communication, and computing have direct PLA applications (submarine detection, secure military communications, cryptography).

Space (Facility Oversight):
Capacity Gap: Dual-use ground stations (like Argentina's Neuquen deep-space facility) often approved without considering military applications (satellite control, space surveillance, anti-satellite operations). Weak interagency coordination between space agencies and defense/intelligence communities.

Biotech (Genomic Governance):
Capacity Gap: Inadequate frameworks for genomic data sharing and synthetic biology research. Chinese military-civil fusion entities (BGI, military medical universities) can access foreign biobanks, patient data, and gain-of-function research. Biosecurity protocols lag behind technology.

Advanced Materials (Process-IP Protection):
Capacity Gap: Process knowledge for advanced materials (carbon fiber, superalloys, composites) often not protected as rigorously as product IP. Joint ventures and technology partnerships can transfer manufacturing expertise critical for aerospace and defense applications.

CROSS-CUTTING CAPACITY NEED: Each domain requires specialized technical expertise to assess dual-use risks—expertise that most due diligence, research security, and export control teams currently lack.""")

print("[OK] Slide 7 complete")

# ===== SLIDE 8: CASE STUDIES =====
print("Creating Slide 8: Case Studies - Domestic Integration")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Case Studies: Domestic MCF Integration"
p = tf.paragraphs[0]
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = WHITE

# 2x3 grid of case boxes
cases = [
    ("SenseTime", "AI/Computer Vision", "Facial recognition → Public security → Military surveillance systems"),
    ("Megvii (Face++)", "AI/Identity Recognition", "Commercial API → Xinjiang deployment → PLA applications"),
    ("BGI Genomics", "Biotechnology", "Global sequencing services → PLA research → Biowarfare concerns"),
    ("USTC", "Quantum Research", "Academic research → Quantum comms → Military-secure networks"),
    ("CASIC", "Defense Conglomerate", "Missile systems ← Commercial drones ← Civilian tech integration")
]

# 2 columns, 3 rows (5 cases + 1 legend)
for idx, (company, domain, flow) in enumerate(cases):
    col = idx % 2
    row = idx // 2

    x_pos = 1 + (col * 4.2)
    y_pos = 1.2 + (row * 1.8)

    case_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_pos), Inches(y_pos),
        Inches(3.8), Inches(1.5)
    )
    case_box.fill.solid()
    case_box.fill.fore_color.rgb = DARK_GRAY
    case_box.line.color.rgb = GOLD
    case_box.line.width = Pt(2)

    case_tf = case_box.text_frame

    # Company name
    p = case_tf.paragraphs[0]
    p.text = company
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.LEFT

    # Domain
    p2 = case_tf.add_paragraph()
    p2.text = domain
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(3)

    # Flow
    p3 = case_tf.add_paragraph()
    p3.text = flow
    p3.font.size = Pt(10)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(6)

# Legend box (6th position)
legend_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(5.2), Inches(4.8),
    Inches(3.8), Inches(1.5)
)
legend_box.fill.solid()
legend_box.fill.fore_color.rgb = RGBColor(40, 50, 65)
legend_box.line.color.rgb = WHITE
legend_box.line.width = Pt(1)

legend_tf = legend_box.text_frame
p = legend_tf.paragraphs[0]
p.text = "Pattern: Civil → Military"
p.font.size = Pt(12)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = legend_tf.add_paragraph()
p2.text = "Commercial technology developed with global investment and partnerships, then integrated into military/security applications. Often involves BIS Entity List designation after military use is documented."
p2.font.size = Pt(9)
p2.font.color.rgb = LIGHT_GRAY
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(6)

add_notes(slide, """Five exemplar cases of China's domestic MCF integration:

SENSETIME:
Commercial origin: Founded 2014 as AI/computer vision startup, received investments from Alibaba, SoftBank, Qualcomm.
Military transition: Facial recognition technology deployed for public security, then integrated into PLA surveillance systems for border monitoring and military base security.
Designation: Added to BIS Entity List (2019) for Xinjiang human rights violations, then Treasury SDN List (2021) for military-industrial complex ties.
Capacity lesson: Foreign investors and technology partners failed to conduct adequate human rights and end-use due diligence.

MEGVII (Face++):
Commercial origin: AI facial recognition company, provided APIs for Alipay, Lenovo, and global customers.
Military transition: Deployed in Xinjiang surveillance network, then provided systems to PLA for military applications.
Designation: Added to BIS Entity List (2019).
Capacity lesson: API-based business models can obscure end-use, requiring deeper supply chain visibility.

BGI GENOMICS:
Commercial origin: World's largest genomics sequencing provider, partnerships with academic institutions globally.
Military transition: PLA military medical universities co-authored research; BGI participated in PLA genetics studies with dual-use biowarfare implications.
Designation: Subsidiaries added to BIS Entity List (2020, 2021).
Capacity lesson: Genomic data sharing agreements lacked biosecurity provisions; bioethics review insufficient for dual-use genetics research.

USTC (University of Science & Technology of China):
Academic status: Prestigious university conducting quantum research.
Military links: Operates under Chinese Academy of Sciences with direct PLA ties. Quantum communication research directly supports military-secure networks. Key researchers have PLA affiliations.
Designation: Not currently designated, but frequently cited in open-source reporting on MCF.
Capacity lesson: Academic prestige can obscure military connections; need systematic research security screening.

CASIC (China Aerospace Science & Industry Corporation):
Origin: State-owned defense conglomerate producing missiles, space systems.
Civil-military flow: Reverse integration—CASIC subsidiaries produce commercial drones and civilian technologies, incorporating innovations from private sector and foreign partnerships into military systems.
Designation: Multiple subsidiaries on various lists.
Capacity lesson: Defense conglomerates with civilian subsidiaries create complex webs requiring entity-resolution and ultimate-ownership analysis.

COMMON THREAD: All cases show how entities can operate in commercial, academic, or civilian contexts while serving military modernization. Foreign partners often discover military connections only after investment, collaboration, or technology transfer.""")

print("[OK] Slide 8 complete")

# ===== SLIDE 9: BRI/DSR GLOBALIZATION =====
print("Creating Slide 9: BRI/DSR Globalization")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Globalization: BRI/DSR as Dual-Use Tech Transfer"
p = tf.paragraphs[0]
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE

# Definition cards on left
bri_card = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(1.2),
    Inches(4), Inches(1.5)
)
bri_card.fill.solid()
bri_card.fill.fore_color.rgb = DARK_GRAY
bri_card.line.color.rgb = GOLD
bri_card.line.width = Pt(2)

bri_tf = bri_card.text_frame
p = bri_tf.paragraphs[0]
p.text = "Belt and Road Initiative (BRI, 2013–)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.LEFT

p2 = bri_tf.add_paragraph()
p2.text = "Infrastructure projects facilitating technology and data flows. Includes telecom networks, data centers, smart cities, and transport corridors that enable dual-use capabilities."
p2.font.size = Pt(10)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(6)

dsr_card = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.5), Inches(3),
    Inches(4), Inches(1.5)
)
dsr_card.fill.solid()
dsr_card.fill.fore_color.rgb = DARK_GRAY
dsr_card.line.color.rgb = GOLD
dsr_card.line.width = Pt(2)

dsr_tf = dsr_card.text_frame
p = dsr_tf.paragraphs[0]
p.text = "Digital Silk Road (DSR, 2015–)"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.LEFT

p2 = dsr_tf.add_paragraph()
p2.text = "ICT and data-driven initiatives enabling dual-use control and surveillance capabilities. Includes 5G networks, cloud infrastructure, fiber-optic cables, and satellite systems."
p2.font.size = Pt(10)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.LEFT
p2.space_before = Pt(6)

# Minimalist world map representation (simplified continents)
# Draw simplified outlines as reference points
map_regions = [
    ("Africa", 6.2, 4),
    ("Europe", 5.8, 2.5),
    ("Asia", 7.5, 3),
    ("Americas", 2.8, 3.2),
    ("Oceania", 8.5, 5)
]

for region, x, y in map_regions:
    # Simple circle as region marker
    marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x), Inches(y),
        Inches(0.25), Inches(0.25)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = RGBColor(100, 100, 120)
    marker.line.width = Pt(0)

    # Label
    label = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + 0.3), Inches(0.8), Inches(0.3))
    tf = label.text_frame
    tf.text = region
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Pins for BRI/DSR presence (gold stars)
pin_locations = [
    (6.3, 3.8, "Kenya"),
    (5.9, 2.4, "Serbia"),
    (3.2, 4.8, "Argentina"),
    (7.6, 3.5, "Pakistan"),
    (7.2, 2.8, "Kazakhstan"),
    (8.3, 4.8, "PNG")
]

for x, y, country in pin_locations:
    pin = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(x), Inches(y),
        Inches(0.15), Inches(0.15)
    )
    pin.fill.solid()
    pin.fill.fore_color.rgb = GOLD
    pin.line.width = Pt(0)
    pin.rotation = 180  # Point downward like a pin

# Bottom context note
context = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
tf = context.text_frame
tf.text = "BRI/DSR inclusion limited to dual-use tech-transfer contexts: telecom infrastructure, data centers, satellite ground stations, cloud services, and smart city surveillance systems. Focus on technologies enabling control, monitoring, and military-relevant capabilities."
tf.paragraphs[0].font.size = Pt(10)
tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(slide, """BRI and DSR as dual-use technology transfer mechanisms:

BELT AND ROAD INITIATIVE (BRI, 2013–):
Officially presented as economic development and infrastructure connectivity, BRI creates pathways for dual-use technology proliferation. Key dual-use components:
- Telecommunications infrastructure (Huawei 5G networks in 70+ countries)
- Port facilities (enabling naval access, surveillance, logistics for PLA power projection)
- Transportation corridors (dual civilian-military use for rapid deployment)
- Energy infrastructure (grid access, critical dependency creation)

DIGITAL SILK ROAD (DSR, 2015–):
The ICT subset of BRI, explicitly focused on digital infrastructure:
- 5G networks (potential backdoors, surveillance capabilities, wartime disruption)
- Fiber-optic cables (data interception, network mapping)
- Data centers and cloud services (access to foreign data, AI training datasets)
- Smart city systems (population monitoring, social control technologies)
- Satellite systems (BeiDou positioning for military use, ground stations for space operations)

SCOPE CLARIFICATION:
This presentation limits BRI/DSR discussion to dual-use technology transfer contexts. Not all BRI projects involve significant dual-use risks (e.g., hydroelectric dams without network connectivity components). Focus is on:
1. Technologies enabling surveillance, control, or military capabilities
2. Infrastructure creating strategic dependencies
3. Data access enabling intelligence collection
4. Systems that could be weaponized in conflict

GEOGRAPHIC SPREAD:
BRI/DSR spans 150+ countries, but dual-use technology concentration is highest in:
- Africa: Smart cities (Kenya, Ethiopia), telecommunications (ZTE/Huawei deployments)
- Central/Eastern Europe: 5G networks (Serbia, Hungary), data centers
- Central Asia: BeiDou ground stations, telecommunications backhaul
- Latin America: Space facilities (Argentina), smart city pilots
- Southeast Asia: Comprehensive digital infrastructure packages (Malaysia, Thailand, Laos)

CAPACITY GAPS ENABLING BRI/DSR DUAL-USE PROLIFERATION:
- Lack of technology impact assessments in procurement decisions
- Insufficient cybersecurity standards for critical infrastructure
- Weak interagency coordination (foreign ministry, defense, intelligence)
- Limited technical expertise to assess dual-use implications of telecom/data systems
- Vulnerability to debt-trap dynamics creating pressure to accept unfavorable terms

MILITARY-CIVIL FUSION CONNECTION:
Many BRI/DSR contractors (Huawei, ZTE, CETC, CASIC subsidiaries) have documented MCF links. Technology deployed through BRI can be remotely accessed, monitored, or disrupted by Chinese state actors under legal compulsion.""")

print("[OK] Slide 9 complete")

# Save progress
print("\nSaving slides 7-9...")
prs.save("C:/Projects/OSINT - Foresight/MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")
print("Progress saved. Slides 1-9 now complete.")
print("=" * 70)
