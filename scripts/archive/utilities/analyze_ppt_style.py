from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from collections import defaultdict
import json

def rgb_to_hex(rgb):
    """Convert RGB tuple to hex color code"""
    if rgb is None:
        return None
    return '#{:02x}{:02x}{:02x}'.format(rgb[0], rgb[1], rgb[2])

def analyze_color(color):
    """Extract color information"""
    if color is None:
        return None

    try:
        if color.type == 1:  # RGB color
            return {
                'type': 'RGB',
                'value': rgb_to_hex((color.rgb[0], color.rgb[1], color.rgb[2])),
                'brightness': color.brightness if hasattr(color, 'brightness') else None
            }
        elif color.type == 2:  # Theme color
            return {
                'type': 'Theme',
                'theme_color': str(color.theme_color) if hasattr(color, 'theme_color') else None
            }
    except:
        return None
    return None

def analyze_font(font):
    """Extract font information"""
    info = {}
    try:
        if font.name:
            info['name'] = font.name
        if font.size:
            info['size_pt'] = font.size.pt
        if font.bold is not None:
            info['bold'] = font.bold
        if font.italic is not None:
            info['italic'] = font.italic
        if font.underline is not None:
            info['underline'] = bool(font.underline)
        if hasattr(font, 'color') and font.color:
            info['color'] = analyze_color(font.color)
    except:
        pass
    return info if info else None

def analyze_slide(slide, slide_num):
    """Analyze a single slide"""
    analysis = {
        'slide_number': slide_num,
        'shapes': [],
        'text_elements': [],
        'colors_used': set(),
        'fonts_used': set(),
        'layout_info': {}
    }

    # Analyze background
    try:
        if slide.background.fill.type:
            analysis['background'] = {
                'fill_type': str(slide.background.fill.type)
            }
            if hasattr(slide.background.fill, 'fore_color'):
                bg_color = analyze_color(slide.background.fill.fore_color)
                if bg_color:
                    analysis['background']['color'] = bg_color
    except:
        pass

    # Analyze each shape
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            'index': idx,
            'type': str(shape.shape_type),
            'name': shape.name if hasattr(shape, 'name') else None,
            'left': shape.left.inches if hasattr(shape.left, 'inches') else None,
            'top': shape.top.inches if hasattr(shape.top, 'inches') else None,
            'width': shape.width.inches if hasattr(shape.width, 'inches') else None,
            'height': shape.height.inches if hasattr(shape.height, 'inches') else None,
        }

        # Handle text frames
        if shape.has_text_frame:
            text_content = shape.text_frame.text
            if text_content.strip():
                shape_info['text'] = text_content[:100]  # First 100 chars
                shape_info['text_length'] = len(text_content)

                # Analyze paragraphs
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    para_info = {
                        'text': para.text[:100] if para.text else '',
                        'level': para.level,
                        'alignment': str(para.alignment) if para.alignment else None,
                    }

                    # Analyze runs (formatted text segments)
                    runs = []
                    for run in para.runs:
                        run_info = {
                            'text': run.text[:50] if run.text else '',
                        }
                        font_info = analyze_font(run.font)
                        if font_info:
                            run_info['font'] = font_info
                            if 'name' in font_info:
                                analysis['fonts_used'].add(font_info['name'])
                        runs.append(run_info)

                    para_info['runs'] = runs
                    paragraphs.append(para_info)

                shape_info['paragraphs'] = paragraphs

        # Handle fills
        if hasattr(shape, 'fill'):
            try:
                fill_info = {
                    'type': str(shape.fill.type)
                }
                if shape.fill.type == 1:  # Solid fill
                    color = analyze_color(shape.fill.fore_color)
                    if color:
                        fill_info['color'] = color
                        if color.get('value'):
                            analysis['colors_used'].add(color['value'])
                shape_info['fill'] = fill_info
            except:
                pass

        # Handle lines
        if hasattr(shape, 'line'):
            try:
                line_info = {
                    'width_pt': shape.line.width.pt if shape.line.width else None,
                }
                if hasattr(shape.line, 'color'):
                    color = analyze_color(shape.line.color)
                    if color:
                        line_info['color'] = color
                        if color.get('value'):
                            analysis['colors_used'].add(color['value'])
                shape_info['line'] = line_info
            except:
                pass

        # Handle tables
        if hasattr(shape, 'table'):
            try:
                shape_info['table'] = {
                    'rows': len(shape.table.rows),
                    'columns': len(shape.table.columns)
                }
            except:
                pass

        # Handle pictures
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info['is_picture'] = True

        analysis['shapes'].append(shape_info)

    # Convert sets to lists for JSON serialization
    analysis['colors_used'] = list(analysis['colors_used'])
    analysis['fonts_used'] = list(analysis['fonts_used'])

    return analysis

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')

# Analyze slides 3-8 (indices 2-7)
results = {
    'presentation_info': {
        'total_slides': len(prs.slides),
        'analyzed_slides': '3-8',
        'slide_width_inches': prs.slide_width.inches,
        'slide_height_inches': prs.slide_height.inches,
    },
    'slides': []
}

for i in range(2, 8):  # Slides 3-8 (0-indexed)
    if i < len(prs.slides):
        slide_analysis = analyze_slide(prs.slides[i], i + 1)
        results['slides'].append(slide_analysis)

# Save results
with open('C:/Projects/OSINT-Foresight/ppt_style_analysis.json', 'w', encoding='utf-8') as f:
    json.dump(results, f, indent=2, ensure_ascii=False)

print("Analysis complete! Results saved to ppt_style_analysis.json")
print(f"\nSlide dimensions: {results['presentation_info']['slide_width_inches']:.2f}\" x {results['presentation_info']['slide_height_inches']:.2f}\"")
print(f"Analyzed slides: {len(results['slides'])}")
