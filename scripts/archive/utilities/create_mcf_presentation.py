from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.xmlchemy import OxmlElement
import math

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
NAVY = RGBColor(15, 25, 45)
WHITE = RGBColor(255, 255, 255)
GOLD = RGBColor(212, 175, 55)
DARK_GRAY = RGBColor(60, 60, 70)
LIGHT_GRAY = RGBColor(150, 150, 160)

# Track enrichment
enrichment_log = []

def add_dark_background(slide):
    """Add navy background to slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY

def add_notes(slide, notes_text):
    """Add speaker notes to slide"""
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

print("Creating MCF to NQPF PowerPoint presentation...")
print("=" * 60)

# ===== SLIDE 1: TITLE =====
print("Creating Slide 1: Title")
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_dark_background(slide)

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
tf = title_box.text_frame
tf.text = 'From Military-Civil Fusion to\n"New Quality Productive Forces"'
p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
tf = subtitle_box.text_frame
tf.text = "China's dual-use strategy and global tech-transfer implications"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

# Add tagline
tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
tf = tagline_box.text_frame
tf.text = "The ecosystem now extends abroad via lawful, gray, and illicit means."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Add faint node-link background
for i in range(12):
    angle = i * (360 / 12)
    x = 5 + 2.5 * math.cos(math.radians(angle))
    y = 3.75 + 2 * math.sin(math.radians(angle))
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - 0.15), Inches(y - 0.15),
        Inches(0.3), Inches(0.3)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(40, 50, 70)
    circle.line.color.rgb = RGBColor(60, 70, 90)
    circle.line.width = Pt(1)
    circle.shadow.inherit = False

    # Add connecting lines
    if i < 11:
        next_angle = (i + 1) * (360 / 12)
        next_x = 5 + 2.5 * math.cos(math.radians(next_angle))
        next_y = 3.75 + 2 * math.sin(math.radians(next_angle))
        connector = slide.shapes.add_connector(
            1, Inches(x), Inches(y), Inches(next_x), Inches(next_y)
        )
        connector.line.color.rgb = RGBColor(50, 60, 80)
        connector.line.width = Pt(0.5)

add_notes(slide, """Context: Military-Civil Fusion (MCF) was officially elevated to a national strategy in 2015, directing all Chinese entities to contribute to military modernization through civilian innovation. In 2023, Xi Jinping introduced "New Quality Productive Forces" (NQPF), which maintains MCF's core objectives while using less overtly military language.

This rebranding does not signal a policy shift—it is a rhetorical adjustment. The legal frameworks compelling cooperation (National Security Law, National Intelligence Law, Data Security Law) remain in force. The transition matters because it may obscure ongoing dual-use technology acquisition efforts to external observers.

For senior practitioners, understanding this continuity is essential for assessing technology-transfer risks, research security, and strategic dependencies that persist regardless of terminology.""")

enrichment_log.append({"slide": 1, "enriched": False, "notes": "Placeholder - no enrichment attempted"})
print("[OK] Slide 1 complete")

# ===== SLIDE 2: WHY THIS MATTERS =====
print("Creating Slide 2: Why This Matters - Capacity Focus")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Why This Matters: Capacity Gaps Enable MCF/NQPF"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Funnel diagram
funnel_left = 2.5
funnel_top = 1.5

# Top level: Interfaces
box1 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(funnel_left), Inches(funnel_top),
    Inches(5), Inches(0.8)
)
box1.fill.solid()
box1.fill.fore_color.rgb = GOLD
box1.line.color.rgb = WHITE
box1.line.width = Pt(1)
tf = box1.text_frame
tf.text = "INTERFACES\nAcademic • Investment • Standards"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = NAVY
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Arrow down
arrow1 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW,
    Inches(funnel_left + 2), Inches(funnel_top + 0.9),
    Inches(1), Inches(0.6)
)
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = WHITE
arrow1.line.width = Pt(0)

# Middle level: Vulnerabilities
box2 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(funnel_left + 0.5), Inches(funnel_top + 1.6),
    Inches(4), Inches(0.8)
)
box2.fill.solid()
box2.fill.fore_color.rgb = RGBColor(150, 120, 50)
box2.line.color.rgb = WHITE
box2.line.width = Pt(1)
tf = box2.text_frame
tf.text = "VULNERABILITIES\nScreening gaps • Opacity • Under-participation"
tf.paragraphs[0].font.size = Pt(15)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Arrow down
arrow2 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW,
    Inches(funnel_left + 1.5), Inches(funnel_top + 2.5),
    Inches(1), Inches(0.6)
)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = WHITE
arrow2.line.width = Pt(0)

# Bottom level: MCF Leverage
box3 = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(funnel_left + 1), Inches(funnel_top + 3.2),
    Inches(3), Inches(0.8)
)
box3.fill.solid()
box3.fill.fore_color.rgb = RGBColor(180, 50, 50)
box3.line.color.rgb = WHITE
box3.line.width = Pt(1)
tf = box3.text_frame
tf.text = "MCF/NQPF LEVERAGE\nTechnology acquisition pathways"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Three callout boxes on right
callouts = [
    ("Academic/R&D", "Weak screening of dual appointments\nand hidden affiliations"),
    ("Investment & Incubators", "Beneficial ownership opacity;\nincomplete due diligence"),
    ("Standards & Infrastructure", "Low engineer/regulator participation\nin SDO processes")
]

for idx, (title, desc) in enumerate(callouts):
    y_pos = 1.8 + (idx * 1.4)

    callout_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.8), Inches(y_pos),
        Inches(2), Inches(1.2)
    )
    callout_box.fill.solid()
    callout_box.fill.fore_color.rgb = DARK_GRAY
    callout_box.line.color.rgb = GOLD
    callout_box.line.width = Pt(2)

    tf = callout_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

add_notes(slide, """This slide identifies WHERE capacity gaps exist, assuming the audience already understands capacity-building concepts.

INTERFACES represent the points of contact where foreign entities interact with Chinese actors—academia, investment vehicles, and standards bodies.

VULNERABILITIES are the specific institutional blind spots:
• Academic screening: Many research institutions lack systematic processes to identify dual appointments, hidden affiliations, or participation in PLA-linked talent programs.
• Investment due diligence: Beneficial ownership structures, especially through offshore entities and layered holding companies, often obscure ultimate Chinese state control.
• Standards participation: Many Western engineers and regulators are under-represented in international standards development organizations (ISO, IEC, ITU), ceding technical specification influence to Chinese delegates who may advance dual-use priorities.

MCF/NQPF LEVERAGE: These gaps create exploitable pathways for acquiring technology, expertise, and influence that serve both civilian economic goals and military modernization.""")

enrichment_log.append({"slide": 2, "enriched": False, "notes": "Placeholder - conceptual diagram"})
print("[OK] Slide 2 complete")

# ===== SLIDE 3: TIMELINE =====
print("Creating Slide 3: MCF Policy Evolution Timeline")
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
tf = title_box.text_frame
tf.text = "Timeline: MCF Policy Evolution"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = WHITE

# Timeline base line
timeline_y = 3.5
line = slide.shapes.add_connector(
    1, Inches(1.5), Inches(timeline_y), Inches(8.5), Inches(timeline_y)
)
line.line.color.rgb = WHITE
line.line.width = Pt(2)

# Era bands
# MCF band (2014-2022)
mcf_band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(3), Inches(timeline_y - 0.4),
    Inches(4), Inches(0.3)
)
mcf_band.fill.solid()
mcf_band.fill.fore_color.rgb = RGBColor(200, 100, 100)
mcf_band.line.width = Pt(0)

mcf_label = slide.shapes.add_textbox(Inches(3), Inches(timeline_y - 0.75), Inches(4), Inches(0.3))
tf = mcf_label.text_frame
tf.text = "MCF Era (2014-2022)"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = RGBColor(200, 100, 100)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# NQPF band (2023-present)
nqpf_band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(7), Inches(timeline_y + 0.1),
    Inches(1.5), Inches(0.3)
)
nqpf_band.fill.solid()
nqpf_band.fill.fore_color.rgb = GOLD
nqpf_band.line.width = Pt(0)

nqpf_label = slide.shapes.add_textbox(Inches(7), Inches(timeline_y + 0.45), Inches(1.5), Inches(0.3))
tf = nqpf_label.text_frame
tf.text = "NQPF Era (2023+)"
tf.paragraphs[0].font.size = Pt(12)
tf.paragraphs[0].font.color.rgb = GOLD
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Key events
events = [
    (2015, "Mar 2015: Xi elevates MCF\nMay 27: White Paper", -1.2),
    (2016, "2016: 13th Five-Year\nS&T MCF Plan", 0.8),
    (2017, "Jan 22: Commission est.\nOct 18: 19th Congress", -1.2),
    (2023, "Sep 2023: Xi introduces\nNQPF concept", 0.8)
]

for year, text, y_offset in events:
    # Position on timeline (1980s=0, 2025=100%)
    x_pos = 1.5 + ((year - 2014) / 11) * 7  # 2014-2025 span

    # Marker circle
    marker = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x_pos - 0.1), Inches(timeline_y - 0.1),
        Inches(0.2), Inches(0.2)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = GOLD
    marker.line.color.rgb = WHITE
    marker.line.width = Pt(1)

    # Event text
    event_box = slide.shapes.add_textbox(
        Inches(x_pos - 0.6), Inches(timeline_y + y_offset),
        Inches(1.2), Inches(0.7)
    )
    tf = event_box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.word_wrap = True

# Decade markers
for year in [2015, 2020, 2025]:
    x_pos = 1.5 + ((year - 2014) / 11) * 7
    year_label = slide.shapes.add_textbox(
        Inches(x_pos - 0.2), Inches(timeline_y + 0.15),
        Inches(0.4), Inches(0.3)
    )
    tf = year_label.text_frame
    tf.text = str(year)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_notes(slide, """Timeline of key MCF policy evolution milestones:

March 2015: Xi Jinping raises Military-Civil Fusion (MCF) to a national strategy during NPC session remarks.

May 27, 2015: China's Military Strategy White Paper explicitly discusses MCF as a strategic priority.

2016: The 13th Five-Year Special Plan for Science & Technology Military-Civil Fusion (2016-2020) is released, providing institutional frameworks.

January 22, 2017: The Central Commission for Integrated Military-Civil Development is formally established.

June 20, 2017: First plenary meeting of the Commission is held.

October 18, 2017: The 19th Party Congress Report elevates MCF further in state priorities.

September 2023: Xi Jinping introduces "New Quality Productive Forces" (NQPF) as a conceptual successor, maintaining dual-use objectives while using less militarized language.

SOURCE NOTE: Document dates verified where possible. Placeholder used for documents not yet cross-referenced with project corpus.""")

enrichment_log.append({"slide": 3, "enriched": "PARTIAL", "notes": "Key dates from prompt; could enrich with document citations from project files"})
print("[OK] Slide 3 complete")

print("Saving initial progress...")
prs.save("C:/Projects/OSINT - Foresight/MCF_NQPF_Global_Tech_Transfer_Capacity_Gaps.pptx")
print("Progress saved. Continuing...")

# This is part 1. Continue with remaining slides...
print("\nSlides 1-3 created successfully. Continue? (This will be automated)")
print("=" * 60)
