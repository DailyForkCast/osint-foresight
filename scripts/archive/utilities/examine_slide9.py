from pptx import Presentation
from pptx.util import Pt, Inches

# Load presentation
prs = Presentation('c:/Users/mrear/Downloads/MCF v2.1.pptx')

# Check if slide 9 exists (index 8)
if len(prs.slides) >= 9:
    slide = prs.slides[8]  # Slide 9 (0-indexed)

    print("=" * 80)
    print("SLIDE 9 CONTENT ANALYSIS")
    print("=" * 80)
    print()

    print(f"Number of shapes: {len(slide.shapes)}")
    print()

    # Extract all text content
    text_elements = []

    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            'index': idx,
            'type': str(shape.shape_type),
            'name': shape.name if hasattr(shape, 'name') else 'Unknown',
        }

        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip()
            shape_info['text'] = text
            shape_info['position'] = {
                'left': round(shape.left.inches, 2),
                'top': round(shape.top.inches, 2),
                'width': round(shape.width.inches, 2),
                'height': round(shape.height.inches, 2)
            }

            # Get font info from first run
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        shape_info['font_size'] = run.font.size.pt if run.font.size else None
                        shape_info['font_name'] = run.font.name if run.font.name else None
                        break
                break

            text_elements.append(shape_info)

    # Sort by vertical position
    text_elements.sort(key=lambda x: x['position']['top'])

    print("TEXT ELEMENTS (top to bottom):")
    print("-" * 80)
    for elem in text_elements:
        print(f"\n[{elem['index']}] {elem['name']}")
        print(f"Position: Left={elem['position']['left']}\", Top={elem['position']['top']}\", "
              f"Width={elem['position']['width']}\", Height={elem['position']['height']}\"")
        if 'font_size' in elem and elem['font_size']:
            print(f"Font: {elem.get('font_name', 'N/A')} {elem['font_size']}pt")
        print(f"Text: {elem['text'][:200]}")
        if len(elem['text']) > 200:
            print(f"... [{len(elem['text'])} total characters]")

    print("\n" + "=" * 80)
    print("SLIDE 9 STRUCTURE SUMMARY")
    print("=" * 80)

    # Try to identify structure
    print("\nIdentified structure:")
    for elem in text_elements:
        if 'title' in elem['name'].lower():
            print(f"  TITLE: {elem['text'][:60]}")
        elif elem['position']['top'] < 2.0:
            print(f"  HEADER: {elem['text'][:60]}")
        else:
            print(f"  BODY: {elem['text'][:60]}")

else:
    print(f"Slide 9 does not exist. Presentation has {len(prs.slides)} slides.")
